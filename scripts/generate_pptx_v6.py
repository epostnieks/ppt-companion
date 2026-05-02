#!/usr/bin/env python3
"""
PPTX Generator V6 — EVERY SLIDE IS A McKINSEY FIGURE

Same pipeline as the 10/10 paper figures: HTML → Playwright → PNG → PPTX.
Every slide rendered at 1400×1050 with full CSS control.

McKinsey 100/100 rules (from SlideUpLift/Slideworks):
  - Title: 28pt ACTION TITLE (complete sentence = the insight)
  - Body: 18pt, chart labels 14pt
  - Charts/visuals: 70-80% of slide area
  - Max 3-5 bullets, each 1-2 lines
  - ONE insight per slide
  - 2-3 colors + gray
  - Source citation at bottom
  - Liberation Sans (our standard)

Standard slide: 10"×7.5" = 1400×1050 at native resolution
"""

import json, os, re, sys, glob, traceback, tempfile, subprocess
from concurrent.futures import ThreadPoolExecutor, as_completed
from pathlib import Path
from html.parser import HTMLParser
from io import BytesIO

from pptx import Presentation
from pptx.util import Inches, Emu

# ── Paths ──
SCRIPT_DIR = Path(__file__).parent
PROJECT_DIR = SCRIPT_DIR.parent
PAPER_DATA_DIR = PROJECT_DIR / "src" / "paperData"
PPTX_DIR = PROJECT_DIR / "public" / "pptx"
AUDIO_DIR = PROJECT_DIR / "public" / "audio"
PPTX_DIR.mkdir(parents=True, exist_ok=True)
AUDIO_DIR.mkdir(parents=True, exist_ok=True)
BASE = Path("/Users/erikpostnieks/Projects")

SLIDE_W = Inches(10)
SLIDE_H = Inches(7.5)
RENDER_W = 1400
RENDER_H = 1050

# ── Colors (McKinsey palette) ──
NAVY = "#1F4E79"
CREAM = "#FBF9F4"
TAN = "#F5EFE0"
GOLD = "#C4972A"
SOURCE_BLUE = "#3D6B8E"
DARK_BG = "#0F172A"
CYAN = "#00D4FF"
RED = "#CC3333"
GREEN = "#228855"
SURFACE = "#F7F5F0"
WHITE = "#FFFFFF"
DARK_TEXT = "#333333"
MID_GRAY = "#888888"
LIGHT_GRAY = "#E8E0D0"


# ══════════════════════════════════════════════════════════════
# HTML TEMPLATES — McKinsey CSS
# ══════════════════════════════════════════════════════════════

CSS_BASE = """
* { margin: 0; padding: 0; box-sizing: border-box; }
body {
  width: 1400px; height: 1050px; overflow: hidden;
  font-family: 'Liberation Sans', 'Helvetica Neue', Arial, sans-serif;
  -webkit-font-smoothing: antialiased;
  background: #FFFFFF;
}
.slide { width: 1400px; height: 1050px; position: relative; overflow: hidden; }

/* Title bar */
.title-bar {
  background: %(NAVY)s; color: white;
  padding: 18px 50px 14px; min-height: 70px;
}
.title-bar h1 {
  font-size: 26px; font-weight: 700; line-height: 1.3;
  letter-spacing: -0.3px;
}
.subtitle-bar {
  background: %(CREAM)s; padding: 8px 50px;
  font-size: 14px; color: %(DARK_TEXT)s; font-style: italic;
}

/* Body content area */
.body { padding: 20px 50px 0; }
.body-text { font-size: 17px; line-height: 1.65; color: %(DARK_TEXT)s; }
.body-small { font-size: 14px; line-height: 1.55; color: %(DARK_TEXT)s; }

/* Bullet list */
.bullets { list-style: none; padding: 0; }
.bullets li {
  font-size: 17px; line-height: 1.55; color: %(DARK_TEXT)s;
  padding: 8px 0 8px 24px; position: relative;
  border-bottom: 1px solid #F0ECE4;
}
.bullets li:last-child { border-bottom: none; }
.bullets li::before {
  content: ''; position: absolute; left: 0; top: 16px;
  width: 8px; height: 8px; background: %(NAVY)s; border-radius: 2px;
}

/* Cards */
.card {
  background: %(SURFACE)s; border-radius: 8px;
  padding: 18px 22px; margin: 8px 0;
}
.card-title { font-size: 14px; font-weight: 700; color: %(NAVY)s; margin-bottom: 6px; text-transform: uppercase; letter-spacing: 1px; }
.card-value { font-size: 28px; font-weight: 700; margin-bottom: 4px; }
.card-desc { font-size: 12px; color: %(MID_GRAY)s; }

/* Metric grid */
.metrics-grid {
  display: grid; grid-template-columns: 1fr 1fr;
  gap: 14px; margin: 16px 0;
}

/* Key Insight box */
.insight-box {
  background: %(TAN)s; border-left: 4px solid %(GOLD)s;
  border-radius: 4px; padding: 14px 22px; margin: 14px 0;
}
.insight-label { font-size: 11px; font-weight: 700; color: %(GOLD)s; letter-spacing: 1.5px; text-transform: uppercase; margin-bottom: 4px; }
.insight-text { font-size: 14px; color: %(DARK_TEXT)s; line-height: 1.5; font-weight: 600; }

/* Source line */
.source {
  position: absolute; bottom: 14px; left: 50px; right: 50px;
  font-size: 11px; color: %(SOURCE_BLUE)s;
}

/* Table */
.data-table { width: 100%%; border-collapse: collapse; font-size: 13px; }
.data-table th {
  background: %(NAVY)s; color: white; padding: 8px 12px;
  text-align: left; font-weight: 700; font-size: 12px;
  letter-spacing: 0.5px; text-transform: uppercase;
}
.data-table td { padding: 7px 12px; border-bottom: 1px solid #E8E0D0; color: %(DARK_TEXT)s; }
.data-table tr:nth-child(even) td { background: %(SURFACE)s; }

/* Two-column layout */
.two-col { display: grid; grid-template-columns: 1fr 1fr; gap: 30px; }
.three-col { display: grid; grid-template-columns: 1fr 1fr 1fr; gap: 20px; }

/* Dark slide */
.dark { background: %(DARK_BG)s; color: white; }
.dark .title-bar { background: transparent; border-bottom: 2px solid rgba(255,255,255,0.1); }

/* Channel row */
.channel-row {
  display: flex; align-items: center; padding: 10px 0;
  border-bottom: 1px solid #F0ECE4;
}
.channel-name { flex: 0 0 200px; font-size: 14px; font-weight: 700; color: %(NAVY)s; }
.channel-bar-wrap { flex: 1; margin: 0 16px; height: 22px; background: #F0ECE4; border-radius: 4px; position: relative; }
.channel-bar { height: 100%%; border-radius: 4px; background: linear-gradient(90deg, %(NAVY)s, #3B7DB5); }
.channel-value { flex: 0 0 80px; font-size: 14px; font-weight: 700; color: %(RED)s; text-align: right; }
.channel-pct { flex: 0 0 50px; font-size: 12px; color: %(MID_GRAY)s; text-align: right; }

/* Agent card */
.agent-card {
  background: %(SURFACE)s; border-radius: 8px; padding: 20px;
  border-left: 4px solid %(NAVY)s;
}
.agent-role { font-size: 13px; font-weight: 700; color: %(GOLD)s; text-transform: uppercase; letter-spacing: 1px; margin-bottom: 6px; }
.agent-subtitle { font-size: 12px; color: %(MID_GRAY)s; margin-bottom: 10px; }
.agent-body { font-size: 14px; color: %(DARK_TEXT)s; line-height: 1.5; }
""" % {
    'NAVY': NAVY, 'CREAM': CREAM, 'TAN': TAN, 'GOLD': GOLD,
    'SOURCE_BLUE': SOURCE_BLUE, 'DARK_BG': DARK_BG, 'CYAN': CYAN,
    'RED': RED, 'GREEN': GREEN, 'SURFACE': SURFACE, 'DARK_TEXT': DARK_TEXT,
    'MID_GRAY': MID_GRAY, 'LIGHT_GRAY': LIGHT_GRAY, 'WHITE': WHITE,
}


# ══════════════════════════════════════════════════════════════
# TEXT UTILITIES
# ══════════════════════════════════════════════════════════════

class HTMLStripper(HTMLParser):
    def __init__(self):
        super().__init__()
        self.result = []
    def handle_data(self, d):
        self.result.append(d)
    def handle_starttag(self, tag, attrs):
        if tag in ('td', 'th'): self.result.append(' | ')
        if tag == 'tr': self.result.append('\n')
        if tag == 'li': self.result.append('\n• ')
        if tag in ('br', 'p'): self.result.append('\n')
    def get_text(self):
        return re.sub(r'\n{3,}', '\n\n', ''.join(self.result)).strip()

def strip(text):
    if not text or not isinstance(text, str): return str(text) if text else ""
    s = HTMLStripper()
    s.feed(text)
    return s.get_text()

def esc(text):
    """HTML-escape text."""
    return str(text).replace('&', '&amp;').replace('<', '&lt;').replace('>', '&gt;').replace('"', '&quot;')

def sentences(text):
    return [s.strip() for s in re.split(r'(?<=[.!?])\s+', strip(text)) if s.strip()]

def paginate(text, max_chars=1200):
    """Split into page-sized chunks for dense slides."""
    sents = sentences(text)
    pages, cur = [], ""
    for s in sents:
        if len(cur) + len(s) > max_chars and cur:
            pages.append(cur.strip())
            cur = s
        else:
            cur = (cur + " " + s).strip()
    if cur: pages.append(cur.strip())
    return pages


# ══════════════════════════════════════════════════════════════
# PLAYWRIGHT RENDERER
# ══════════════════════════════════════════════════════════════

def render_html_to_png(browser, html_content):
    """Render HTML slide to PNG using Playwright."""
    page = browser.new_page(viewport={"width": RENDER_W, "height": RENDER_H})
    page.set_content(html_content, wait_until="networkidle")
    png_bytes = page.screenshot(type="png")
    page.close()
    return png_bytes

def make_slide_html(body_content):
    """Wrap body content in full slide HTML."""
    return f"""<!DOCTYPE html>
<html><head><meta charset="utf-8"><style>{CSS_BASE}</style></head>
<body><div class="slide">{body_content}</div></body></html>"""


# ══════════════════════════════════════════════════════════════
# SLIDE HTML BUILDERS
# ══════════════════════════════════════════════════════════════

def html_title_slide(data):
    title = esc(data.get("title", data.get("slug", "")))
    tt = data.get("theoremType", "")
    tn = esc(data.get("theoremName", ""))
    ep = esc(strip(data.get("epigraph", "")))

    badge = ""
    if tt and tn:
        color = RED if tt == "Impossibility" else GOLD
        label = tn if tt.lower() in tn.lower() else f"{tn} ({tt})"
        badge = f'<div style="font-size:18px;font-weight:700;color:{color};margin-top:20px;">{esc(label)}</div>'

    epigraph = f'<div style="font-size:16px;color:{MID_GRAY};font-style:italic;margin-top:30px;max-width:800px;">"{ep}"</div>' if ep else ""

    return make_slide_html(f"""
    <div class="dark" style="width:1400px;height:1050px;padding:60px;">
      <div style="border-top:3px solid {CYAN};width:80px;margin-bottom:30px;"></div>
      <div style="font-size:13px;font-weight:700;color:{CYAN};letter-spacing:2px;margin-bottom:20px;">SYSTEM ASSET PRICING MODEL — MASTERCLASS</div>
      <div style="font-size:38px;font-weight:700;color:white;line-height:1.2;max-width:900px;">{title}</div>
      {badge}
      {epigraph}
      <div style="position:absolute;bottom:60px;left:60px;font-size:13px;color:{MID_GRAY};">
        Erik Postnieks &nbsp;|&nbsp; Independent Researcher &nbsp;|&nbsp; 2026
      </div>
    </div>""")


def html_divider(num, title, subtitle=""):
    sub = f'<div style="font-size:16px;color:{CREAM};font-style:italic;margin-top:20px;max-width:700px;">{esc(subtitle)}</div>' if subtitle else ""
    return make_slide_html(f"""
    <div style="width:1400px;height:1050px;background:{NAVY};padding:60px;display:flex;flex-direction:column;justify-content:center;">
      <div style="font-size:14px;font-weight:700;color:{CYAN};letter-spacing:2px;margin-bottom:16px;">CHAPTER {num}</div>
      <div style="font-size:36px;font-weight:700;color:white;line-height:1.2;max-width:900px;">{esc(title)}</div>
      {sub}
    </div>""")


def html_metrics(data):
    wa = data.get("welfareAnalysis")
    wa = wa if isinstance(wa, dict) else {}
    beta = wa.get("betaW", data.get("betaW", "—"))
    ci = wa.get("ci90", data.get("ci90", "—"))
    pi_raw = wa.get("pi", "—")
    dw_raw = wa.get("deltaW", "—")
    tt = data.get("theoremType", "")
    tn = data.get("theoremName", "")

    def money(v):
        if not v or v == "—": return "—"
        m = re.search(r'[-−]?\$[\d,.]+\s*(billion|trillion|million)?', str(v), re.I)
        return m.group(0) if m else str(v)[:50]

    insight = ""
    if tt:
        desc = ("Physical or biological constraints prevent internalization. No policy reform closes this gap."
                if tt == "Impossibility" else
                "Institutional constraints. At least one country has demonstrated a proven reform.")
        insight = f"""<div class="insight-box">
          <div class="insight-label">Classification</div>
          <div class="insight-text">{esc(tn)}: {desc}</div>
        </div>"""

    return make_slide_html(f"""
    <div class="title-bar"><h1>For every dollar of revenue, this industry destroys ${esc(str(beta))} of social welfare</h1></div>
    <div class="subtitle-bar">System Beta and Welfare Analysis — Monte Carlo confirmed (N=10,000, seed=42)</div>
    <div class="body">
      <div class="metrics-grid">
        <div class="card">
          <div class="card-title">System Beta (βW)</div>
          <div class="card-value" style="color:{NAVY};">{esc(str(beta))}</div>
          <div class="card-desc">Welfare destruction per dollar of revenue</div>
        </div>
        <div class="card">
          <div class="card-title">90% Confidence Interval</div>
          <div class="card-value" style="color:{DARK_TEXT};font-size:22px;">{esc(str(ci))}</div>
          <div class="card-desc">Monte Carlo simulation bounds</div>
        </div>
        <div class="card">
          <div class="card-title">Π (Annual Revenue)</div>
          <div class="card-value" style="color:{GREEN};">{esc(money(pi_raw))}</div>
          <div class="card-desc">Private payoff — always revenue, never profit (Iron Law)</div>
        </div>
        <div class="card">
          <div class="card-title">ΔW (Welfare Cost)</div>
          <div class="card-value" style="color:{RED};">{esc(money(dw_raw))}</div>
          <div class="card-desc">Annual system welfare destruction</div>
        </div>
      </div>
      {insight}
    </div>
    <div class="source">Source: Monte Carlo simulation (N=10,000, seed=42). Π = annual industry revenue (Iron Law).</div>
    """)


def html_text_slide(action_title, text, subtitle="", source_text=""):
    """Dense text slide with action title. Text is pre-paginated to fit."""
    paragraphs = text.split('\n')
    body = ""
    for p in paragraphs:
        p = p.strip()
        if p:
            body += f'<p style="margin-bottom:12px;">{esc(p)}</p>'

    sub = f'<div class="subtitle-bar">{esc(subtitle)}</div>' if subtitle else ""
    src = f'<div class="source">Source: {esc(source_text)}</div>' if source_text else ""

    return make_slide_html(f"""
    <div class="title-bar"><h1>{esc(action_title)}</h1></div>
    {sub}
    <div class="body">
      <div class="body-text">{body}</div>
    </div>
    {src}""")


def html_bullets_slide(action_title, bullets, subtitle="", source_text=""):
    """Bullet list slide — 3-7 bullets with action title."""
    items = "".join(f"<li>{esc(b)}</li>" for b in bullets[:7])
    sub = f'<div class="subtitle-bar">{esc(subtitle)}</div>' if subtitle else ""
    src = f'<div class="source">Source: {esc(source_text)}</div>' if source_text else ""

    return make_slide_html(f"""
    <div class="title-bar"><h1>{esc(action_title)}</h1></div>
    {sub}
    <div class="body">
      <ul class="bullets">{items}</ul>
    </div>
    {src}""")


def html_channels(channels, slug):
    """Welfare channel overview with bars."""
    rows = ""
    max_val = 0
    for ch in channels:
        v = strip(str(ch.get("value", "0")))
        num = float(re.sub(r'[^0-9.]', '', v) or 0)
        if num > max_val: max_val = num

    for ch in channels[:8]:
        name = esc(ch.get("name", ""))
        value = esc(strip(str(ch.get("value", "—"))))
        v_num = float(re.sub(r'[^0-9.]', '', strip(str(ch.get("value", "0")))) or 0)
        pct = (v_num / max_val * 100) if max_val > 0 else 10

        rows += f"""<div class="channel-row">
          <div class="channel-name">{name}</div>
          <div class="channel-bar-wrap"><div class="channel-bar" style="width:{pct:.0f}%;"></div></div>
          <div class="channel-value">{value}</div>
        </div>"""

    return make_slide_html(f"""
    <div class="title-bar"><h1>The welfare cost decomposes into {len(channels)} distinct channels</h1></div>
    <div class="subtitle-bar">Welfare Channel Decomposition — {esc(slug)}</div>
    <div class="body">{rows}</div>
    <div class="source">Source: Monte Carlo simulation (N=10,000, seed=42)</div>""")


def html_agents_overview(agents):
    """Six-agent framework overview — all 6 on one slide."""
    meta = {
        "whistleblower": ("Whistleblower", "Break information asymmetry"),
        "plaintiff": ("Plaintiff", "Monetize liability"),
        "regulator": ("Regulator", "Redesign the game"),
        "legislator": ("Legislator", "Change the law"),
        "investor": ("Investor", "Reprice capital"),
        "supranational": ("Supranational", "Coordinate jurisdictions"),
    }
    cards = ""
    for key in ["whistleblower", "plaintiff", "regulator", "legislator", "investor", "supranational"]:
        content = agents.get(key) or (agents.get("insider") if key == "whistleblower" else None)
        if not content: continue
        label, sub = meta.get(key, (key.title(), ""))
        first = sentences(content)[0] if sentences(content) else ""
        cards += f"""<div class="agent-card">
          <div class="agent-role">{esc(label)}</div>
          <div class="agent-subtitle">{esc(sub)}</div>
          <div class="agent-body">{esc(first[:200])}</div>
        </div>"""

    return make_slide_html(f"""
    <div class="title-bar"><h1>Six agents can break the Private-Systemic Tension in this domain</h1></div>
    <div class="subtitle-bar">Conflictoring Framework — Each agent type has a specific intervention pathway</div>
    <div class="body">
      <div class="two-col">{cards}</div>
    </div>""")


def html_table(header, rows, title, source=""):
    """Data table slide."""
    th = "".join(f"<th>{esc(h[:30])}</th>" for h in header[:7])
    trs = ""
    for row in rows[:16]:
        tds = "".join(f"<td>{esc(str(c)[:35])}</td>" for c in row[:7])
        trs += f"<tr>{tds}</tr>"
    src = f'<div class="source">Source: {esc(source)}</div>' if source else ""

    return make_slide_html(f"""
    <div class="title-bar"><h1>{esc(title)}</h1></div>
    <div class="body">
      <table class="data-table"><thead><tr>{th}</tr></thead><tbody>{trs}</tbody></table>
    </div>
    {src}""")


def html_quiz(data):
    wa = data.get("welfareAnalysis")
    wa = wa if isinstance(wa, dict) else {}
    beta = wa.get("betaW", "?")
    tt = data.get("theoremType", "?")
    tn = data.get("theoremName", "?")

    qs = [
        (f"What is βW for this domain?", f"βW = {beta}"),
        (f"Impossibility or Intractability?", f"{tt}"),
        (f"What is the theorem called?", f"{tn}"),
        ("Can policy alone solve this?",
         "No — physical constraints" if tt == "Impossibility" else "Yes — proven models exist"),
    ]
    items = ""
    for q, a in qs:
        items += f"""<div style="margin-bottom:28px;">
          <div style="font-size:18px;font-weight:700;color:white;">{esc(q)}</div>
          <div style="font-size:16px;color:{CYAN};margin-top:6px;">{esc(a)}</div>
        </div>"""

    return make_slide_html(f"""
    <div class="dark" style="width:1400px;height:1050px;padding:60px;">
      <div style="font-size:13px;font-weight:700;color:{CYAN};letter-spacing:2px;margin-bottom:30px;">COMPREHENSION CHECK</div>
      {items}
    </div>""")


def html_closing(data):
    wa = data.get("welfareAnalysis")
    wa = wa if isinstance(wa, dict) else {}
    beta = wa.get("betaW", "")
    tt = data.get("theoremType", "")
    title = esc(data.get("title", ""))
    color = RED if tt == "Impossibility" else GOLD

    return make_slide_html(f"""
    <div class="dark" style="width:1400px;height:1050px;display:flex;flex-direction:column;align-items:center;justify-content:center;text-align:center;">
      <div style="font-size:20px;color:white;font-weight:700;max-width:800px;margin-bottom:30px;">{title}</div>
      <div style="width:100px;height:2px;background:{CYAN};margin-bottom:30px;"></div>
      <div style="font-size:52px;font-weight:700;color:{CYAN};">βW = {esc(str(beta))}</div>
      <div style="font-size:20px;color:{color};margin-top:16px;">{esc(tt)} Theorem</div>
      <div style="font-size:13px;color:{MID_GRAY};margin-top:40px;">Erik Postnieks &nbsp;|&nbsp; Independent Researcher &nbsp;|&nbsp; postnieks.com</div>
      <div style="font-size:11px;color:{MID_GRAY};margin-top:8px;">Paper: SSRN &nbsp;|&nbsp; MC Replication: github.com/epostnieks</div>
      <div style="font-size:11px;color:{CYAN};font-weight:700;margin-top:20px;letter-spacing:2px;">END OF MASTERCLASS</div>
    </div>""")


# ══════════════════════════════════════════════════════════════
# ASSET EXTRACTION (same as V5)
# ══════════════════════════════════════════════════════════════

def find_figures(slug):
    sapm_png = BASE / f"sapm-{slug}" / "figures" / "png"
    if sapm_png.exists():
        pngs = sorted(sapm_png.glob("*.png"))
        if pngs:
            seen = set()
            out = []
            for p in pngs:
                base = re.sub(r'_v\d+\.png$', '.png', p.name)
                if base not in seen:
                    out.append(p)
                    seen.add(base)
            return out
    sapm_figs = BASE / f"sapm-{slug}" / "figures"
    if sapm_figs.exists():
        pngs = sorted(sapm_figs.glob("*.png"))
        if pngs: return pngs
    return []

def find_docx(slug):
    for pat in [f"{BASE}/sapm-{slug}/paper/{slug}-current.docx",
                f"{BASE}/sapm-{slug}/paper/*-current.docx",
                f"{BASE}/sapm-{slug}/paper/*-final.docx",
                f"{BASE}/sapm-{slug}/paper/*.docx"]:
        matches = [m for m in glob.glob(pat) if '~' not in m and 'template' not in m.lower()]
        if matches: return matches[0]
    return None

def get_tables(path, max_t=50):
    try:
        from docx import Document as DD
        doc = DD(path)
        tables = []
        for i, t in enumerate(doc.tables[:max_t]):
            rows = [[c.text.strip() for c in r.cells] for r in t.rows]
            if len(rows) > 1:
                tables.append({"title": f"Table {i+1}", "header": rows[0], "rows": rows[1:]})
        return tables
    except: return []

def fig_label(p):
    return re.sub(r'^fig\d+_', '', p.stem).replace('_', ' ').title()


# ══════════════════════════════════════════════════════════════
# MAIN GENERATOR
# ══════════════════════════════════════════════════════════════

def generate(json_path, browser):
    """Generate one masterclass PPTX."""
    with open(json_path, 'r') as f:
        data = json.load(f)

    slug = json_path.stem
    figures = find_figures(slug)
    docx_path = find_docx(slug)
    tables = get_tables(docx_path) if docx_path else []

    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    def add_rendered_slide(html, notes_text=""):
        png = render_html_to_png(browser, html)
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        slide.shapes.add_picture(BytesIO(png), 0, 0, SLIDE_W, SLIDE_H)
        if notes_text:
            slide.notes_slide.notes_text_frame.text = str(notes_text)[:5000]

    def add_figure_slide(fig_path, notes_text=""):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        try:
            slide.shapes.add_picture(str(fig_path), 0, 0, SLIDE_W, SLIDE_H)
        except:
            pass
        if notes_text:
            slide.notes_slide.notes_text_frame.text = notes_text[:5000]

    total = 0

    # Title
    add_rendered_slide(html_title_slide(data),
        f"Welcome to the SAPM Masterclass on {data.get('title', slug)}.")
    total += 1

    # Ch 1: Numbers
    add_rendered_slide(html_divider(1, "THE NUMBERS", "System beta, welfare cost, and confidence intervals"))
    add_rendered_slide(html_metrics(data),
        f"Key metrics: beta W = {(data.get('welfareAnalysis') or {}).get('betaW', '?')}")
    total += 2

    # Insert MC + sensitivity figures
    for fig in figures:
        if any(k in fig.name for k in ['mc_dist', 'sensitivity', 'channel_breakdown']):
            add_figure_slide(fig, f"Figure: {fig_label(fig)}")
            total += 1

    # Ch 2: Executive Summary
    es = data.get("executiveSummary", "")
    if es:
        add_rendered_slide(html_divider(2, "EXECUTIVE SUMMARY", "The complete argument"))
        total += 1
        for page in paginate(es, 1000):
            sents = sentences(page)
            action = sents[0] if sents else "Executive Summary"
            add_rendered_slide(html_text_slide(action, page), page)
            total += 1

    # Ch 3: Key Findings
    findings = data.get("keyFindings", [])
    if findings:
        add_rendered_slide(html_divider(3, "KEY FINDINGS", "What the evidence shows"))
        total += 1
        # Group 3-5 findings per slide as bullets
        for i in range(0, len(findings), 4):
            batch = findings[i:i+4]
            bullets = [strip(f)[:200] if isinstance(f, str) else str(f)[:200] for f in batch]
            action = sentences(bullets[0])[0] if sentences(bullets[0]) else "Key findings from the analysis"
            add_rendered_slide(html_bullets_slide(action, bullets), "\n".join(bullets))
            total += 1

    # Ch 4: Theorem
    theorem = data.get("theorem", {})
    if theorem and isinstance(theorem, dict):
        tn = data.get("theoremName", "Formal Result")
        add_rendered_slide(html_divider(4, tn, "Formal proof and structural assumptions"))
        total += 1

        plain = theorem.get("plain", "")
        if plain:
            for page in paginate(plain, 1000):
                sents = sentences(page)
                action = sents[0] if sents else "The theorem in plain language"
                add_rendered_slide(html_text_slide(action, page, tn), page)
                total += 1

        axioms = theorem.get("axioms", [])
        if axioms:
            bullets = [f"{a.get('id','')}: {a.get('name','')} — {strip(a.get('statement',''))[:150]}" for a in axioms]
            add_rendered_slide(html_bullets_slide(
                f"The result rests on {len(axioms)} structural axioms", bullets,
                "Each axiom is empirically verifiable"), "\n".join(bullets))
            total += 1

        formal = theorem.get("formal", "")
        if formal:
            for page in paginate(formal, 1000):
                sents = sentences(page)
                action = sents[0][:120] if sents else "Formal statement"
                add_rendered_slide(html_text_slide(action, page), page)
                total += 1

        proof = theorem.get("proofSketch", "")
        if proof:
            for page in paginate(proof, 1000):
                sents = sentences(page)
                action = sents[0][:120] if sents else "Proof sketch"
                add_rendered_slide(html_text_slide(action, page, "Proof Sketch"), page)
                total += 1

    # Ch 5: Welfare Channels
    wa = data.get("welfareAnalysis")
    wa = wa if isinstance(wa, dict) else {}
    channels = wa.get("channels", [])
    if channels:
        add_rendered_slide(html_divider(5, "WELFARE CHANNELS", "Where the damage comes from"))
        add_rendered_slide(html_channels(channels, slug),
            "Welfare channel decomposition showing all channels.")
        total += 2

        for ch in channels:
            desc = strip(ch.get("description", ""))
            if len(desc) > 200:
                name = ch.get("name", "Channel")
                for page in paginate(desc, 1000):
                    sents = sentences(page)
                    action = sents[0][:120] if sents else f"{name} channel detail"
                    add_rendered_slide(html_text_slide(action, page, name), page)
                    total += 1

    # Ch 6: ALL remaining McKinsey figures
    used = {f.name for f in figures if any(k in f.name for k in ['mc_dist', 'sensitivity', 'channel_breakdown'])}
    remaining = [f for f in figures if f.name not in used]
    if remaining:
        add_rendered_slide(html_divider(6, "PAPER FIGURES",
            f"{len(remaining)} McKinsey-quality figures from the published research"))
        total += 1
        for fig in remaining:
            add_figure_slide(fig, f"{fig_label(fig)} — from the published paper.")
            total += 1

    # Ch 7: Evidence — cases + tables
    cases = data.get("caseStudies", [])
    if cases or tables:
        add_rendered_slide(html_divider(7, "EVIDENCE",
            f"{len(cases)} case studies, {len(tables)} data tables"))
        total += 1

        for case in cases:
            title = case.get("title", "Case Study")
            content = strip(case.get("content", case.get("description", "")))
            if content:
                for page in paginate(content, 1000):
                    sents = sentences(page)
                    action = sents[0][:120] if sents else title
                    add_rendered_slide(html_text_slide(action, page, title), page)
                    total += 1

        for t in tables[:20]:
            add_rendered_slide(html_table(t["header"], t["rows"],
                f"Data from the published paper ({len(t['rows'])} rows)",
                f"Extracted from {slug} Word document"))
            total += 1

    # Ch 8: Policy
    policy = data.get("policyAnalysis")
    if policy and isinstance(policy, dict):
        add_rendered_slide(html_divider(8, "POLICY ANALYSIS", "What works, what fails, and why"))
        total += 1
        for key, label in [("currentFramework", "Current Framework"),
                           ("failures", "Why It Fails"), ("reform", "Reform Pathway")]:
            content = policy.get(key, "")
            if content:
                for page in paginate(content, 1000):
                    sents = sentences(page)
                    action = sents[0][:120] if sents else label
                    add_rendered_slide(html_text_slide(action, page, label), page)
                    total += 1

    # Ch 9: Six Agents
    agents = data.get("agents", {})
    if agents and isinstance(agents, dict):
        add_rendered_slide(html_divider(9, "SIX-AGENT FRAMEWORK",
            "Who can break the Private-Systemic Tension"))
        add_rendered_slide(html_agents_overview(agents), "Six-agent framework overview.")
        total += 2

        meta = {"whistleblower": "Whistleblower", "plaintiff": "Plaintiff",
                "regulator": "Regulator", "legislator": "Legislator",
                "investor": "Investor", "supranational": "Supranational"}
        for key, label in meta.items():
            content = agents.get(key) or (agents.get("insider") if key == "whistleblower" else None)
            if not content: continue
            for page in paginate(content, 1000):
                sents = sentences(page)
                action = sents[0][:120] if sents else f"{label} intervention"
                add_rendered_slide(html_text_slide(action, page, f"Agent: {label}"), page)
                total += 1

    # Ch 10: Context
    add_rendered_slide(html_divider(10, "CONNECTIONS & CONTEXT",
        "Cross-domain links, literature, methodology"))
    total += 1

    for section_key, section_label in [("literatureContext", "Literature"),
                                        ("methodology", "Methodology")]:
        content = data.get(section_key, "")
        if content:
            for page in paginate(content, 1000):
                sents = sentences(page)
                action = sents[0][:120] if sents else section_label
                add_rendered_slide(html_text_slide(action, page, section_label), page)
                total += 1

    # Ch 11: FAQ
    faq = data.get("faq", [])
    if faq:
        add_rendered_slide(html_divider(11, "FAQ", "Anticipating the hard questions"))
        total += 1
        for item in faq:
            q = strip(item.get("q", ""))
            a = strip(item.get("a", ""))
            if q and a:
                add_rendered_slide(html_text_slide(q, a[:1000]), f"Q: {q}\nA: {a}")
                total += 1

    # Quiz + Closing
    add_rendered_slide(html_quiz(data), "Comprehension check.")
    add_rendered_slide(html_closing(data), "End of masterclass.")
    total += 2

    out_path = PPTX_DIR / f"{slug}.pptx"
    prs.save(str(out_path))
    return slug, total, len(figures), len(tables)


def main():
    from playwright.sync_api import sync_playwright

    json_files = sorted(PAPER_DATA_DIR.glob("*.json"))
    json_files = [f for f in json_files if '_raw' not in f.name]
    print(f"Found {len(json_files)} papers")

    results = []
    errors = []

    with sync_playwright() as pw:
        browser = pw.chromium.launch(headless=True)

        for f in json_files:
            try:
                slug, slides, figs, tables = generate(f, browser)
                results.append((slug, slides, figs, tables))
                print(f"  OK  {slug}: {slides} slides, {figs} figs, {tables} tables")
            except Exception as e:
                errors.append((f.stem, str(e)))
                print(f"  FAIL {f.stem}: {e}")
                traceback.print_exc()

        browser.close()

    total_slides = sum(r[1] for r in results)
    print(f"\nDone: {len(results)} OK, {len(errors)} failed")
    print(f"Total: {total_slides} slides (avg {total_slides // max(len(results),1)}/paper)")
    print(f"Output: {PPTX_DIR}")

    if errors:
        print(f"\nFailed:")
        for s, e in errors:
            print(f"  {s}: {e}")


if __name__ == "__main__":
    main()
