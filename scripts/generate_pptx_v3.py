#!/usr/bin/env python3
"""
PPTX Generator V3 — HTML→Playwright→PNG→PPTX pipeline.
Each slide is a fully designed HTML page rendered at 1920×1080, then embedded as a slide image.
McKinsey-quality design with gradients, shadows, typography, data visualization.
"""

import json, os, re, sys, glob, tempfile, shutil
from concurrent.futures import ThreadPoolExecutor, as_completed
from pathlib import Path
from html import escape as html_escape
from io import BytesIO

from pptx import Presentation
from pptx.util import Inches, Emu
from playwright.sync_api import sync_playwright

# ── Paths ──
PAPER_DATA_DIR = Path(__file__).parent.parent / "src" / "paperData"
OUTPUT_DIR = Path(__file__).parent.parent / "pptx_output"
OUTPUT_DIR.mkdir(exist_ok=True)
BASE = "/Users/erikpostnieks/Projects"

# ── HTML Templates ──

CSS_BASE = """
* { margin: 0; padding: 0; box-sizing: border-box; }
body {
  width: 1920px; height: 1080px; overflow: hidden;
  font-family: 'Liberation Sans', 'Helvetica Neue', Arial, sans-serif;
  -webkit-font-smoothing: antialiased;
}
.slide { width: 1920px; height: 1080px; position: relative; overflow: hidden; }

/* Navy gradient background */
.bg-navy { background: linear-gradient(135deg, #0B1D3A 0%, #1A3A5C 40%, #0F2847 100%); }
.bg-white { background: #FAFAF8; }
.bg-cream { background: linear-gradient(180deg, #FAF8F3 0%, #F5F0E8 100%); }
.bg-dark { background: linear-gradient(135deg, #080E1A 0%, #0F1B2D 50%, #0A1420 100%); }

/* Typography */
.title-hero { font-size: 52px; font-weight: 700; line-height: 1.15; letter-spacing: -1px; }
.title-section { font-size: 40px; font-weight: 700; line-height: 1.2; letter-spacing: -0.5px; }
.title-slide { font-size: 32px; font-weight: 700; line-height: 1.25; }
.subtitle { font-size: 20px; font-weight: 400; line-height: 1.5; opacity: 0.7; }
.body-text { font-size: 18px; font-weight: 400; line-height: 1.7; }
.body-small { font-size: 15px; font-weight: 400; line-height: 1.65; }
.label { font-size: 13px; font-weight: 700; letter-spacing: 2.5px; text-transform: uppercase; }
.mono { font-family: 'JetBrains Mono', 'Fira Code', 'Consolas', monospace; }

/* Colors */
.text-white { color: #FFFFFF; }
.text-cream { color: #F5EFE0; }
.text-dark { color: #1A1A2E; }
.text-navy { color: #1F4E79; }
.text-gold { color: #C4972A; }
.text-cyan { color: #00D4FF; }
.text-red { color: #E8453C; }
.text-green { color: #2D8F5E; }
.text-muted { color: rgba(255,255,255,0.45); }
.text-muted-dark { color: rgba(0,0,0,0.4); }

/* Cards */
.card {
  background: rgba(255,255,255,0.06);
  border: 1px solid rgba(255,255,255,0.08);
  border-radius: 12px;
  backdrop-filter: blur(10px);
}
.card-light {
  background: #FFFFFF;
  border: 1px solid rgba(0,0,0,0.06);
  border-radius: 10px;
  box-shadow: 0 2px 12px rgba(0,0,0,0.04);
}
.card-metric {
  background: rgba(255,255,255,0.04);
  border: 1px solid rgba(255,255,255,0.1);
  border-radius: 16px;
  text-align: center;
  padding: 32px 24px;
}

/* Accent elements */
.accent-bar { width: 60px; height: 4px; background: linear-gradient(90deg, #C4972A, #E8C547); border-radius: 2px; }
.accent-line { height: 1px; background: rgba(255,255,255,0.1); }
.accent-line-dark { height: 1px; background: rgba(0,0,0,0.08); }
.gold-border-left { border-left: 4px solid #C4972A; }

/* Key Insight box */
.insight-box {
  background: linear-gradient(135deg, rgba(196,151,42,0.08) 0%, rgba(196,151,42,0.03) 100%);
  border: 1px solid rgba(196,151,42,0.2);
  border-left: 4px solid #C4972A;
  border-radius: 8px;
  padding: 20px 28px;
}

/* Source line */
.source { font-size: 12px; color: #3D6B8E; position: absolute; bottom: 24px; left: 80px; right: 80px; }

/* Nav dots */
.nav-dots { position: absolute; bottom: 28px; right: 80px; display: flex; gap: 6px; }
.dot { width: 8px; height: 8px; border-radius: 50%; background: rgba(255,255,255,0.2); }
.dot.active { background: #C4972A; }

/* Table styling - McKinsey clean */
table { width: 100%; border-collapse: collapse; font-size: 15px; }
th { text-align: left; padding: 12px 16px; font-weight: 700; color: #1F4E79; border-bottom: 2px solid #1F4E79; }
td { padding: 10px 16px; border-bottom: 1px solid #E8E0D0; color: #333; }
tr:last-child td { border-bottom: none; }
"""

def strip_html(text):
    if not text:
        return ""
    text = str(text)
    text = re.sub(r'<[^>]+>', ' ', text)
    text = re.sub(r'\s+', ' ', text)
    return text.strip()

def esc(text):
    return html_escape(strip_html(str(text))) if text else ""

def truncate(text, n=400):
    text = strip_html(text)
    return text[:n-3] + "..." if len(text) > n else text

def first_sentences(text, n=2):
    text = strip_html(text)
    sents = re.split(r'(?<=[.!?])\s+', text)
    result = ' '.join(sents[:n])
    return result[:350] + "..." if len(result) > 350 else result


def make_slide_html(body_html):
    return f"""<!DOCTYPE html>
<html><head><meta charset="utf-8"><style>{CSS_BASE}</style></head>
<body>{body_html}</body></html>"""


# ══════════════════════════════════════════════════════════════
# SLIDE HTML GENERATORS
# ══════════════════════════════════════════════════════════════

def html_title(data):
    title = esc(data.get("title", ""))
    tt = data.get("theoremType", "")
    tn = esc(data.get("theoremName", ""))
    epigraph = esc(data.get("epigraph", ""))
    badge_color = "#E8453C" if tt == "Impossibility" else "#C4972A" if tt == "Intractability" else "#00D4FF"

    return make_slide_html(f"""
    <div class="slide bg-dark" style="display:flex; flex-direction:column; justify-content:center; padding: 80px 100px;">
      <div class="label text-cyan" style="margin-bottom: 20px;">SYSTEM ASSET PRICING MODEL</div>
      <div class="title-hero text-white" style="max-width: 1400px; margin-bottom: 24px;">{title}</div>
      <div style="display:flex; align-items:center; gap: 16px; margin-bottom: 40px;">
        <div style="background:{badge_color}; color:white; padding: 6px 18px; border-radius: 20px; font-size: 14px; font-weight: 700; letter-spacing: 1px;">{esc(tt).upper()} THEOREM</div>
        <div class="body-text text-cream" style="opacity:0.6;">{tn}</div>
      </div>
      {"<div style='max-width:1100px; padding: 24px 32px; border-left: 3px solid rgba(255,255,255,0.15); margin-bottom: 40px;'><div class=&quot;body-text text-cream&quot; style=&quot;font-style:italic; opacity:0.5;&quot;>&ldquo;" + epigraph + "&rdquo;</div></div>" if epigraph else ""}
      <div style="display:flex; align-items:center; gap: 16px; margin-top: auto;">
        <div class="accent-bar"></div>
        <div class="body-small text-muted">Erik Postnieks &nbsp;|&nbsp; Independent Researcher &nbsp;|&nbsp; 2026</div>
      </div>
    </div>""")


def html_metrics(data):
    wa = data.get("welfareAnalysis")
    wa = wa if isinstance(wa, dict) else {}
    beta = wa.get("betaW", "—")
    ci = wa.get("ci90", "—")
    pi_raw = str(wa.get("pi", "—"))
    dw_raw = str(wa.get("deltaW", "—"))

    def short_money(val):
        m = re.search(r'[-−]?\$[\d,.]+\s*(billion|trillion|million)?', val, re.IGNORECASE)
        return m.group(0) if m else val[:40] if val != "—" else "—"

    pi = short_money(pi_raw)
    dw = short_money(dw_raw)
    tt = data.get("theoremType", "")
    badge_color = "#E8453C" if tt == "Impossibility" else "#C4972A"

    return make_slide_html(f"""
    <div class="slide bg-navy" style="padding: 60px 80px;">
      <div class="label text-cyan" style="margin-bottom: 12px;">KEY METRICS</div>
      <div class="title-section text-white" style="margin-bottom: 48px;">Welfare Analysis Dashboard</div>
      <div style="display: grid; grid-template-columns: repeat(4, 1fr); gap: 24px; margin-bottom: 40px;">
        <div class="card-metric">
          <div class="label text-cyan" style="margin-bottom: 12px;">βW</div>
          <div class="mono text-white" style="font-size: 56px; font-weight: 700;">{esc(str(beta))}</div>
          <div class="body-small text-muted" style="margin-top: 8px;">System Beta</div>
        </div>
        <div class="card-metric">
          <div class="label text-gold" style="margin-bottom: 12px;">90% CI</div>
          <div class="mono text-white" style="font-size: 32px; font-weight: 600; padding-top: 12px;">{esc(str(ci))}</div>
          <div class="body-small text-muted" style="margin-top: 8px;">Confidence Interval</div>
        </div>
        <div class="card-metric">
          <div class="label text-green" style="margin-bottom: 12px;">Π (REVENUE)</div>
          <div class="mono text-green" style="font-size: 32px; font-weight: 700; padding-top: 12px;">{esc(pi)}</div>
          <div class="body-small text-muted" style="margin-top: 8px;">Annual Private Payoff</div>
        </div>
        <div class="card-metric">
          <div class="label text-red" style="margin-bottom: 12px;">ΔW (COST)</div>
          <div class="mono text-red" style="font-size: 32px; font-weight: 700; padding-top: 12px;">{esc(dw)}</div>
          <div class="body-small text-muted" style="margin-top: 8px;">Annual Welfare Destruction</div>
        </div>
      </div>
      <div class="insight-box" style="background: rgba(196,151,42,0.06); border-color: rgba(196,151,42,0.15);">
        <div class="label text-gold" style="margin-bottom: 8px;">CLASSIFICATION</div>
        <div class="body-text text-cream">{esc(tt)} Theorem — {"Physical/biological constraints prevent welfare internalization. No policy can solve this." if tt == "Impossibility" else "Institutional constraints are removable. Proven policy models exist." if tt == "Intractability" else "Foundational result underpinning all domain theorems."}</div>
      </div>
      <div class="source" style="color: rgba(255,255,255,0.3);">Source: Monte Carlo simulation (N=10,000, seed=42). Π = annual industry revenue (Iron Law).</div>
    </div>""")


def html_chapter_divider(num, title, subtitle=""):
    return make_slide_html(f"""
    <div class="slide bg-dark" style="display:flex; flex-direction:column; justify-content:center; padding: 80px 100px;">
      <div style="display:flex; align-items:center; gap: 20px; margin-bottom: 24px;">
        <div style="width:64px; height:64px; border-radius:50%; border: 2px solid #C4972A; display:flex; align-items:center; justify-content:center;">
          <span class="mono text-gold" style="font-size: 28px; font-weight: 700;">{num}</span>
        </div>
        <div class="label text-gold">CHAPTER {num}</div>
      </div>
      <div class="title-hero text-white" style="max-width: 1300px;">{esc(title)}</div>
      {"<div class='subtitle text-cream' style='margin-top: 16px; max-width: 1000px;'>" + esc(subtitle) + "</div>" if subtitle else ""}
      <div class="accent-bar" style="margin-top: 40px;"></div>
    </div>""")


def html_text_slide(heading, body, subheading="", page_info=""):
    return make_slide_html(f"""
    <div class="slide bg-cream" style="padding: 0;">
      <div style="background: #1F4E79; padding: 28px 80px 20px;">
        <div class="title-slide text-white">{esc(heading)}</div>
        {"<div class='body-small' style='color:rgba(255,255,255,0.6); margin-top:4px;'>" + esc(subheading) + "</div>" if subheading else ""}
      </div>
      {"<div style='background:#F5EFE0; padding: 4px 80px; font-size:12px; color:#666;'>" + esc(page_info) + "</div>" if page_info else ""}
      <div style="padding: 36px 80px 80px; overflow:hidden;">
        <div class="body-text text-dark" style="max-width: 1600px; column-count: 2; column-gap: 48px; line-height: 1.75;">
          {body}
        </div>
      </div>
    </div>""")


def html_finding_slide(num, total, headline, body):
    return make_slide_html(f"""
    <div class="slide bg-cream" style="padding: 0;">
      <div style="background: #1F4E79; padding: 24px 80px 18px; display:flex; align-items:center; justify-content:space-between;">
        <div class="title-slide text-white">KEY FINDING</div>
        <div class="mono" style="color:rgba(255,255,255,0.4); font-size: 16px;">{num} / {total}</div>
      </div>
      <div style="padding: 48px 80px 80px;">
        <div style="display:flex; gap: 32px; align-items:flex-start;">
          <div style="min-width:80px; width:80px; height:80px; border-radius:50%; background:linear-gradient(135deg, #1F4E79, #2D6FA3); display:flex; align-items:center; justify-content:center;">
            <span class="mono text-white" style="font-size: 32px; font-weight: 700;">{num}</span>
          </div>
          <div style="flex:1;">
            <div class="title-slide text-navy" style="margin-bottom: 20px; line-height: 1.3;">{esc(headline)}</div>
            <div class="body-text text-dark" style="line-height: 1.75; max-width: 1400px;">{esc(truncate(body, 700))}</div>
          </div>
        </div>
      </div>
    </div>""")


def html_agent_slide(agent_name, role, content):
    icons = {"WHISTLEBLOWER": "🔔", "PLAINTIFF": "⚖️", "REGULATOR": "🏛️",
             "LEGISLATOR": "📜", "INVESTOR": "📊", "SUPRANATIONAL": "🌍"}
    icon = icons.get(agent_name, "●")
    return make_slide_html(f"""
    <div class="slide bg-cream" style="padding: 0;">
      <div style="background: #1F4E79; padding: 24px 80px 18px;">
        <div class="label" style="color:rgba(255,255,255,0.5); margin-bottom:4px;">SIX-AGENT CONFLICTORING FRAMEWORK</div>
        <div class="title-slide text-white">{icon} {esc(agent_name)}</div>
      </div>
      <div style="background: #F5EFE0; padding: 6px 80px;">
        <div class="body-small text-navy" style="font-weight:600;">{esc(role)}</div>
      </div>
      <div style="padding: 36px 80px 80px;">
        <div class="body-text text-dark" style="line-height: 1.75; column-count: 2; column-gap: 48px;">
          {esc(truncate(content, 1200))}
        </div>
      </div>
    </div>""")


def html_quiz_slide(data):
    wa = data.get("welfareAnalysis")
    wa = wa if isinstance(wa, dict) else {}
    beta = wa.get("betaW", "?")
    tt = data.get("theoremType", "?")
    tn = esc(data.get("theoremName", "?"))

    return make_slide_html(f"""
    <div class="slide bg-navy" style="padding: 60px 80px;">
      <div style="display:flex; align-items:center; gap:16px; margin-bottom: 40px;">
        <div style="width:48px; height:48px; border-radius:50%; background: #C4972A; display:flex; align-items:center; justify-content:center;">
          <span style="font-size:24px;">?</span>
        </div>
        <div class="title-section text-white">Comprehension Check</div>
      </div>
      <div style="display: grid; grid-template-columns: 1fr 1fr; gap: 24px;">
        <div class="card" style="padding: 28px 32px;">
          <div class="label text-cyan" style="margin-bottom: 12px;">QUESTION 1</div>
          <div class="body-text text-white" style="margin-bottom: 16px;">What is the system beta (βW) for this domain?</div>
          <div class="mono text-cyan" style="font-size: 28px; font-weight: 700;">{esc(str(beta))}</div>
        </div>
        <div class="card" style="padding: 28px 32px;">
          <div class="label text-cyan" style="margin-bottom: 12px;">QUESTION 2</div>
          <div class="body-text text-white" style="margin-bottom: 16px;">Impossibility or Intractability?</div>
          <div class="mono text-gold" style="font-size: 24px; font-weight: 700;">{esc(tt)}</div>
        </div>
        <div class="card" style="padding: 28px 32px;">
          <div class="label text-cyan" style="margin-bottom: 12px;">QUESTION 3</div>
          <div class="body-text text-white" style="margin-bottom: 16px;">What is the theorem called?</div>
          <div class="text-white" style="font-size: 20px; font-weight: 600;">{tn}</div>
        </div>
        <div class="card" style="padding: 28px 32px;">
          <div class="label text-cyan" style="margin-bottom: 12px;">QUESTION 4</div>
          <div class="body-text text-white" style="margin-bottom: 16px;">Can policy alone solve this?</div>
          <div class="text-gold" style="font-size: 20px; font-weight: 600;">{"No — physical law prevents it" if tt == "Impossibility" else "Yes — proven models exist" if tt == "Intractability" else "See theorem"}</div>
        </div>
      </div>
    </div>""")


def html_closing(data):
    wa = data.get("welfareAnalysis")
    wa = wa if isinstance(wa, dict) else {}
    beta = wa.get("betaW", "")
    tt = data.get("theoremType", "")
    title = esc(data.get("title", ""))
    badge_color = "#E8453C" if tt == "Impossibility" else "#C4972A"

    return make_slide_html(f"""
    <div class="slide bg-dark" style="display:flex; flex-direction:column; align-items:center; justify-content:center; text-align:center; padding: 80px;">
      <div class="title-section text-white" style="max-width: 1200px; margin-bottom: 32px;">{title}</div>
      <div style="height:2px; width:200px; background: linear-gradient(90deg, transparent, #00D4FF, transparent); margin-bottom: 32px;"></div>
      {"<div class='mono text-cyan' style='font-size: 72px; font-weight: 700; margin-bottom: 8px;'>βW = " + esc(str(beta)) + "</div>" if beta else ""}
      {"<div style='display:inline-block; background:" + badge_color + "; color:white; padding: 6px 24px; border-radius: 20px; font-size: 16px; font-weight: 700; letter-spacing: 1px; margin-bottom: 40px;'>" + esc(tt).upper() + " THEOREM</div>" if tt else ""}
      <div style="margin-top: 40px;">
        <div class="body-text text-muted" style="margin-bottom: 8px;">Erik Postnieks &nbsp;|&nbsp; Independent Researcher &nbsp;|&nbsp; postnieks.com</div>
        <div class="body-small text-muted">Full paper and Monte Carlo replication: github.com/epostnieks</div>
      </div>
      <div class="label text-cyan" style="margin-top: 32px; opacity: 0.5;">END OF MASTERCLASS</div>
    </div>""")


# ══════════════════════════════════════════════════════════════
# RENDERING
# ══════════════════════════════════════════════════════════════

def render_slides_to_pptx(html_slides, output_path, pw_browser):
    """Render list of HTML strings to PNG screenshots, assemble into PPTX."""
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    context = pw_browser.new_context(viewport={"width": 1920, "height": 1080})
    page = context.new_page()

    for i, html in enumerate(html_slides):
        page.set_content(html, wait_until="networkidle")
        png_bytes = page.screenshot(type="png")

        slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
        # Full-bleed image
        slide.shapes.add_picture(
            BytesIO(png_bytes),
            Emu(0), Emu(0),
            prs.slide_width, prs.slide_height
        )

    context.close()
    prs.save(str(output_path))
    return len(html_slides)


def build_paper_slides(data):
    """Build all HTML slides for a paper."""
    slides = []

    # 1. Title
    slides.append(html_title(data))

    # 2. Metrics
    slides.append(html_metrics(data))

    # 3. Executive Summary
    es = strip_html(data.get("executiveSummary", ""))
    if es:
        chunks = []
        sents = re.split(r'(?<=[.!?])\s+', es)
        current = ""
        for s in sents:
            if len(current) + len(s) > 600 and current:
                chunks.append(current.strip())
                current = s
            else:
                current = (current + " " + s).strip()
        if current:
            chunks.append(current.strip())

        slides.append(html_chapter_divider(1, "EXECUTIVE SUMMARY", "The complete argument"))
        for i, chunk in enumerate(chunks[:4]):
            slides.append(html_text_slide(
                "EXECUTIVE SUMMARY",
                esc(chunk),
                page_info=f"Page {i+1} of {min(len(chunks),4)}"
            ))

    # 4. Key Findings
    findings = data.get("keyFindings", [])
    if findings:
        slides.append(html_chapter_divider(2, "KEY FINDINGS", f"{len(findings)} evidence-based results"))
        for i, f in enumerate(findings[:12]):
            text = strip_html(f) if isinstance(f, str) else str(f)
            sents = re.split(r'(?<=[.!?])\s+', text)
            headline = sents[0] if sents else text
            body = ' '.join(sents[1:]) if len(sents) > 1 else ""
            slides.append(html_finding_slide(i+1, min(len(findings), 12), headline, body))

    # 5. Theorem
    theorem = data.get("theorem", {})
    if isinstance(theorem, dict) and theorem:
        slides.append(html_chapter_divider(3, "THE THEOREM", data.get("theoremName", "")))

        plain = strip_html(theorem.get("plain", ""))
        if plain:
            chunks = []
            sents = re.split(r'(?<=[.!?])\s+', plain)
            current = ""
            for s in sents:
                if len(current) + len(s) > 600 and current:
                    chunks.append(current.strip())
                    current = s
                else:
                    current = (current + " " + s).strip()
            if current:
                chunks.append(current.strip())
            for i, chunk in enumerate(chunks[:3]):
                slides.append(html_text_slide("PLAIN LANGUAGE", esc(chunk),
                    subheading=data.get("theoremName", ""),
                    page_info=f"Page {i+1} of {min(len(chunks),3)}"))

        formal = strip_html(theorem.get("formal", ""))
        if formal:
            chunks = []
            sents = re.split(r'(?<=[.!?])\s+', formal)
            current = ""
            for s in sents:
                if len(current) + len(s) > 600 and current:
                    chunks.append(current.strip())
                    current = s
                else:
                    current = (current + " " + s).strip()
            if current:
                chunks.append(current.strip())
            for i, chunk in enumerate(chunks[:3]):
                slides.append(html_text_slide("FORMAL STATEMENT", esc(chunk),
                    page_info=f"Page {i+1} of {min(len(chunks),3)}"))

    # 6. Welfare Channels
    wa = data.get("welfareAnalysis")
    wa = wa if isinstance(wa, dict) else {}
    channels = wa.get("channels", [])
    if channels:
        slides.append(html_chapter_divider(4, "WELFARE CHANNELS", "Where the damage comes from"))
        for ch in channels[:6]:
            name = strip_html(ch.get("name", ""))
            desc = strip_html(ch.get("description", ""))
            value = strip_html(str(ch.get("value", "")))
            if desc:
                slides.append(html_text_slide(f"CHANNEL: {name}", esc(truncate(desc, 1200)),
                    subheading=value))

    # 7. Case Studies
    cases = data.get("caseStudies", [])
    if cases:
        slides.append(html_chapter_divider(5, "CASE STUDIES", "Real-world evidence"))
        for cs in cases[:5]:
            title = strip_html(cs.get("title", ""))
            content = strip_html(cs.get("content", ""))
            if content:
                slides.append(html_text_slide(f"CASE: {title}", esc(truncate(content, 1200))))

    # 8. Policy
    policy = data.get("policyAnalysis")
    if isinstance(policy, dict):
        slides.append(html_chapter_divider(6, "POLICY ANALYSIS", "What works, what fails, and why"))
        for label, key in [("CURRENT FRAMEWORK", "currentFramework"), ("WHY IT FAILS", "failures"), ("REFORM PATHWAY", "reform")]:
            text = strip_html(policy.get(key, ""))
            if text:
                slides.append(html_text_slide(label, esc(truncate(text, 1200))))

    # 9. Six-Agent Framework
    agents = data.get("agents", {})
    if isinstance(agents, dict) and agents:
        slides.append(html_chapter_divider(7, "SIX-AGENT FRAMEWORK", "Who can break the Private-Systemic Tension"))
        agent_roles = {
            "insider": ("WHISTLEBLOWER", "Break Information Asymmetry"),
            "whistleblower": ("WHISTLEBLOWER", "Break Information Asymmetry"),
            "plaintiff": ("PLAINTIFF", "Monetize Liability"),
            "regulator": ("REGULATOR", "Redesign the Game"),
            "legislator": ("LEGISLATOR", "Change the Law"),
            "investor": ("INVESTOR", "Reprice Capital"),
            "supranational": ("SUPRANATIONAL", "Coordinate Jurisdictions"),
        }
        for key in ["insider", "whistleblower", "plaintiff", "regulator", "legislator", "investor", "supranational"]:
            content = agents.get(key)
            if not content:
                continue
            name, role = agent_roles.get(key, (key.upper(), ""))
            slides.append(html_agent_slide(name, role, str(content)))

    # 10. Monte Carlo
    mc = wa.get("monteCarlo", "")
    if mc:
        slides.append(html_chapter_divider(8, "MONTE CARLO SIMULATION", "Quantifying uncertainty"))
        text = strip_html(mc)
        chunks = []
        sents = re.split(r'(?<=[.!?])\s+', text)
        current = ""
        for s in sents:
            if len(current) + len(s) > 600 and current:
                chunks.append(current.strip())
                current = s
            else:
                current = (current + " " + s).strip()
        if current:
            chunks.append(current.strip())
        for i, chunk in enumerate(chunks[:2]):
            slides.append(html_text_slide("MONTE CARLO METHODOLOGY", esc(chunk),
                subheading="N=10,000 | Seed=42 | Reproducible",
                page_info=f"Page {i+1} of {min(len(chunks),2)}"))

    # 11. FAQ
    faq = data.get("faq", [])
    if faq:
        slides.append(html_chapter_divider(9, "SKEPTIC'S FAQ", "Hard questions, evidence-based answers"))
        for item in faq[:6]:
            q = strip_html(item.get("q", ""))
            a = strip_html(item.get("a", ""))
            if q and a:
                slides.append(html_text_slide(f"Q: {q}", esc(truncate(a, 1000))))

    # 12. Quiz
    slides.append(html_quiz_slide(data))

    # 13. Closing
    slides.append(html_closing(data))

    return slides


def generate_one(json_path, pw_browser):
    with open(json_path) as f:
        data = json.load(f)
    slug = data.get("slug", json_path.stem)
    slides = build_paper_slides(data)
    out = OUTPUT_DIR / f"{slug}.pptx"
    count = render_slides_to_pptx(slides, out, pw_browser)
    return slug, count


def main():
    json_files = sorted(PAPER_DATA_DIR.glob("*.json"))
    print(f"Found {len(json_files)} papers. Rendering with Playwright...")

    with sync_playwright() as p:
        browser = p.chromium.launch(headless=True)

        results = []
        errors = []
        for jf in json_files:
            try:
                slug, count = generate_one(jf, browser)
                results.append((slug, count))
                print(f"  {slug}: {count} slides")
            except Exception as e:
                errors.append((jf.stem, str(e)[:120]))
                print(f"  ERROR {jf.stem}: {str(e)[:120]}")

        browser.close()

    print(f"\n{'='*60}")
    print(f"Generated: {len(results)} decks")
    if errors:
        print(f"Errors: {len(errors)}")
        for n, e in errors:
            print(f"  {n}: {e}")
    total = sum(c for _, c in results)
    print(f"Total slides: {total}")
    print(f"Output: {OUTPUT_DIR}")


if __name__ == "__main__":
    main()
