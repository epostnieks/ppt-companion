#!/usr/bin/env python3
"""
PPTX Generator V2 — Masterclass-quality decks.
- Extracts tables from Word docs as slide content
- Uses full 20K-word Gemini summaries (all sections)
- 40-60 slides per paper, chapter-structured
- McKinsey formatting: Navy headers, cream accents, Liberation Sans
- Quiz slides, chapter summaries, metric cards
"""

import json
import os
import re
import sys
import glob
import copy
from concurrent.futures import ThreadPoolExecutor, as_completed
from pathlib import Path
from html.parser import HTMLParser
from io import BytesIO

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# Try to import docx for table extraction
try:
    from docx import Document as DocxDocument
    HAS_DOCX = True
except ImportError:
    HAS_DOCX = False

# ── Paths ──
PAPER_DATA_DIR = Path(__file__).parent.parent / "src" / "paperData"
OUTPUT_DIR = Path(__file__).parent.parent / "pptx_output"
OUTPUT_DIR.mkdir(exist_ok=True)
BASE = "/Users/erikpostnieks/Projects"

# ── McKinsey Colors ──
NAVY = RGBColor(0x1F, 0x4E, 0x79)          # Title bars
CREAM = RGBColor(0xFB, 0xF9, 0xF4)         # Subtitle bars
TAN = RGBColor(0xF5, 0xEF, 0xE0)           # KEY INSIGHT box bg
GOLD_BORDER = RGBColor(0xC4, 0x97, 0x2A)   # KEY INSIGHT left border
SOURCE_BLUE = RGBColor(0x3D, 0x6B, 0x8E)   # Source line
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x00, 0x00, 0x00)
DARK_TEXT = RGBColor(0x33, 0x33, 0x33)
LIGHT_GRAY = RGBColor(0xE8, 0xE0, 0xD0)    # Table row dividers
MID_GRAY = RGBColor(0x88, 0x88, 0x88)
RED = RGBColor(0xCC, 0x33, 0x33)
GREEN = RGBColor(0x22, 0x88, 0x55)
SLIDE_BG = RGBColor(0xFF, 0xFF, 0xFF)       # White background
DARK_BG = RGBColor(0x0F, 0x17, 0x2A)        # Dark slides (title, closing)
ACCENT_CYAN = RGBColor(0x00, 0xD4, 0xFF)
SURFACE = RGBColor(0xF7, 0xF5, 0xF0)        # Light surface

FONT = "Liberation Sans"

# ── HTML stripping ──
class HTMLStripper(HTMLParser):
    def __init__(self):
        super().__init__()
        self.result = []
    def handle_data(self, d):
        self.result.append(d)
    def handle_starttag(self, tag, attrs):
        if tag == 'td' or tag == 'th':
            self.result.append(' | ')
        if tag == 'tr':
            self.result.append('\n')
    def get_text(self):
        return ''.join(self.result).strip()

def strip_html(text):
    if not text or not isinstance(text, str):
        return str(text) if text else ""
    s = HTMLStripper()
    s.feed(text)
    return s.get_text()

def truncate(text, max_chars=600):
    text = strip_html(text)
    if len(text) <= max_chars:
        return text
    return text[:max_chars-3] + "..."

def split_into_chunks(text, max_chars=700):
    """Split text into presentation-sized chunks at sentence boundaries."""
    text = strip_html(text)
    sentences = re.split(r'(?<=[.!?])\s+', text)
    chunks = []
    current = ""
    for s in sentences:
        if len(current) + len(s) > max_chars and current:
            chunks.append(current.strip())
            current = s
        else:
            current = (current + " " + s).strip()
    if current:
        chunks.append(current.strip())
    return chunks

def set_slide_bg(slide, color=SLIDE_BG):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_rect(slide, left, top, width, height, color):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    shape.shadow.inherit = False
    return shape

def add_text(slide, left, top, width, height, text, size=14,
             color=DARK_TEXT, bold=False, italic=False, align=PP_ALIGN.LEFT, font=FONT):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.italic = italic
    p.font.name = font
    p.alignment = align
    return txBox

def add_multiline(slide, left, top, width, height, text, size=12,
                  color=DARK_TEXT, line_spacing=1.5, font=FONT):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    paragraphs = text.split('\n')
    for i, para_text in enumerate(paragraphs):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = para_text
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.font.name = font
        p.space_after = Pt(size * 0.5)
    return txBox

def mckinsey_header(slide, title, subtitle=None):
    """Five-layer McKinsey header: navy title bar + cream subtitle."""
    # Navy title bar
    add_rect(slide, 0, 0, Inches(13.333), Inches(0.9), NAVY)
    add_text(slide, Inches(0.6), Inches(0.15), Inches(12), Inches(0.6),
             title, size=26, color=WHITE, bold=True)

    if subtitle:
        add_rect(slide, 0, Inches(0.9), Inches(13.333), Inches(0.5), CREAM)
        add_text(slide, Inches(0.6), Inches(0.95), Inches(12), Inches(0.4),
                 subtitle, size=14, color=DARK_TEXT, italic=True)

def key_insight_box(slide, text, y_pos=Inches(5.8)):
    """KEY INSIGHT box with gold left border on tan background."""
    # Tan background
    add_rect(slide, Inches(0.5), y_pos, Inches(12.3), Inches(1.2), TAN)
    # Gold left border (thin rect)
    add_rect(slide, Inches(0.5), y_pos, Inches(0.06), Inches(1.2), GOLD_BORDER)
    # Label
    add_text(slide, Inches(0.8), y_pos + Inches(0.05), Inches(2), Inches(0.3),
             "KEY INSIGHT", size=11, color=GOLD_BORDER, bold=True)
    # Text
    add_text(slide, Inches(0.8), y_pos + Inches(0.3), Inches(11.5), Inches(0.8),
             truncate(text, 300), size=13, color=DARK_TEXT)

def source_line(slide, text, y_pos=Inches(7.1)):
    """Source line at bottom."""
    add_text(slide, Inches(0.6), y_pos, Inches(12), Inches(0.3),
             f"Source: {text}", size=11, color=SOURCE_BLUE)


# ══════════════════════════════════════════════════════════════
# WORD DOC TABLE EXTRACTION
# ══════════════════════════════════════════════════════════════

def find_docx(slug):
    """Find the Word document for a paper."""
    patterns = [
        f"{BASE}/sapm-{slug}/paper/{slug}-current.docx",
        f"{BASE}/sapm-{slug}/paper/*-current.docx",
        f"{BASE}/sapm-{slug}/paper/*-final.docx",
        f"{BASE}/sapm-{slug}/paper/*.docx",
    ]
    for pat in patterns:
        matches = glob.glob(pat)
        matches = [m for m in matches if '~' not in m and 'template' not in m.lower()]
        if matches:
            return matches[0]
    return None

def extract_docx_tables(docx_path, max_tables=20):
    """Extract tables from a Word document as structured data."""
    if not HAS_DOCX or not docx_path:
        return []

    try:
        doc = DocxDocument(docx_path)
    except Exception:
        return []

    tables = []
    for i, table in enumerate(doc.tables[:max_tables]):
        rows = []
        for row in table.rows:
            cells = [cell.text.strip() for cell in row.cells]
            rows.append(cells)
        if rows and len(rows) > 1:  # Need header + at least one data row
            # Try to determine a title from the paragraph before the table
            title = f"Table {i+1}"
            tables.append({"title": title, "header": rows[0], "rows": rows[1:]})

    return tables


# ══════════════════════════════════════════════════════════════
# SLIDE BUILDERS
# ══════════════════════════════════════════════════════════════

def slide_title(prs, data):
    """Title slide — dark background, paper title, epigraph."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)

    # Thin accent line
    add_rect(slide, 0, 0, Inches(13.333), Inches(0.06), ACCENT_CYAN)

    # SAPM label
    add_text(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(0.4),
             "SYSTEM ASSET PRICING MODEL — MASTERCLASS", size=14, color=ACCENT_CYAN, bold=True)

    # Title
    title = data.get("title", data.get("slug", ""))
    add_text(slide, Inches(0.8), Inches(1.3), Inches(11.5), Inches(2.5),
             title, size=36, color=WHITE, bold=True)

    # Theorem badge
    tt = data.get("theoremType", "")
    tn = data.get("theoremName", "")
    if tt:
        badge_color = RED if tt == "Impossibility" else GOLD_BORDER
        add_text(slide, Inches(0.8), Inches(4.0), Inches(6), Inches(0.5),
                 f"{tt} Theorem: {tn}", size=18, color=badge_color, bold=True)

    # Epigraph
    epigraph = data.get("epigraph", "")
    if epigraph:
        add_text(slide, Inches(0.8), Inches(5.2), Inches(10), Inches(1.0),
                 f'"{epigraph}"', size=16, color=MID_GRAY, italic=True)

    # Author
    add_text(slide, Inches(0.8), Inches(6.6), Inches(10), Inches(0.4),
             "Erik Postnieks  |  Independent Researcher  |  2026", size=12, color=MID_GRAY)


def slide_chapter_divider(prs, chapter_num, chapter_title, subtitle=""):
    """Chapter divider slide — navy background."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, NAVY)

    add_text(slide, Inches(0.8), Inches(2.0), Inches(11), Inches(0.6),
             f"CHAPTER {chapter_num}", size=16, color=ACCENT_CYAN, bold=True)
    add_text(slide, Inches(0.8), Inches(2.8), Inches(11), Inches(1.5),
             chapter_title, size=36, color=WHITE, bold=True)
    if subtitle:
        add_text(slide, Inches(0.8), Inches(4.5), Inches(10), Inches(1.0),
                 subtitle, size=16, color=CREAM, italic=True)


def slides_metrics(prs, data):
    """Key metrics dashboard — 1-2 slides."""
    wa = data.get("welfareAnalysis")
    wa = wa if isinstance(wa, dict) else {}

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    mckinsey_header(slide, "KEY METRICS DASHBOARD", "System Beta and Welfare Analysis")

    beta = wa.get("betaW", data.get("betaW", "—"))
    ci = wa.get("ci90", data.get("ci90", "—"))
    pi_raw = wa.get("pi", "—")
    dw_raw = wa.get("deltaW", "—")

    def short_money(val):
        if not val or val == "—":
            return "—"
        val = str(val)
        m = re.search(r'[-−]?\$[\d,.]+\s*(billion|trillion|million)?', val, re.IGNORECASE)
        return m.group(0) if m else val[:50]

    metrics = [
        ("βW (System Beta)", str(beta), "Welfare destruction per dollar of private gain", NAVY),
        ("90% Confidence Interval", str(ci), "Monte Carlo simulation bounds", DARK_TEXT),
        ("Π (Annual Revenue)", short_money(pi_raw), "Private payoff — always revenue, never profit", GREEN),
        ("ΔW (Welfare Cost)", short_money(dw_raw), "Annual system welfare destruction", RED),
    ]

    y_start = Inches(1.8)
    card_w = Inches(5.8)
    card_h = Inches(1.1)

    for i, (label, value, desc, accent) in enumerate(metrics):
        col = i % 2
        row = i // 2
        x = Inches(0.5) + col * (card_w + Inches(0.4))
        y = y_start + row * (card_h + Inches(0.3))

        add_rect(slide, x, y, card_w, card_h, SURFACE)
        add_text(slide, x + Inches(0.2), y + Inches(0.08), card_w - Inches(0.4), Inches(0.3),
                 label, size=12, color=MID_GRAY, bold=True)
        val_size = 28 if len(value) < 15 else 20 if len(value) < 30 else 14
        add_text(slide, x + Inches(0.2), y + Inches(0.35), card_w - Inches(0.4), Inches(0.5),
                 value, size=val_size, color=accent, bold=True)
        add_text(slide, x + Inches(0.2), y + Inches(0.8), card_w - Inches(0.4), Inches(0.25),
                 desc, size=10, color=MID_GRAY)

    # Theorem type
    tt = data.get("theoremType", "")
    if tt:
        key_insight_box(slide,
            f"Classification: {tt} Theorem — {'Physical/biological constraints prevent welfare internalization' if tt == 'Impossibility' else 'Institutional constraints are removable with proven policy reform'}",
            y_pos=Inches(5.0))

    source_line(slide, "Monte Carlo simulation (N=10,000, seed=42). Π = annual industry revenue (Iron Law).")


def slides_executive_summary(prs, data):
    """Executive summary — multi-slide, full text."""
    es = data.get("executiveSummary", "")
    if not es:
        return

    text = strip_html(es)
    chunks = split_into_chunks(text, 900)

    for i, chunk in enumerate(chunks[:5]):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        set_slide_bg(slide)
        suffix = f" ({i+1}/{min(len(chunks),5)})" if len(chunks) > 1 else ""
        mckinsey_header(slide, f"EXECUTIVE SUMMARY{suffix}",
                       "What this paper proves and why it matters")

        add_multiline(slide, Inches(0.6), Inches(1.6), Inches(12), Inches(5.0),
                      chunk, size=14, color=DARK_TEXT, line_spacing=1.6)

        if i == 0 and len(chunks) > 1:
            source_line(slide, f"Paper: {data.get('title', '')}")


def slides_key_findings(prs, data):
    """Key findings — one finding per slide with full text."""
    findings = data.get("keyFindings", [])
    if not findings:
        return

    for i, finding in enumerate(findings[:12]):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        set_slide_bg(slide)
        mckinsey_header(slide, f"KEY FINDING {i+1} OF {min(len(findings), 12)}")

        text = strip_html(finding) if isinstance(finding, str) else str(finding)

        # Extract first sentence as the "headline"
        sentences = re.split(r'(?<=[.!?])\s+', text)
        headline = sentences[0] if sentences else ""
        body = ' '.join(sentences[1:]) if len(sentences) > 1 else ""

        # Bold headline
        add_text(slide, Inches(0.6), Inches(1.3), Inches(12), Inches(0.8),
                 headline, size=18, color=NAVY, bold=True)

        # Full body text
        if body:
            add_multiline(slide, Inches(0.6), Inches(2.3), Inches(12), Inches(3.5),
                          truncate(body, 900), size=13, color=DARK_TEXT)

        # Key insight from the finding
        if len(headline) > 20:
            key_insight_box(slide, headline, y_pos=Inches(6.0))


def slides_theorem(prs, data):
    """Theorem — formal, plain language, axioms, proof sketch. Multiple slides."""
    theorem = data.get("theorem", {})
    if not theorem or not isinstance(theorem, dict):
        return

    # Slide 1: Plain language
    plain = theorem.get("plain", "")
    if plain:
        chunks = split_into_chunks(strip_html(plain), 900)
        for i, chunk in enumerate(chunks[:3]):
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            set_slide_bg(slide)
            suffix = f" — PLAIN LANGUAGE ({i+1}/{min(len(chunks),3)})" if len(chunks) > 1 else " — PLAIN LANGUAGE"
            mckinsey_header(slide, f"THE THEOREM{suffix}",
                           data.get("theoremName", ""))
            add_multiline(slide, Inches(0.6), Inches(1.6), Inches(12), Inches(5.0),
                          chunk, size=14, color=DARK_TEXT)

    # Slide 2: Axioms
    axioms = theorem.get("axioms", [])
    if axioms:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        set_slide_bg(slide)
        mckinsey_header(slide, "AXIOMS", "The structural assumptions that drive the result")

        for j, ax in enumerate(axioms[:5]):
            y = Inches(1.6) + j * Inches(1.1)
            add_rect(slide, Inches(0.5), y, Inches(12.3), Inches(0.95), SURFACE)

            ax_id = ax.get("id", f"A{j+1}")
            ax_name = ax.get("name", "")
            add_text(slide, Inches(0.7), y + Inches(0.05), Inches(3), Inches(0.3),
                     f"{ax_id}: {ax_name}", size=14, color=NAVY, bold=True)

            statement = strip_html(ax.get("statement", ""))
            add_text(slide, Inches(0.7), y + Inches(0.35), Inches(11.5), Inches(0.55),
                     truncate(statement, 250), size=11, color=DARK_TEXT)

    # Slide 3+: Formal statement
    formal = theorem.get("formal", "")
    if formal:
        chunks = split_into_chunks(strip_html(formal), 900)
        for i, chunk in enumerate(chunks[:3]):
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            set_slide_bg(slide)
            suffix = f" ({i+1}/{min(len(chunks),3)})" if len(chunks) > 1 else ""
            mckinsey_header(slide, f"FORMAL THEOREM STATEMENT{suffix}")
            add_multiline(slide, Inches(0.6), Inches(1.4), Inches(12), Inches(5.5),
                          chunk, size=12, color=DARK_TEXT, font=FONT)

    # Slide 4+: Proof sketch
    proof = theorem.get("proofSketch", "")
    if proof:
        chunks = split_into_chunks(strip_html(proof), 900)
        for i, chunk in enumerate(chunks[:4]):
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            set_slide_bg(slide)
            suffix = f" ({i+1}/{min(len(chunks),4)})" if len(chunks) > 1 else ""
            mckinsey_header(slide, f"PROOF SKETCH{suffix}",
                           "Step-by-step construction of the impossibility result")
            add_multiline(slide, Inches(0.6), Inches(1.6), Inches(12), Inches(5.0),
                          chunk, size=13, color=DARK_TEXT)


def slides_welfare_channels(prs, data):
    """Welfare channel breakdown — one slide per channel."""
    wa = data.get("welfareAnalysis")
    wa = wa if isinstance(wa, dict) else {}
    channels = wa.get("channels", [])
    if not channels:
        return

    # Overview slide
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    mckinsey_header(slide, "WELFARE CHANNEL DECOMPOSITION",
                   "How the total welfare cost breaks down by channel")

    for j, ch in enumerate(channels[:8]):
        col = j % 2
        row = j // 2
        x = Inches(0.5) + col * Inches(6.4)
        y = Inches(1.6) + row * Inches(1.3)

        add_rect(slide, x, y, Inches(6.1), Inches(1.15), SURFACE)
        name = ch.get("name", f"Channel {j+1}")
        value = strip_html(str(ch.get("value", "—")))
        add_text(slide, x + Inches(0.15), y + Inches(0.05), Inches(4), Inches(0.3),
                 name, size=13, color=NAVY, bold=True)
        add_text(slide, x + Inches(4.2), y + Inches(0.05), Inches(1.7), Inches(0.3),
                 value, size=13, color=RED, bold=True, align=PP_ALIGN.RIGHT)

        desc = strip_html(ch.get("description", ""))
        first_sent = desc.split('. ')[0] + '.' if '. ' in desc else desc
        add_text(slide, x + Inches(0.15), y + Inches(0.4), Inches(5.8), Inches(0.7),
                 truncate(first_sent, 200), size=10, color=MID_GRAY)

    # Individual channel slides with full detail
    for j, ch in enumerate(channels[:6]):
        desc = strip_html(ch.get("description", ""))
        if len(desc) < 100:
            continue
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        set_slide_bg(slide)
        name = ch.get("name", f"Channel {j+1}")
        value = strip_html(str(ch.get("value", "")))
        mckinsey_header(slide, f"CHANNEL: {name.upper()}", value)
        add_multiline(slide, Inches(0.6), Inches(1.6), Inches(12), Inches(5.0),
                      truncate(desc, 1200), size=13, color=DARK_TEXT)


def slides_case_studies(prs, data):
    """Case studies — one slide per case."""
    cases = data.get("caseStudies", [])
    if not cases:
        return

    for i, case in enumerate(cases[:6]):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        set_slide_bg(slide)
        title = case.get("title", f"Case Study {i+1}")
        mckinsey_header(slide, f"CASE STUDY: {title.upper()}")

        content = strip_html(case.get("content", ""))
        add_multiline(slide, Inches(0.6), Inches(1.4), Inches(12), Inches(5.2),
                      truncate(content, 1200), size=13, color=DARK_TEXT)


def slides_policy(prs, data):
    """Policy analysis — current framework, failures, reform. 3 slides."""
    policy = data.get("policyAnalysis")
    if not policy or not isinstance(policy, dict):
        return

    sections = [
        ("CURRENT REGULATORY FRAMEWORK", policy.get("currentFramework", "")),
        ("WHY THE CURRENT FRAMEWORK FAILS", policy.get("failures", "")),
        ("REFORM PATHWAY", policy.get("reform", "")),
    ]

    for title, content in sections:
        if not content:
            continue
        text = strip_html(content)
        chunks = split_into_chunks(text, 900)
        for i, chunk in enumerate(chunks[:2]):
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            set_slide_bg(slide)
            suffix = f" ({i+1}/{min(len(chunks),2)})" if len(chunks) > 1 else ""
            mckinsey_header(slide, f"{title}{suffix}")
            add_multiline(slide, Inches(0.6), Inches(1.4), Inches(12), Inches(5.3),
                          chunk, size=13, color=DARK_TEXT)


def slides_agents(prs, data):
    """Six-agent framework — one slide per agent with full detail."""
    agents = data.get("agents", {})
    if not agents or not isinstance(agents, dict):
        return

    agent_meta = {
        "insider": ("WHISTLEBLOWER", "Break Information Asymmetry"),
        "whistleblower": ("WHISTLEBLOWER", "Break Information Asymmetry"),
        "plaintiff": ("PLAINTIFF", "Monetize Liability"),
        "regulator": ("REGULATOR", "Redesign the Game"),
        "legislator": ("LEGISLATOR", "Change the Law"),
        "investor": ("INVESTOR", "Reprice Capital"),
        "supranational": ("SUPRANATIONAL", "Coordinate Jurisdictions"),
    }

    for key in ["insider", "whistleblower", "plaintiff", "regulator", "legislator", "investor", "supranational"]:
        content = agents.get(key)
        if not content:
            continue
        label, subtitle = agent_meta.get(key, (key.upper(), ""))
        text = strip_html(str(content))
        chunks = split_into_chunks(text, 900)

        for i, chunk in enumerate(chunks[:2]):
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            set_slide_bg(slide)
            suffix = f" ({i+1}/{min(len(chunks),2)})" if len(chunks) > 1 else ""
            mckinsey_header(slide, f"AGENT: {label}{suffix}", subtitle)
            add_multiline(slide, Inches(0.6), Inches(1.6), Inches(12), Inches(5.0),
                          chunk, size=13, color=DARK_TEXT)


def slides_cross_domain(prs, data):
    """Cross-domain connections."""
    connections = data.get("crossDomainConnections", [])
    if not connections:
        return

    # 2 connections per slide
    for page in range(0, min(len(connections), 8), 2):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        set_slide_bg(slide)
        mckinsey_header(slide, "CROSS-DOMAIN CONNECTIONS",
                       "How this domain links to others in the SAPM framework")

        batch = connections[page:page+2]
        for j, conn in enumerate(batch):
            y = Inches(1.6) + j * Inches(2.8)
            domain = conn.get("domain", "")
            desc = strip_html(conn.get("connection", ""))

            add_rect(slide, Inches(0.5), y, Inches(12.3), Inches(2.5), SURFACE)
            add_text(slide, Inches(0.7), y + Inches(0.1), Inches(11.5), Inches(0.4),
                     domain, size=16, color=NAVY, bold=True)
            add_multiline(slide, Inches(0.7), y + Inches(0.5), Inches(11.5), Inches(1.8),
                          truncate(desc, 600), size=12, color=DARK_TEXT)


def slides_monte_carlo(prs, data):
    """Monte Carlo methodology."""
    wa = data.get("welfareAnalysis")
    wa = wa if isinstance(wa, dict) else {}
    mc = wa.get("monteCarlo", "")
    if not mc:
        return

    text = strip_html(mc)
    chunks = split_into_chunks(text, 900)

    for i, chunk in enumerate(chunks[:3]):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        set_slide_bg(slide)
        suffix = f" ({i+1}/{min(len(chunks),3)})" if len(chunks) > 1 else ""
        mckinsey_header(slide, f"MONTE CARLO SIMULATION{suffix}",
                       "N=10,000 draws | Seed=42 | Reproducible")
        add_multiline(slide, Inches(0.6), Inches(1.6), Inches(12), Inches(5.0),
                      chunk, size=13, color=DARK_TEXT)


def slides_faq(prs, data):
    """FAQ — one slide per Q&A."""
    faq = data.get("faq", [])
    if not faq:
        return

    for i, item in enumerate(faq[:8]):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        set_slide_bg(slide)
        mckinsey_header(slide, f"FREQUENTLY ASKED QUESTION {i+1}")

        q = strip_html(item.get("q", ""))
        a = strip_html(item.get("a", ""))

        add_text(slide, Inches(0.6), Inches(1.3), Inches(12), Inches(0.8),
                 q, size=18, color=NAVY, bold=True, italic=True)

        add_multiline(slide, Inches(0.6), Inches(2.3), Inches(12), Inches(4.0),
                      truncate(a, 1000), size=13, color=DARK_TEXT)


def slides_methodology(prs, data):
    """Methodology."""
    meth = data.get("methodology", "")
    if not meth:
        return

    text = strip_html(meth)
    chunks = split_into_chunks(text, 900)

    for i, chunk in enumerate(chunks[:3]):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        set_slide_bg(slide)
        suffix = f" ({i+1}/{min(len(chunks),3)})" if len(chunks) > 1 else ""
        mckinsey_header(slide, f"METHODOLOGY{suffix}",
                       "Data sources, calibration, and limitations")
        add_multiline(slide, Inches(0.6), Inches(1.6), Inches(12), Inches(5.0),
                      chunk, size=13, color=DARK_TEXT)


def slides_literature(prs, data):
    """Literature context."""
    lit = data.get("literatureContext", "")
    if not lit:
        return

    text = strip_html(lit)
    chunks = split_into_chunks(text, 900)

    for i, chunk in enumerate(chunks[:2]):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        set_slide_bg(slide)
        suffix = f" ({i+1}/{min(len(chunks),2)})" if len(chunks) > 1 else ""
        mckinsey_header(slide, f"LITERATURE CONTEXT{suffix}",
                       "Where this paper sits in the academic landscape")
        add_multiline(slide, Inches(0.6), Inches(1.6), Inches(12), Inches(5.0),
                      chunk, size=13, color=DARK_TEXT)


def slides_timeline(prs, data):
    """Timeline events."""
    timeline = data.get("timeline", [])
    if not timeline:
        return

    # 6 events per slide
    for page in range(0, min(len(timeline), 18), 6):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        set_slide_bg(slide)
        mckinsey_header(slide, "TIMELINE", "Key events in the domain's history")

        batch = timeline[page:page+6]
        for j, event in enumerate(batch):
            y = Inches(1.5) + j * Inches(0.9)
            year = str(event.get("year", ""))
            evt = strip_html(event.get("event", ""))
            sig = strip_html(event.get("significance", ""))

            add_text(slide, Inches(0.6), y, Inches(1.2), Inches(0.4),
                     year, size=14, color=NAVY, bold=True)
            add_text(slide, Inches(2.0), y, Inches(4), Inches(0.4),
                     truncate(evt, 80), size=12, color=DARK_TEXT, bold=True)
            add_text(slide, Inches(6.2), y, Inches(6.5), Inches(0.8),
                     truncate(sig, 150), size=11, color=MID_GRAY)


def slides_quiz(prs, data):
    """Quiz slide — tests comprehension."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, NAVY)

    add_text(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(0.5),
             "COMPREHENSION CHECK", size=14, color=ACCENT_CYAN, bold=True)

    wa = data.get("welfareAnalysis")
    wa = wa if isinstance(wa, dict) else {}
    beta = wa.get("betaW", "?")
    tt = data.get("theoremType", "?")
    tn = data.get("theoremName", "?")

    questions = [
        (f"Q1: What is the system beta (βW) for this domain?", f"A: {beta}"),
        (f"Q2: Is this an Impossibility or Intractability theorem?", f"A: {tt}"),
        (f"Q3: What is the theorem called?", f"A: {tn}"),
        ("Q4: Can this problem be solved with policy reform alone?",
         f"A: {'No — physical/biological constraints prevent it' if tt == 'Impossibility' else 'Yes — proven models exist (see Reform Pathway)' if tt == 'Intractability' else 'See theorem classification'}"),
    ]

    for i, (q, a) in enumerate(questions):
        y = Inches(1.5) + i * Inches(1.3)
        add_text(slide, Inches(0.8), y, Inches(11), Inches(0.5),
                 q, size=16, color=WHITE, bold=True)
        add_text(slide, Inches(1.2), y + Inches(0.5), Inches(10.5), Inches(0.5),
                 a, size=14, color=ACCENT_CYAN)


def slides_docx_tables(prs, data, docx_tables):
    """Render extracted Word doc tables as slides."""
    if not docx_tables:
        return

    for i, table in enumerate(docx_tables[:12]):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        set_slide_bg(slide)
        mckinsey_header(slide, f"DATA TABLE {i+1} (from paper)",
                       f"{len(table['rows'])} rows")

        header = table["header"]
        rows = table["rows"][:15]  # Max 15 rows per slide

        # Calculate column positions
        n_cols = len(header)
        if n_cols == 0:
            continue
        col_w = min(Inches(12) / n_cols, Inches(3))
        start_x = Inches(0.5)

        # Header row
        y = Inches(1.5)
        add_rect(slide, start_x, y, col_w * n_cols + Inches(0.4), Inches(0.35), NAVY)
        for j, h in enumerate(header[:8]):
            add_text(slide, start_x + Inches(0.1) + j * col_w, y + Inches(0.02),
                     col_w - Inches(0.1), Inches(0.3),
                     truncate(h, 40), size=10, color=WHITE, bold=True)

        # Data rows
        for r, row in enumerate(rows):
            y = Inches(1.85) + r * Inches(0.32)
            if r % 2 == 0:
                add_rect(slide, start_x, y, col_w * n_cols + Inches(0.4), Inches(0.32), SURFACE)
            for j, cell in enumerate(row[:8]):
                add_text(slide, start_x + Inches(0.1) + j * col_w, y + Inches(0.02),
                         col_w - Inches(0.1), Inches(0.28),
                         truncate(str(cell), 40), size=9, color=DARK_TEXT)

        source_line(slide, f"Extracted from {data.get('slug', 'paper')} Word document")


def slide_closing(prs, data):
    """Closing slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)

    add_rect(slide, 0, Inches(3.2), Inches(13.333), Inches(0.04), ACCENT_CYAN)

    add_text(slide, Inches(1), Inches(1.0), Inches(11), Inches(1.5),
             data.get("title", ""), size=20, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

    wa = data.get("welfareAnalysis")
    wa = wa if isinstance(wa, dict) else {}
    beta = wa.get("betaW", "")
    if beta and beta != "—":
        add_text(slide, Inches(1), Inches(3.5), Inches(11), Inches(0.8),
                 f"βW = {beta}", size=48, color=ACCENT_CYAN, bold=True, align=PP_ALIGN.CENTER)

    tt = data.get("theoremType", "")
    if tt:
        color = RED if tt == "Impossibility" else GOLD_BORDER
        add_text(slide, Inches(1), Inches(4.5), Inches(11), Inches(0.5),
                 f"{tt} Theorem", size=20, color=color, align=PP_ALIGN.CENTER)

    add_text(slide, Inches(1), Inches(5.5), Inches(11), Inches(0.4),
             "Erik Postnieks  |  Independent Researcher  |  postnieks.com",
             size=14, color=MID_GRAY, align=PP_ALIGN.CENTER)
    add_text(slide, Inches(1), Inches(6.0), Inches(11), Inches(0.4),
             "Full paper and Monte Carlo replication: github.com/epostnieks",
             size=12, color=MID_GRAY, align=PP_ALIGN.CENTER)
    add_text(slide, Inches(1), Inches(6.5), Inches(11), Inches(0.3),
             "END OF MASTERCLASS", size=11, color=ACCENT_CYAN, bold=True, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════
# MAIN GENERATOR
# ══════════════════════════════════════════════════════════════

def generate_pptx(json_path):
    """Generate a masterclass PPTX from a paper JSON file."""
    with open(json_path, 'r') as f:
        data = json.load(f)

    slug = data.get("slug", json_path.stem)

    # Find and extract Word doc tables
    docx_path = find_docx(slug)
    docx_tables = extract_docx_tables(docx_path) if docx_path else []

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Build masterclass structure
    slide_title(prs, data)

    # Chapter 1: The Numbers
    slide_chapter_divider(prs, 1, "THE NUMBERS", "System beta, welfare cost, and what they mean")
    slides_metrics(prs, data)

    # Chapter 2: Executive Summary
    slide_chapter_divider(prs, 2, "EXECUTIVE SUMMARY", "The complete argument in 2,000 words")
    slides_executive_summary(prs, data)

    # Chapter 3: Key Findings
    slide_chapter_divider(prs, 3, "KEY FINDINGS", "What the evidence shows")
    slides_key_findings(prs, data)

    # Chapter 4: The Theorem
    slide_chapter_divider(prs, 4, "THE THEOREM", f"{data.get('theoremName', 'Formal Result')}")
    slides_theorem(prs, data)

    # Chapter 5: Welfare Analysis
    slide_chapter_divider(prs, 5, "WELFARE CHANNEL DECOMPOSITION", "Where the damage comes from")
    slides_welfare_channels(prs, data)

    # Chapter 6: Monte Carlo
    slide_chapter_divider(prs, 6, "MONTE CARLO SIMULATION", "Quantifying uncertainty")
    slides_monte_carlo(prs, data)

    # Chapter 7: Evidence — Case Studies + Tables from Word doc
    slide_chapter_divider(prs, 7, "EVIDENCE", "Case studies and data tables from the paper")
    slides_case_studies(prs, data)
    slides_docx_tables(prs, data, docx_tables)

    # Chapter 8: Policy
    slide_chapter_divider(prs, 8, "POLICY ANALYSIS", "What works, what fails, and why")
    slides_policy(prs, data)

    # Chapter 9: Six-Agent Framework
    slide_chapter_divider(prs, 9, "SIX-AGENT CONFLICTORING FRAMEWORK",
                         "Who can break the Private-Systemic Tension")
    slides_agents(prs, data)

    # Chapter 10: Connections & Context
    slide_chapter_divider(prs, 10, "CONNECTIONS & CONTEXT",
                         "Cross-domain links, methodology, and literature")
    slides_cross_domain(prs, data)
    slides_methodology(prs, data)
    slides_literature(prs, data)
    slides_timeline(prs, data)

    # Chapter 11: FAQ
    slide_chapter_divider(prs, 11, "SKEPTIC'S FAQ", "Hard questions, evidence-based answers")
    slides_faq(prs, data)

    # Comprehension Check
    slides_quiz(prs, data)

    # Closing
    slide_closing(prs, data)

    out_path = OUTPUT_DIR / f"{slug}.pptx"
    prs.save(str(out_path))
    slide_count = len(prs.slides)
    return slug, slide_count, str(out_path), len(docx_tables)


def main():
    json_files = sorted(PAPER_DATA_DIR.glob("*.json"))
    print(f"Found {len(json_files)} paper JSON files")

    results = []
    errors = []

    with ThreadPoolExecutor(max_workers=8) as executor:
        futures = {executor.submit(generate_pptx, f): f for f in json_files}
        for future in as_completed(futures):
            path = futures[future]
            try:
                slug, slides, out, tables = future.result()
                results.append((slug, slides, tables))
                table_tag = f" + {tables} tables" if tables else ""
                print(f"  {slug}: {slides} slides{table_tag}")
            except Exception as e:
                errors.append((path.stem, str(e)))
                print(f"  ERROR {path.stem}: {e}")

    print(f"\n{'='*60}")
    print(f"Generated: {len(results)} decks")
    if errors:
        print(f"Errors: {len(errors)}")
        for name, err in errors:
            print(f"  - {name}: {err}")

    total_slides = sum(s for _, s, _ in results)
    total_tables = sum(t for _, _, t in results)
    print(f"Total slides: {total_slides}")
    print(f"Word doc tables embedded: {total_tables}")
    print(f"Output: {OUTPUT_DIR}")


if __name__ == "__main__":
    main()
