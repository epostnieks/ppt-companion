#!/usr/bin/env python3
"""
Master script: Generate best-in-class PowerPoint decks for all papers.
One script, all papers, parallel execution. Token-efficient.
"""

import json
import os
import re
import sys
from concurrent.futures import ThreadPoolExecutor, as_completed
from pathlib import Path
from html.parser import HTMLParser

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ── Paths ──
PAPER_DATA_DIR = Path(__file__).parent.parent / "src" / "paperData"
OUTPUT_DIR = Path(__file__).parent.parent / "pptx_output"
OUTPUT_DIR.mkdir(exist_ok=True)

# ── Colors ──
DARK_BG = RGBColor(0x0F, 0x17, 0x2A)       # Deep navy
ACCENT_CYAN = RGBColor(0x00, 0xD4, 0xFF)    # Cyan accent
ACCENT_GOLD = RGBColor(0xFF, 0xD7, 0x00)    # Gold accent
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xCC, 0xCC, 0xCC)
MID_GRAY = RGBColor(0x88, 0x88, 0x88)
DARK_CARD = RGBColor(0x1A, 0x24, 0x3B)      # Card background
RED_ACCENT = RGBColor(0xFF, 0x44, 0x44)      # For warnings/costs
GREEN_ACCENT = RGBColor(0x00, 0xCC, 0x88)    # For positive
SLIDE_BG = RGBColor(0x0A, 0x0F, 0x1E)       # Darkest

# ── HTML stripping ──
class HTMLStripper(HTMLParser):
    def __init__(self):
        super().__init__()
        self.result = []
        self.skip = False
    def handle_data(self, d):
        if not self.skip:
            self.result.append(d)
    def handle_starttag(self, tag, attrs):
        if tag in ('table', 'thead', 'tbody', 'tr', 'th', 'td'):
            self.skip = False
        if tag == 'td' or tag == 'th':
            self.result.append(' | ')
        if tag == 'tr':
            self.result.append('\n')
    def get_text(self):
        return ''.join(self.result).strip()

def strip_html(text):
    if not text or not isinstance(text, str):
        return str(text) if text else ""
    s = HTMLStripper()
    s.feed(text)
    return s.get_text()

def truncate(text, max_chars=600):
    text = strip_html(text)
    if len(text) <= max_chars:
        return text
    return text[:max_chars-3] + "..."

def set_slide_bg(slide, color=SLIDE_BG):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_shape_bg(slide, left, top, width, height, color=DARK_CARD, radius=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    shape.shadow.inherit = False
    return shape

def add_text_box(slide, left, top, width, height, text, font_size=14,
                 color=WHITE, bold=False, alignment=PP_ALIGN.LEFT, font_name="Calibri"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox

def add_bullet_list(slide, left, top, width, height, items, font_size=12,
                    color=WHITE, bullet_color=ACCENT_CYAN, max_items=8):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items[:max_items]):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        text = strip_html(item) if isinstance(item, str) else str(item)
        # Extract first sentence for brevity
        first_sent = text.split('. ')[0] + '.' if '. ' in text else text
        if len(first_sent) > 200:
            first_sent = first_sent[:197] + "..."
        p.text = f"▸ {first_sent}"
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = "Calibri"
        p.space_after = Pt(6)
    return txBox


def make_title_slide(prs, data):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
    set_slide_bg(slide, DARK_BG)

    # Accent bar at top
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.33), Inches(0.08))
    bar.fill.solid()
    bar.fill.fore_color.rgb = ACCENT_CYAN
    bar.line.fill.background()

    # "SYSTEM ASSET PRICING MODEL" label
    add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.5),
                 "SYSTEM ASSET PRICING MODEL", font_size=14, color=ACCENT_CYAN,
                 bold=True, font_name="Calibri")

    # Title
    title = data.get("title", data.get("slug", "Untitled"))
    add_text_box(slide, Inches(0.8), Inches(1.0), Inches(11), Inches(2.0),
                 title, font_size=36, color=WHITE, bold=True, font_name="Calibri")

    # Theorem type badge
    theorem_type = data.get("theoremType", "")
    if theorem_type:
        badge_color = RED_ACCENT if theorem_type == "Impossibility" else ACCENT_GOLD
        add_text_box(slide, Inches(0.8), Inches(3.2), Inches(4), Inches(0.5),
                     f"⬢ {theorem_type} Theorem", font_size=16, color=badge_color, bold=True)

    # Theorem name
    theorem_name = data.get("theoremName", "")
    if theorem_name:
        add_text_box(slide, Inches(0.8), Inches(3.8), Inches(10), Inches(0.5),
                     theorem_name, font_size=20, color=LIGHT_GRAY, bold=False)

    # Epigraph
    epigraph = data.get("epigraph", "")
    if epigraph:
        add_text_box(slide, Inches(0.8), Inches(5.0), Inches(10), Inches(1.2),
                     f'"{epigraph}"', font_size=16, color=MID_GRAY,
                     alignment=PP_ALIGN.LEFT, font_name="Calibri")

    # Author
    add_text_box(slide, Inches(0.8), Inches(6.5), Inches(10), Inches(0.4),
                 "Erik Postnieks  |  Independent Researcher  |  2026",
                 font_size=12, color=MID_GRAY)


def make_metrics_slide(prs, data):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)

    add_text_box(slide, Inches(0.8), Inches(0.3), Inches(11), Inches(0.6),
                 "KEY METRICS", font_size=28, color=ACCENT_CYAN, bold=True)

    wa = data.get("welfareAnalysis")
    wa = wa if isinstance(wa, dict) else {}
    if not wa and not data.get("executiveSummary"):
        add_text_box(slide, Inches(0.8), Inches(1.5), Inches(11), Inches(1.0),
                     "Welfare analysis data not yet available for this paper.",
                     font_size=18, color=MID_GRAY)
        return

    # Extract metrics
    beta = wa.get("betaW", "—")
    ci = wa.get("ci90", "—")
    pi_raw = wa.get("pi", "—")
    dw_raw = wa.get("deltaW", "—")

    # Clean pi/dw to short form
    def short_money(val):
        if not val or val == "—":
            return "—"
        val = str(val)
        # Extract first dollar amount
        m = re.search(r'[-−]?\$[\d,.]+\s*(billion|trillion|million)?', val, re.IGNORECASE)
        if m:
            return m.group(0)
        return val[:60]

    pi_short = short_money(pi_raw)
    dw_short = short_money(dw_raw)

    # 4 metric cards
    metrics = [
        ("βW", str(beta), "System Beta", ACCENT_CYAN),
        ("90% CI", str(ci), "Confidence Interval", ACCENT_GOLD),
        ("Π", pi_short, "Annual Revenue", GREEN_ACCENT),
        ("ΔW", dw_short, "Welfare Cost", RED_ACCENT),
    ]

    card_w = Inches(2.8)
    card_h = Inches(2.5)
    start_x = Inches(0.5)
    gap = Inches(0.3)

    for i, (label, value, subtitle, accent) in enumerate(metrics):
        x = start_x + i * (card_w + gap)
        y = Inches(1.5)

        card = add_shape_bg(slide, x, y, card_w, card_h, DARK_CARD)

        # Label
        add_text_box(slide, x + Inches(0.2), y + Inches(0.2), card_w - Inches(0.4), Inches(0.4),
                     label, font_size=14, color=accent, bold=True)
        # Value
        val_size = 28 if len(value) < 15 else 18 if len(value) < 30 else 14
        add_text_box(slide, x + Inches(0.2), y + Inches(0.7), card_w - Inches(0.4), Inches(1.0),
                     value, font_size=val_size, color=WHITE, bold=True)
        # Subtitle
        add_text_box(slide, x + Inches(0.2), y + Inches(1.8), card_w - Inches(0.4), Inches(0.4),
                     subtitle, font_size=11, color=MID_GRAY)

    # Theorem type at bottom
    tt = data.get("theoremType", "")
    if tt:
        badge_text = f"Classification: {tt} Theorem"
        badge_color = RED_ACCENT if tt == "Impossibility" else ACCENT_GOLD
        add_text_box(slide, Inches(0.8), Inches(4.5), Inches(11), Inches(0.5),
                     badge_text, font_size=16, color=badge_color, bold=True)


def make_executive_summary_slides(prs, data):
    es = data.get("executiveSummary", "")
    if not es:
        return

    text = strip_html(es)
    # Split into ~2 slides worth
    paras = [p.strip() for p in text.split('\n') if p.strip()]
    if not paras:
        paras = [text]

    # Combine short paragraphs, split into chunks of ~800 chars
    chunks = []
    current = ""
    for p in paras:
        if len(current) + len(p) > 900 and current:
            chunks.append(current.strip())
            current = p
        else:
            current = current + "\n\n" + p if current else p
    if current:
        chunks.append(current.strip())

    for i, chunk in enumerate(chunks[:3]):  # Max 3 slides
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        set_slide_bg(slide)

        suffix = f" ({i+1}/{min(len(chunks),3)})" if len(chunks) > 1 else ""
        add_text_box(slide, Inches(0.8), Inches(0.3), Inches(11), Inches(0.6),
                     f"EXECUTIVE SUMMARY{suffix}", font_size=28, color=ACCENT_CYAN, bold=True)

        add_text_box(slide, Inches(0.8), Inches(1.2), Inches(11.5), Inches(5.8),
                     truncate(chunk, 1200), font_size=13, color=LIGHT_GRAY)


def make_key_findings_slides(prs, data):
    findings = data.get("keyFindings", [])
    if not findings:
        return

    # 4 findings per slide
    per_slide = 4
    for page in range(0, min(len(findings), 12), per_slide):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        set_slide_bg(slide)

        page_num = page // per_slide + 1
        total_pages = min(len(findings), 12) // per_slide + (1 if min(len(findings), 12) % per_slide else 0)
        suffix = f" ({page_num}/{total_pages})" if total_pages > 1 else ""
        add_text_box(slide, Inches(0.8), Inches(0.3), Inches(11), Inches(0.6),
                     f"KEY FINDINGS{suffix}", font_size=28, color=ACCENT_CYAN, bold=True)

        batch = findings[page:page+per_slide]
        card_h = Inches(1.2)
        for j, finding in enumerate(batch):
            y = Inches(1.2) + j * (card_h + Inches(0.15))
            # Card background
            add_shape_bg(slide, Inches(0.5), y, Inches(12.2), card_h, DARK_CARD)

            # Number badge
            add_text_box(slide, Inches(0.7), y + Inches(0.15), Inches(0.5), Inches(0.5),
                         str(page + j + 1), font_size=20, color=ACCENT_CYAN, bold=True)

            # Finding text - first 2 sentences
            text = strip_html(finding) if isinstance(finding, str) else str(finding)
            sentences = re.split(r'(?<=[.!?])\s+', text)
            short = ' '.join(sentences[:2])
            if len(short) > 250:
                short = short[:247] + "..."

            add_text_box(slide, Inches(1.3), y + Inches(0.1), Inches(11.2), card_h - Inches(0.2),
                         short, font_size=12, color=WHITE)


def make_theorem_slide(prs, data):
    theorem = data.get("theorem", {})
    if not theorem:
        return

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)

    add_text_box(slide, Inches(0.8), Inches(0.3), Inches(11), Inches(0.6),
                 "THEOREM", font_size=28, color=ACCENT_GOLD, bold=True)

    theorem_name = data.get("theoremName", "")
    if theorem_name:
        add_text_box(slide, Inches(0.8), Inches(0.9), Inches(11), Inches(0.5),
                     theorem_name, font_size=20, color=WHITE, bold=True)

    # Plain language explanation
    plain = theorem.get("plain", "")
    if plain:
        add_shape_bg(slide, Inches(0.5), Inches(1.6), Inches(12.2), Inches(3.5), DARK_CARD)
        add_text_box(slide, Inches(0.8), Inches(1.7), Inches(11.5), Inches(0.3),
                     "PLAIN LANGUAGE", font_size=11, color=ACCENT_CYAN, bold=True)
        add_text_box(slide, Inches(0.8), Inches(2.1), Inches(11.5), Inches(2.8),
                     truncate(plain, 800), font_size=12, color=LIGHT_GRAY)

    # Axioms
    axioms = theorem.get("axioms", [])
    if axioms:
        y = Inches(5.3)
        add_text_box(slide, Inches(0.8), y, Inches(11), Inches(0.4),
                     "AXIOMS", font_size=14, color=ACCENT_GOLD, bold=True)
        for k, ax in enumerate(axioms[:4]):
            ax_text = f"{ax.get('id', '')}: {ax.get('name', '')} — {ax.get('statement', '')}"
            add_text_box(slide, Inches(0.8), y + Inches(0.4) + k * Inches(0.35),
                         Inches(11.5), Inches(0.35),
                         truncate(ax_text, 150), font_size=10, color=LIGHT_GRAY)


def make_welfare_channels_slide(prs, data):
    wa = data.get("welfareAnalysis")
    wa = wa if isinstance(wa, dict) else {}
    channels = wa.get("channels", [])
    if not channels:
        return

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)

    add_text_box(slide, Inches(0.8), Inches(0.3), Inches(11), Inches(0.6),
                 "WELFARE CHANNELS", font_size=28, color=RED_ACCENT, bold=True)

    # Channel cards - 2 columns
    col_w = Inches(6.0)
    per_col = (len(channels) + 1) // 2
    for i, ch in enumerate(channels[:8]):
        col = i // per_col
        row = i % per_col
        x = Inches(0.5) + col * (col_w + Inches(0.3))
        y = Inches(1.2) + row * Inches(1.3)

        add_shape_bg(slide, x, y, col_w, Inches(1.15), DARK_CARD)

        name = ch.get("name", f"Channel {i+1}")
        value = ch.get("value", "—")
        desc = ch.get("description", "")

        # Channel name + value
        add_text_box(slide, x + Inches(0.15), y + Inches(0.08), col_w * 0.6, Inches(0.35),
                     name, font_size=13, color=ACCENT_CYAN, bold=True)
        add_text_box(slide, x + col_w * 0.6, y + Inches(0.08), col_w * 0.35, Inches(0.35),
                     strip_html(str(value)), font_size=13, color=RED_ACCENT, bold=True,
                     alignment=PP_ALIGN.RIGHT)

        # Description (truncated)
        first_sent = strip_html(desc).split('. ')[0] + '.' if desc else ""
        if len(first_sent) > 180:
            first_sent = first_sent[:177] + "..."
        add_text_box(slide, x + Inches(0.15), y + Inches(0.45), col_w - Inches(0.3), Inches(0.65),
                     first_sent, font_size=10, color=MID_GRAY)


def make_case_studies_slide(prs, data):
    cases = data.get("caseStudies", [])
    if not cases:
        return

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)

    add_text_box(slide, Inches(0.8), Inches(0.3), Inches(11), Inches(0.6),
                 "CASE STUDIES", font_size=28, color=ACCENT_CYAN, bold=True)

    for i, case in enumerate(cases[:4]):
        y = Inches(1.2) + i * Inches(1.4)
        add_shape_bg(slide, Inches(0.5), y, Inches(12.2), Inches(1.25), DARK_CARD)

        title = case.get("title", f"Case {i+1}")
        content = strip_html(case.get("content", ""))

        add_text_box(slide, Inches(0.7), y + Inches(0.08), Inches(11.5), Inches(0.35),
                     title, font_size=14, color=ACCENT_GOLD, bold=True)

        # First 2 sentences
        sentences = re.split(r'(?<=[.!?])\s+', content)
        short = ' '.join(sentences[:3])
        if len(short) > 280:
            short = short[:277] + "..."
        add_text_box(slide, Inches(0.7), y + Inches(0.45), Inches(11.5), Inches(0.75),
                     short, font_size=11, color=LIGHT_GRAY)


def make_policy_slide(prs, data):
    policy = data.get("policyAnalysis")
    if not policy or not isinstance(policy, dict):
        return

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)

    add_text_box(slide, Inches(0.8), Inches(0.3), Inches(11), Inches(0.6),
                 "POLICY ANALYSIS", font_size=28, color=ACCENT_GOLD, bold=True)

    sections = [
        ("Current Framework", policy.get("currentFramework", ""), ACCENT_CYAN),
        ("Key Failures", policy.get("failures", ""), RED_ACCENT),
        ("Reform Pathway", policy.get("reform", ""), GREEN_ACCENT),
    ]

    for i, (label, content, accent) in enumerate(sections):
        if not content:
            continue
        y = Inches(1.1) + i * Inches(1.8)
        add_shape_bg(slide, Inches(0.5), y, Inches(12.2), Inches(1.6), DARK_CARD)
        add_text_box(slide, Inches(0.7), y + Inches(0.08), Inches(11), Inches(0.35),
                     label.upper(), font_size=12, color=accent, bold=True)
        add_text_box(slide, Inches(0.7), y + Inches(0.4), Inches(11.5), Inches(1.1),
                     truncate(strip_html(content), 500), font_size=11, color=LIGHT_GRAY)


def make_agents_slide(prs, data):
    agents = data.get("agents", {})
    if not agents:
        return

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)

    add_text_box(slide, Inches(0.8), Inches(0.3), Inches(11), Inches(0.6),
                 "SIX-AGENT ACTION FRAMEWORK", font_size=28, color=ACCENT_CYAN, bold=True)

    agent_labels = {
        "insider": ("Whistleblower", "🔔"),
        "whistleblower": ("Whistleblower", "🔔"),
        "plaintiff": ("Plaintiff", "⚖️"),
        "regulator": ("Regulator", "🏛️"),
        "legislator": ("Legislator", "📜"),
        "investor": ("Investor", "📊"),
        "supranational": ("Supranational", "🌍"),
    }

    keys = [k for k in agents.keys() if k in agent_labels]
    # 3x2 grid
    card_w = Inches(4.0)
    card_h = Inches(1.6)

    for i, key in enumerate(keys[:6]):
        col = i % 3
        row = i // 3
        x = Inches(0.3) + col * (card_w + Inches(0.15))
        y = Inches(1.2) + row * (card_h + Inches(0.15))

        label, emoji = agent_labels.get(key, (key.title(), "●"))
        content = strip_html(str(agents[key]))

        add_shape_bg(slide, x, y, card_w, card_h, DARK_CARD)
        add_text_box(slide, x + Inches(0.15), y + Inches(0.08), card_w - Inches(0.3), Inches(0.35),
                     f"{emoji}  {label.upper()}", font_size=12, color=ACCENT_CYAN, bold=True)

        # First 3 sentences
        sentences = re.split(r'(?<=[.!?])\s+', content)
        short = ' '.join(sentences[:3])
        if len(short) > 250:
            short = short[:247] + "..."
        add_text_box(slide, x + Inches(0.15), y + Inches(0.4), card_w - Inches(0.3), card_h - Inches(0.5),
                     short, font_size=10, color=LIGHT_GRAY)


def make_cross_domain_slide(prs, data):
    connections = data.get("crossDomainConnections", [])
    if not connections:
        return

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)

    add_text_box(slide, Inches(0.8), Inches(0.3), Inches(11), Inches(0.6),
                 "CROSS-DOMAIN CONNECTIONS", font_size=28, color=ACCENT_CYAN, bold=True)

    for i, conn in enumerate(connections[:6]):
        col = i % 2
        row = i // 2
        x = Inches(0.5) + col * Inches(6.3)
        y = Inches(1.2) + row * Inches(1.5)

        add_shape_bg(slide, x, y, Inches(6.0), Inches(1.35), DARK_CARD)

        domain = conn.get("domain", f"Domain {i+1}")
        desc = strip_html(conn.get("connection", ""))

        add_text_box(slide, x + Inches(0.15), y + Inches(0.08), Inches(5.6), Inches(0.3),
                     domain, font_size=13, color=ACCENT_GOLD, bold=True)

        sentences = re.split(r'(?<=[.!?])\s+', desc)
        short = ' '.join(sentences[:2])
        if len(short) > 200:
            short = short[:197] + "..."
        add_text_box(slide, x + Inches(0.15), y + Inches(0.4), Inches(5.6), Inches(0.85),
                     short, font_size=10, color=LIGHT_GRAY)


def make_monte_carlo_slide(prs, data):
    wa = data.get("welfareAnalysis")
    wa = wa if isinstance(wa, dict) else {}
    mc = wa.get("monteCarlo", "")
    if not mc:
        return

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)

    add_text_box(slide, Inches(0.8), Inches(0.3), Inches(11), Inches(0.6),
                 "MONTE CARLO SIMULATION", font_size=28, color=ACCENT_GOLD, bold=True)

    add_shape_bg(slide, Inches(0.5), Inches(1.1), Inches(12.2), Inches(5.5), DARK_CARD)

    text = strip_html(mc)
    add_text_box(slide, Inches(0.8), Inches(1.3), Inches(11.5), Inches(5.0),
                 truncate(text, 1400), font_size=12, color=LIGHT_GRAY)


def make_closing_slide(prs, data):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)

    # Accent bar
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(3.2), Inches(13.33), Inches(0.04))
    bar.fill.solid()
    bar.fill.fore_color.rgb = ACCENT_CYAN
    bar.line.fill.background()

    title = data.get("title", "")
    add_text_box(slide, Inches(1), Inches(1.5), Inches(11), Inches(1.5),
                 title, font_size=24, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

    wa = data.get("welfareAnalysis")
    wa = wa if isinstance(wa, dict) else {}
    beta = wa.get("betaW", "")
    if beta:
        add_text_box(slide, Inches(1), Inches(3.5), Inches(11), Inches(0.8),
                     f"βW = {beta}", font_size=48, color=ACCENT_CYAN, bold=True,
                     alignment=PP_ALIGN.CENTER)

    tt = data.get("theoremType", "")
    if tt:
        color = RED_ACCENT if tt == "Impossibility" else ACCENT_GOLD
        add_text_box(slide, Inches(1), Inches(4.5), Inches(11), Inches(0.5),
                     f"{tt} Theorem", font_size=20, color=color, bold=False,
                     alignment=PP_ALIGN.CENTER)

    add_text_box(slide, Inches(1), Inches(5.8), Inches(11), Inches(0.4),
                 "Erik Postnieks  |  Independent Researcher  |  postnieks.com",
                 font_size=14, color=MID_GRAY, alignment=PP_ALIGN.CENTER)

    add_text_box(slide, Inches(1), Inches(6.3), Inches(11), Inches(0.4),
                 "Full paper and Monte Carlo replication: github.com/epostnieks",
                 font_size=12, color=MID_GRAY, alignment=PP_ALIGN.CENTER)


def generate_pptx(json_path):
    """Generate a single PPTX from a paper JSON file."""
    with open(json_path, 'r') as f:
        data = json.load(f)

    slug = data.get("slug", json_path.stem)

    prs = Presentation()
    # Widescreen 16:9
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Build slides
    make_title_slide(prs, data)
    make_metrics_slide(prs, data)
    make_executive_summary_slides(prs, data)
    make_key_findings_slides(prs, data)
    make_theorem_slide(prs, data)
    make_welfare_channels_slide(prs, data)
    make_case_studies_slide(prs, data)
    make_policy_slide(prs, data)
    make_agents_slide(prs, data)
    make_cross_domain_slide(prs, data)
    make_monte_carlo_slide(prs, data)
    make_closing_slide(prs, data)

    out_path = OUTPUT_DIR / f"{slug}.pptx"
    prs.save(str(out_path))
    slide_count = len(prs.slides)
    return slug, slide_count, str(out_path)


def main():
    json_files = sorted(PAPER_DATA_DIR.glob("*.json"))
    print(f"Found {len(json_files)} paper JSON files")

    results = []
    errors = []

    with ThreadPoolExecutor(max_workers=8) as executor:
        futures = {executor.submit(generate_pptx, f): f for f in json_files}
        for future in as_completed(futures):
            path = futures[future]
            try:
                slug, slides, out = future.result()
                results.append((slug, slides, out))
                print(f"  ✓ {slug}: {slides} slides")
            except Exception as e:
                errors.append((path.stem, str(e)))
                print(f"  ✗ {path.stem}: {e}")

    print(f"\n{'='*60}")
    print(f"Generated: {len(results)} decks")
    if errors:
        print(f"Errors: {len(errors)}")
        for name, err in errors:
            print(f"  - {name}: {err}")

    total_slides = sum(s for _, s, _ in results)
    print(f"Total slides: {total_slides}")
    print(f"Output directory: {OUTPUT_DIR}")


if __name__ == "__main__":
    main()
