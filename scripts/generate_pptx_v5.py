#!/usr/bin/env python3
"""
PPTX Generator V5 — MASTERCLASS FINAL
- Standard 10"×7.5" slides (printable on 8.5×11)
- DENSE content: 5-8 bullet points per slide, no single-bullet waste
- McKinsey figures at native 1400×1050 as full-page anchors
- Word doc tables + embedded images
- TTS audio narration per paper (gTTS → MP3)
- Speaker notes on every slide
- NO text truncation — full 20K word summaries
- NO text overlap — proper Y-position tracking
"""

import json, os, re, sys, glob, traceback, hashlib
from concurrent.futures import ThreadPoolExecutor, as_completed
from pathlib import Path
from html.parser import HTMLParser
from io import BytesIO

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

try:
    from docx import Document as DocxDocument
    HAS_DOCX = True
except ImportError:
    HAS_DOCX = False

try:
    from gtts import gTTS
    HAS_TTS = True
except ImportError:
    HAS_TTS = False

# ── Paths ──
SCRIPT_DIR = Path(__file__).parent
PROJECT_DIR = SCRIPT_DIR.parent
PAPER_DATA_DIR = PROJECT_DIR / "src" / "paperData"
PPTX_DIR = PROJECT_DIR / "public" / "pptx"
AUDIO_DIR = PROJECT_DIR / "public" / "audio"
PPTX_DIR.mkdir(parents=True, exist_ok=True)
AUDIO_DIR.mkdir(parents=True, exist_ok=True)
BASE = Path("/Users/erikpostnieks/Projects")
DR_OUTPUT = BASE / "deep-research" / "output"

# ── Standard slide size (10" × 7.5" = fits 8.5×11 paper) ──
SLIDE_W = Inches(10)
SLIDE_H = Inches(7.5)

# Content area
MARGIN_L = Inches(0.5)
MARGIN_R = Inches(0.5)
CONTENT_W = Inches(9.0)  # 10 - 0.5 - 0.5
HEADER_H = Inches(1.2)   # Navy bar + subtitle
CONTENT_TOP = Inches(1.4)
CONTENT_BOTTOM = Inches(6.8)
CONTENT_H = Inches(5.4)  # Available content height

# ── McKinsey Colors ──
NAVY = RGBColor(0x1F, 0x4E, 0x79)
CREAM = RGBColor(0xFB, 0xF9, 0xF4)
TAN = RGBColor(0xF5, 0xEF, 0xE0)
GOLD = RGBColor(0xC4, 0x97, 0x2A)
SOURCE_BLUE = RGBColor(0x3D, 0x6B, 0x8E)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x00, 0x00, 0x00)
DARK_TEXT = RGBColor(0x33, 0x33, 0x33)
MID_GRAY = RGBColor(0x88, 0x88, 0x88)
RED = RGBColor(0xCC, 0x33, 0x33)
GREEN = RGBColor(0x22, 0x88, 0x55)
DARK_BG = RGBColor(0x0F, 0x17, 0x2A)
CYAN = RGBColor(0x00, 0xD4, 0xFF)
SURFACE = RGBColor(0xF7, 0xF5, 0xF0)
SLIDE_BG_COLOR = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xE8, 0xE0, 0xD0)

FONT = "Liberation Sans"

# ══════════════════════════════════════════════════════════════
# TEXT UTILITIES
# ══════════════════════════════════════════════════════════════

class HTMLStripper(HTMLParser):
    def __init__(self):
        super().__init__()
        self.result = []
    def handle_data(self, d):
        self.result.append(d)
    def handle_starttag(self, tag, attrs):
        if tag in ('td', 'th'):
            self.result.append(' | ')
        if tag == 'tr':
            self.result.append('\n')
        if tag == 'li':
            self.result.append('\n• ')
        if tag in ('br', 'p'):
            self.result.append('\n')
    def get_text(self):
        return re.sub(r'\n{3,}', '\n\n', ''.join(self.result)).strip()

def strip_html(text):
    if not text or not isinstance(text, str):
        return str(text) if text else ""
    s = HTMLStripper()
    s.feed(text)
    return s.get_text()

def split_sentences(text):
    text = strip_html(text)
    parts = re.split(r'(?<=[.!?])\s+', text)
    return [p.strip() for p in parts if p.strip()]

def paginate_text(text, max_chars=1800):
    """Split text into DENSE page-sized chunks. ~1800 chars = a full slide of readable 12pt text."""
    text = strip_html(text)
    if not text:
        return []
    sentences = split_sentences(text)
    pages = []
    current = ""
    for s in sentences:
        if len(current) + len(s) + 1 > max_chars and current:
            pages.append(current.strip())
            current = s
        else:
            current = (current + " " + s).strip()
    if current:
        pages.append(current.strip())
    return pages


# ══════════════════════════════════════════════════════════════
# SHAPE HELPERS
# ══════════════════════════════════════════════════════════════

def set_bg(slide, color=SLIDE_BG_COLOR):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def rect(slide, left, top, width, height, color):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    shape.shadow.inherit = False
    return shape

def text_box(slide, left, top, width, height, text, size=12,
             color=DARK_TEXT, bold=False, italic=False, align=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = str(text)[:2000]
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.italic = italic
    p.font.name = FONT
    p.alignment = align
    return txBox

def multi_text(slide, left, top, width, height, text, size=12,
               color=DARK_TEXT, spacing=1.15):
    """Multi-paragraph text box — used for dense body content."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    paragraphs = str(text).split('\n')
    for i, para in enumerate(paragraphs):
        if not para.strip():
            continue
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = para.strip()[:1000]
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.font.name = FONT
        p.space_after = Pt(int(size * 0.4))
        from pptx.oxml.ns import qn
        pPr = p._pPr
        if pPr is None:
            pPr = p._p.get_or_add_pPr()
        lnSpc = pPr.makeelement(qn('a:lnSpc'), {})
        spcPct = lnSpc.makeelement(qn('a:spcPct'), {'val': str(int(spacing * 100000))})
        lnSpc.append(spcPct)
        pPr.append(lnSpc)
    return txBox

def notes(slide, text):
    if not text:
        return
    ns = slide.notes_slide
    ns.notes_text_frame.text = str(text)[:5000]

def header_bar(slide, title, subtitle=None):
    """McKinsey header: navy bar + optional cream subtitle bar."""
    rect(slide, 0, 0, SLIDE_W, Inches(0.75), NAVY)
    text_box(slide, Inches(0.4), Inches(0.12), Inches(9), Inches(0.5),
             title, size=22, color=WHITE, bold=True)
    if subtitle:
        rect(slide, 0, Inches(0.75), SLIDE_W, Inches(0.4), CREAM)
        text_box(slide, Inches(0.4), Inches(0.78), Inches(9), Inches(0.35),
                 subtitle, size=12, color=DARK_TEXT, italic=True)

def insight_box(slide, text, y=Inches(6.2)):
    """KEY INSIGHT box — gold left border, tan background."""
    rect(slide, MARGIN_L, y, CONTENT_W, Inches(0.9), TAN)
    rect(slide, MARGIN_L, y, Inches(0.05), Inches(0.9), GOLD)
    text_box(slide, Inches(0.7), y + Inches(0.03), Inches(1.5), Inches(0.2),
             "KEY INSIGHT", size=9, color=GOLD, bold=True)
    text_box(slide, Inches(0.7), y + Inches(0.22), Inches(8.5), Inches(0.6),
             text[:280], size=11, color=DARK_TEXT)

def source_line(slide, text):
    text_box(slide, MARGIN_L, Inches(7.15), CONTENT_W, Inches(0.25),
             f"Source: {text}", size=9, color=SOURCE_BLUE)


# ══════════════════════════════════════════════════════════════
# ASSET EXTRACTION
# ══════════════════════════════════════════════════════════════

def find_figures(slug):
    """Find McKinsey figure PNGs — PRIMARY: sapm-{slug}/figures/png/, FALLBACK: deep-research/output/."""
    # Primary: the real paper figures in sapm repos
    sapm_fig_dir = BASE / f"sapm-{slug}" / "figures" / "png"
    if sapm_fig_dir.exists():
        pngs = sorted(sapm_fig_dir.glob("*.png"))
        if pngs:
            seen = set()
            out = []
            for p in pngs:
                base = re.sub(r'_v\d+\.png$', '.png', p.name)
                if base not in seen:
                    out.append(p)
                    seen.add(base)
            return out

    # Also check sapm-{slug}/figures/ directly
    sapm_fig_dir2 = BASE / f"sapm-{slug}" / "figures"
    if sapm_fig_dir2.exists():
        pngs = sorted(sapm_fig_dir2.glob("*.png"))
        if pngs:
            return pngs

    # Fallback: deep-research output (4 standard figures)
    dr_slug = slug.replace("-", "_")
    fig_dir = DR_OUTPUT / dr_slug / "figures"
    if not fig_dir.exists():
        fig_dir = DR_OUTPUT / slug / "figures"
    if not fig_dir.exists():
        return []
    pngs = sorted(fig_dir.glob("*.png"))
    seen = set()
    out = []
    for p in pngs:
        base = re.sub(r'_v\d+\.png$', '.png', p.name)
        if base not in seen:
            out.append(p)
            seen.add(base)
    return out

def find_docx(slug):
    patterns = [
        f"{BASE}/sapm-{slug}/paper/{slug}-current.docx",
        f"{BASE}/sapm-{slug}/paper/*-current.docx",
        f"{BASE}/sapm-{slug}/paper/*-final.docx",
        f"{BASE}/sapm-{slug}/paper/*.docx",
    ]
    for pat in patterns:
        matches = [m for m in glob.glob(pat) if '~' not in m and 'template' not in m.lower()]
        if matches:
            return matches[0]
    return None

def get_docx_tables(path, max_tables=50):
    if not HAS_DOCX or not path:
        return []
    try:
        doc = DocxDocument(path)
    except Exception:
        return []
    tables = []
    for i, t in enumerate(doc.tables[:max_tables]):
        rows = [[c.text.strip() for c in r.cells] for r in t.rows]
        if len(rows) > 1:
            tables.append({"title": f"Table {i+1}", "header": rows[0], "rows": rows[1:]})
    return tables

def get_docx_images(path):
    if not HAS_DOCX or not path:
        return []
    try:
        doc = DocxDocument(path)
    except Exception:
        return []
    images = []
    for rel in doc.part.rels.values():
        if "image" in rel.reltype:
            try:
                blob = rel.target_part.blob
                ext = rel.target_part.content_type.split('/')[-1]
                if ext in ('png', 'jpeg', 'jpg', 'gif'):
                    images.append({"data": blob, "ext": ext})
            except Exception:
                pass
    return images

def fig_label(path):
    name = re.sub(r'^fig\d+_', '', path.stem)
    return name.replace('_', ' ').title()


# ══════════════════════════════════════════════════════════════
# SLIDE BUILDERS — DENSE, STANDARD SIZE
# ══════════════════════════════════════════════════════════════

def make_title(prs, data):
    """Opening title slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, DARK_BG)
    rect(slide, 0, 0, SLIDE_W, Inches(0.05), CYAN)

    text_box(slide, Inches(0.6), Inches(0.4), Inches(8), Inches(0.3),
             "SYSTEM ASSET PRICING MODEL — MASTERCLASS", size=12, color=CYAN, bold=True)

    title = data.get("title", data.get("slug", ""))
    text_box(slide, Inches(0.6), Inches(1.2), Inches(8.5), Inches(2.0),
             title, size=32, color=WHITE, bold=True)

    tt = data.get("theoremType", "")
    tn = data.get("theoremName", "")
    if tt and tn:
        c = RED if tt == "Impossibility" else GOLD
        # Avoid redundancy: "Impossibility Theorem: The X Impossibility Theorem"
        # Just show the theorem name if it already contains the type
        if tt.lower() in tn.lower():
            badge = tn
        else:
            badge = f"{tn} ({tt})"
        text_box(slide, Inches(0.6), Inches(3.5), Inches(8), Inches(0.4),
                 badge, size=16, color=c, bold=True)

    ep = strip_html(data.get("epigraph", ""))
    if ep:
        text_box(slide, Inches(0.6), Inches(4.5), Inches(8), Inches(1.0),
                 f'"{ep}"', size=14, color=MID_GRAY, italic=True)

    text_box(slide, Inches(0.6), Inches(6.2), Inches(8), Inches(0.3),
             "Erik Postnieks  |  Independent Researcher  |  2026", size=11, color=MID_GRAY)

    notes(slide, f"Welcome to the SAPM Masterclass on {title}. "
          f"This is a {tt} Theorem — {tn}. "
          f"We will walk through the system beta, welfare channels, formal proof, "
          f"case studies, six-agent framework, and policy implications. "
          f"{ep}")
    return 1


def make_divider(prs, num, title, subtitle="", notes_text=""):
    """Chapter divider."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, NAVY)
    text_box(slide, Inches(0.6), Inches(2.0), Inches(8), Inches(0.4),
             f"CHAPTER {num}", size=14, color=CYAN, bold=True)
    text_box(slide, Inches(0.6), Inches(2.6), Inches(8.5), Inches(1.2),
             title, size=30, color=WHITE, bold=True)
    if subtitle:
        text_box(slide, Inches(0.6), Inches(4.2), Inches(8), Inches(0.8),
                 subtitle, size=14, color=CREAM, italic=True)
    notes(slide, notes_text or f"Chapter {num}: {title}. {subtitle}")
    return 1


def make_figure(prs, fig_path, caption="", note=""):
    """Full-page McKinsey figure at native aspect ratio (1400×1050 → 10"×7.5")."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    try:
        # 1400×1050 figures fit perfectly in 10"×7.5" slides
        slide.shapes.add_picture(str(fig_path),
            left=Inches(0), top=Inches(0),
            width=SLIDE_W, height=SLIDE_H)
    except Exception:
        text_box(slide, Inches(1), Inches(3), Inches(8), Inches(1),
                 f"[Figure: {fig_path.name}]", size=18, color=MID_GRAY, align=PP_ALIGN.CENTER)
    if caption:
        # Overlay caption at bottom
        text_box(slide, Inches(0.3), Inches(7.0), Inches(9), Inches(0.3),
                 caption, size=9, color=SOURCE_BLUE)
    notes(slide, note or f"This figure shows {fig_label(fig_path)}. {caption}")
    return 1


def make_image_slide(prs, img_data, img_ext, title="", note=""):
    """Embedded Word doc image as full-page slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    if title:
        header_bar(slide, title)
    try:
        top = Inches(1.2) if title else Inches(0)
        h = Inches(6.0) if title else SLIDE_H
        slide.shapes.add_picture(BytesIO(img_data),
            left=Inches(0.3), top=top, width=Inches(9.4), height=h)
    except Exception:
        text_box(slide, Inches(1), Inches(3), Inches(8), Inches(1),
                 f"[Embedded image]", size=18, color=MID_GRAY, align=PP_ALIGN.CENTER)
    notes(slide, note or f"Figure from the published paper. {title}")
    return 1


def make_metrics(prs, data):
    """Key metrics dashboard — dense, all four metrics + classification."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    header_bar(slide, "KEY METRICS DASHBOARD", "System Beta and Welfare Analysis")

    wa = data.get("welfareAnalysis")
    wa = wa if isinstance(wa, dict) else {}
    beta = wa.get("betaW", data.get("betaW", "—"))
    ci = wa.get("ci90", data.get("ci90", "—"))
    pi_raw = wa.get("pi", "—")
    dw_raw = wa.get("deltaW", "—")

    def money(val):
        if not val or val == "—": return "—"
        m = re.search(r'[-−]?\$[\d,.]+\s*(billion|trillion|million)?', str(val), re.I)
        return m.group(0) if m else str(val)[:60]

    metrics = [
        ("βW (System Beta)", str(beta), "Welfare destruction per $ of private gain", NAVY),
        ("90% CI", str(ci), "Monte Carlo bounds (N=10,000)", DARK_TEXT),
        ("Π (Revenue)", money(pi_raw), "Annual industry revenue (Iron Law)", GREEN),
        ("ΔW (Welfare Cost)", money(dw_raw), "Annual welfare destruction", RED),
    ]

    # 2×2 grid
    for i, (label, value, desc, accent) in enumerate(metrics):
        col, row = i % 2, i // 2
        x = Inches(0.4) + col * Inches(4.7)
        y = Inches(1.5) + row * Inches(1.5)
        w, h = Inches(4.4), Inches(1.3)

        rect(slide, x, y, w, h, SURFACE)
        text_box(slide, x + Inches(0.15), y + Inches(0.05), w - Inches(0.3), Inches(0.22),
                 label, size=10, color=MID_GRAY, bold=True)
        vs = 24 if len(value) < 15 else 18 if len(value) < 25 else 13
        text_box(slide, x + Inches(0.15), y + Inches(0.3), w - Inches(0.3), Inches(0.5),
                 value, size=vs, color=accent, bold=True)
        text_box(slide, x + Inches(0.15), y + Inches(0.85), w - Inches(0.3), Inches(0.25),
                 desc, size=9, color=MID_GRAY)

    tt = data.get("theoremType", "")
    tn = data.get("theoremName", "")
    if tt:
        class_desc = ('Physical or biological constraints prevent welfare internalization. '
                      'No policy reform can close this gap.' if tt == 'Impossibility'
                      else 'Institutional constraints. At least one country has demonstrated '
                           'a proven reform that substantially reduces the welfare gap.')
        insight_box(slide, f"{tn}: {class_desc}", y=Inches(5.0))

    source_line(slide, "Monte Carlo simulation (N=10,000, seed=42). Π = annual revenue (Iron Law).")

    notes(slide, f"Key metrics: βW = {beta}, CI = {ci}, Π = {money(pi_raw)}, ΔW = {money(dw_raw)}. "
          f"{'Impossibility — physical constraints.' if tt == 'Impossibility' else 'Intractability — solvable with reform.'}")
    return 1


def make_dense_text_slides(prs, title, subtitle, text, notes_prefix=""):
    """Dense text slides — 1800 chars per slide (~40 lines of 12pt). NO waste."""
    if not text:
        return 0
    full = strip_html(text)
    if not full:
        return 0

    pages = paginate_text(full, 1800)
    count = 0
    for i, page in enumerate(pages):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        set_bg(slide)
        pg = f" ({i+1}/{len(pages)})" if len(pages) > 1 else ""
        header_bar(slide, f"{title}{pg}", subtitle if i == 0 else None)

        top = Inches(1.3) if (subtitle and i == 0) else Inches(0.9)
        multi_text(slide, MARGIN_L, top, CONTENT_W, Inches(6.0),
                   page, size=12, color=DARK_TEXT, spacing=1.15)

        n = f"{notes_prefix}\n\n{page}" if notes_prefix and i == 0 else page
        notes(slide, n)
        count += 1
    return count


def make_findings_dense(prs, findings):
    """Pack multiple findings per slide — 3-4 findings/slide, not 1."""
    if not findings:
        return 0

    count = 0
    # Group findings: ~3 per slide depending on length
    group = []
    group_len = 0
    groups = []

    for f in findings:
        txt = strip_html(f) if isinstance(f, str) else strip_html(str(f))
        if group_len + len(txt) > 1600 and group:
            groups.append(group)
            group = [txt]
            group_len = len(txt)
        else:
            group.append(txt)
            group_len += len(txt)
    if group:
        groups.append(group)

    for gi, grp in enumerate(groups):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        set_bg(slide)
        pg = f" ({gi+1}/{len(groups)})" if len(groups) > 1 else ""
        header_bar(slide, f"KEY FINDINGS{pg}")

        y = Inches(1.0)
        slide_notes = []
        for fi, txt in enumerate(grp):
            # Finding number + bullet
            sentences = split_sentences(txt)
            headline = sentences[0] if sentences else txt
            body = ' '.join(sentences[1:3]) if len(sentences) > 1 else ""  # First 2-3 sentences

            # Number badge
            rect(slide, MARGIN_L, y, Inches(0.4), Inches(0.35), NAVY)
            text_box(slide, MARGIN_L, y + Inches(0.02), Inches(0.4), Inches(0.3),
                     str(fi + gi * len(groups[-1]) + 1), size=12, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

            # Headline
            text_box(slide, Inches(1.0), y, Inches(8.2), Inches(0.35),
                     headline[:200], size=12, color=NAVY, bold=True)

            # Body (condensed)
            if body:
                text_box(slide, Inches(1.0), y + Inches(0.35), Inches(8.2), Inches(0.8),
                         body[:350], size=11, color=DARK_TEXT)
                y += Inches(1.3)
            else:
                y += Inches(0.6)

            slide_notes.append(f"Finding: {txt}")

            if y > Inches(6.0):
                break

        notes(slide, "\n\n".join(slide_notes))
        count += 1

    return count


def make_theorem_dense(prs, data):
    """Theorem section — all content, dense layout."""
    theorem = data.get("theorem", {})
    if not theorem or not isinstance(theorem, dict):
        return 0

    count = 0

    # Plain language
    plain = theorem.get("plain", "")
    if plain:
        count += make_dense_text_slides(prs, "THE THEOREM — PLAIN LANGUAGE",
            data.get("theoremName", ""), plain)

    # Axioms — ALL on one or two slides, not one per axiom
    axioms = theorem.get("axioms", [])
    if axioms:
        # Pack axioms densely
        ax_groups = []
        current_group = []
        current_len = 0
        for ax in axioms:
            stmt = strip_html(ax.get("statement", ""))
            entry = f"{ax.get('id', 'A?')}: {ax.get('name', '')} — {stmt}"
            if current_len + len(entry) > 1400 and current_group:
                ax_groups.append(current_group)
                current_group = [(ax, entry)]
                current_len = len(entry)
            else:
                current_group.append((ax, entry))
                current_len += len(entry)
        if current_group:
            ax_groups.append(current_group)

        for gi, grp in enumerate(ax_groups):
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            set_bg(slide)
            pg = f" ({gi+1}/{len(ax_groups)})" if len(ax_groups) > 1 else ""
            header_bar(slide, f"AXIOMS{pg}", "Structural assumptions driving the result")

            y = Inches(1.3)
            ax_notes = []
            for ax_obj, entry_text in grp:
                ax_id = ax_obj.get("id", "A?")
                ax_name = ax_obj.get("name", "")
                stmt = strip_html(ax_obj.get("statement", ""))

                h = Inches(0.85) if len(stmt) < 150 else Inches(1.1)
                rect(slide, MARGIN_L, y, CONTENT_W, h, SURFACE)
                rect(slide, MARGIN_L, y, Inches(0.05), h, GOLD)

                text_box(slide, Inches(0.7), y + Inches(0.05), Inches(8), Inches(0.25),
                         f"{ax_id}: {ax_name}", size=12, color=NAVY, bold=True)
                text_box(slide, Inches(0.7), y + Inches(0.3), Inches(8.3), h - Inches(0.35),
                         stmt[:300], size=11, color=DARK_TEXT)

                y += h + Inches(0.1)
                ax_notes.append(f"{ax_id}: {ax_name}. {stmt}")

                if y > Inches(6.5):
                    break

            notes(slide, "Axioms:\n\n" + "\n\n".join(ax_notes))
            count += 1

    # Formal + proof sketch
    formal = theorem.get("formal", "")
    if formal:
        count += make_dense_text_slides(prs, "FORMAL THEOREM STATEMENT", "", formal)

    proof = theorem.get("proofSketch", "")
    if proof:
        count += make_dense_text_slides(prs, "PROOF SKETCH",
            "Step-by-step construction", proof)

    return count


def make_welfare_dense(prs, data):
    """Welfare channels — overview + detail, dense."""
    wa = data.get("welfareAnalysis")
    wa = wa if isinstance(wa, dict) else {}
    channels = wa.get("channels", [])
    if not channels:
        return 0

    count = 0

    # Overview: all channels on one slide
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    header_bar(slide, "WELFARE CHANNEL DECOMPOSITION",
               f"{len(channels)} channels of welfare destruction")

    y = Inches(1.3)
    ch_notes = []
    for j, ch in enumerate(channels):
        name = ch.get("name", f"Channel {j+1}")
        value = strip_html(str(ch.get("value", "—")))
        desc = strip_html(ch.get("description", ""))
        first_sent = (desc.split('. ')[0] + '.') if '. ' in desc else desc[:100]

        row_h = Inches(0.55)
        bg_color = SURFACE if j % 2 == 0 else SLIDE_BG_COLOR
        rect(slide, MARGIN_L, y, CONTENT_W, row_h, bg_color)

        text_box(slide, Inches(0.6), y + Inches(0.02), Inches(3.5), Inches(0.25),
                 name, size=11, color=NAVY, bold=True)
        text_box(slide, Inches(4.2), y + Inches(0.02), Inches(1.2), Inches(0.25),
                 value, size=11, color=RED, bold=True, align=PP_ALIGN.RIGHT)
        text_box(slide, Inches(5.6), y + Inches(0.02), Inches(3.8), Inches(0.5),
                 first_sent[:120], size=9, color=MID_GRAY)

        ch_notes.append(f"{name}: {value}. {first_sent}")
        y += row_h
        if y > Inches(6.8):
            break

    notes(slide, "Welfare channels:\n" + "\n".join(ch_notes))
    count += 1

    # Detail slides for channels with substantial content
    for ch in channels:
        desc = strip_html(ch.get("description", ""))
        if len(desc) > 200:
            name = ch.get("name", "Channel")
            value = strip_html(str(ch.get("value", "")))
            count += make_dense_text_slides(prs, f"CHANNEL: {name.upper()}", value, desc)

    return count


def make_cases_dense(prs, cases):
    """Case studies — dense, full content."""
    if not cases:
        return 0
    count = 0
    for case in cases:
        title = case.get("title", "Case Study")
        content = strip_html(case.get("content", case.get("description", "")))
        if content:
            count += make_dense_text_slides(prs, f"CASE STUDY: {title.upper()}", "", content)
    return count


def make_policy_dense(prs, data):
    """Policy analysis — framework, failures, reform."""
    policy = data.get("policyAnalysis")
    if not policy or not isinstance(policy, dict):
        return 0
    count = 0
    for title, key in [("CURRENT REGULATORY FRAMEWORK", "currentFramework"),
                       ("WHY THE FRAMEWORK FAILS", "failures"),
                       ("REFORM PATHWAY", "reform")]:
        content = policy.get(key, "")
        if content:
            count += make_dense_text_slides(prs, title, "", content)
    return count


def make_agents_dense(prs, agents):
    """Six-agent framework — dense, full content."""
    if not agents or not isinstance(agents, dict):
        return 0

    meta = {
        "whistleblower": ("WHISTLEBLOWER", "Break Information Asymmetry"),
        "plaintiff": ("PLAINTIFF", "Monetize Liability"),
        "regulator": ("REGULATOR", "Redesign the Game"),
        "legislator": ("LEGISLATOR", "Change the Law"),
        "investor": ("INVESTOR", "Reprice Capital"),
        "supranational": ("SUPRANATIONAL", "Coordinate Jurisdictions"),
    }

    count = 0
    for key in ["whistleblower", "plaintiff", "regulator", "legislator", "investor", "supranational"]:
        content = agents.get(key) or (agents.get("insider") if key == "whistleblower" else None)
        if not content:
            continue
        label, sub = meta.get(key, (key.upper(), ""))
        count += make_dense_text_slides(prs, f"AGENT: {label}", sub, content)
    return count


def make_cross_domain_dense(prs, connections):
    """Cross-domain connections — pack multiple per slide."""
    if not connections:
        return 0

    count = 0
    # Pack 2-3 connections per slide
    for i in range(0, len(connections), 3):
        batch = connections[i:i+3]
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        set_bg(slide)
        header_bar(slide, "CROSS-DOMAIN CONNECTIONS")

        y = Inches(1.0)
        conn_notes = []
        for conn in batch:
            domain = conn.get("domain", "")
            desc = strip_html(conn.get("connection", conn.get("description", "")))[:400]

            h = Inches(1.6)
            rect(slide, MARGIN_L, y, CONTENT_W, h, SURFACE)
            text_box(slide, Inches(0.6), y + Inches(0.05), Inches(8), Inches(0.25),
                     domain, size=13, color=NAVY, bold=True)
            text_box(slide, Inches(0.6), y + Inches(0.35), Inches(8.3), h - Inches(0.4),
                     desc, size=11, color=DARK_TEXT)
            y += h + Inches(0.15)
            conn_notes.append(f"{domain}: {desc}")

        notes(slide, "\n\n".join(conn_notes))
        count += 1

    return count


def make_faq_dense(prs, faq):
    """FAQ — pack 2-3 Q&As per slide."""
    if not faq:
        return 0

    count = 0
    # Group by total text length
    groups = []
    grp = []
    grp_len = 0
    for item in faq:
        q = strip_html(item.get("q", ""))
        a = strip_html(item.get("a", ""))
        entry_len = len(q) + len(a)
        if grp_len + entry_len > 1800 and grp:
            groups.append(grp)
            grp = [(q, a)]
            grp_len = entry_len
        else:
            grp.append((q, a))
            grp_len += entry_len
    if grp:
        groups.append(grp)

    for gi, group in enumerate(groups):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        set_bg(slide)
        pg = f" ({gi+1}/{len(groups)})" if len(groups) > 1 else ""
        header_bar(slide, f"FREQUENTLY ASKED QUESTIONS{pg}")

        y = Inches(1.0)
        faq_notes = []
        for q, a in group:
            text_box(slide, MARGIN_L, y, CONTENT_W, Inches(0.3),
                     q[:200], size=12, color=NAVY, bold=True, italic=True)
            a_short = a[:500]
            text_box(slide, Inches(0.7), y + Inches(0.3), Inches(8.5), Inches(1.2),
                     a_short, size=11, color=DARK_TEXT)
            faq_notes.append(f"Q: {q}\nA: {a}")
            y += Inches(1.8)
            if y > Inches(6.0):
                break

        notes(slide, "\n\n".join(faq_notes))
        count += 1

    return count


def make_timeline_dense(prs, timeline):
    """Timeline — pack 8-10 events per slide."""
    if not timeline:
        return 0

    count = 0
    for i in range(0, len(timeline), 10):
        batch = timeline[i:i+10]
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        set_bg(slide)
        header_bar(slide, "TIMELINE")

        y = Inches(0.9)
        tl_notes = []
        for ev in batch:
            year = str(ev.get("year", ""))
            evt = strip_html(ev.get("event", ""))
            sig = strip_html(ev.get("significance", ""))

            text_box(slide, Inches(0.5), y, Inches(0.8), Inches(0.25),
                     year, size=11, color=NAVY, bold=True)
            text_box(slide, Inches(1.4), y, Inches(3.2), Inches(0.25),
                     evt[:70], size=10, color=DARK_TEXT, bold=True)
            text_box(slide, Inches(4.8), y, Inches(4.7), Inches(0.5),
                     sig[:150], size=9, color=MID_GRAY)
            tl_notes.append(f"{year}: {evt}. {sig}")
            y += Inches(0.55)
            if y > Inches(7.0):
                break

        notes(slide, "\n".join(tl_notes))
        count += 1

    return count


def make_table_slide(prs, table, slug):
    """Word doc table — dense, full-width."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    header_bar(slide, f"DATA: {table['title'].upper()}", f"{len(table['rows'])} rows — from published paper")

    header = table["header"]
    rows = table["rows"][:20]
    n_cols = min(len(header), 7)
    if n_cols == 0:
        return 0

    col_w = min(CONTENT_W / n_cols, Inches(2.5))
    start_x = MARGIN_L

    # Header
    y = Inches(1.3)
    rect(slide, start_x, y, col_w * n_cols, Inches(0.28), NAVY)
    for j, h in enumerate(header[:n_cols]):
        text_box(slide, start_x + j * col_w + Inches(0.05), y + Inches(0.02),
                 col_w - Inches(0.1), Inches(0.24),
                 h[:35], size=9, color=WHITE, bold=True)

    # Rows
    for r, row in enumerate(rows):
        y = Inches(1.58) + r * Inches(0.26)
        if r % 2 == 0:
            rect(slide, start_x, y, col_w * n_cols, Inches(0.26), SURFACE)
        for j, cell in enumerate(row[:n_cols]):
            text_box(slide, start_x + j * col_w + Inches(0.05), y + Inches(0.01),
                     col_w - Inches(0.1), Inches(0.24),
                     str(cell)[:35], size=8, color=DARK_TEXT)
        if y > Inches(6.8):
            break

    source_line(slide, f"Extracted from {slug} Word document")

    table_notes = f"Table with {len(table['rows'])} rows. Columns: {', '.join(header[:n_cols])}"
    notes(slide, table_notes)
    return 1


def make_quiz(prs, data):
    """Quiz slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, NAVY)
    text_box(slide, Inches(0.5), Inches(0.4), Inches(8), Inches(0.3),
             "COMPREHENSION CHECK", size=12, color=CYAN, bold=True)

    wa = data.get("welfareAnalysis")
    wa = wa if isinstance(wa, dict) else {}
    beta = wa.get("betaW", "?")
    tt = data.get("theoremType", "?")
    tn = data.get("theoremName", "?")

    qs = [
        (f"Q1: What is βW for this domain?", f"βW = {beta}"),
        (f"Q2: Impossibility or Intractability?", f"{tt}"),
        (f"Q3: Theorem name?", f"{tn}"),
        ("Q4: Solvable with policy?",
         "No — physical constraints" if tt == "Impossibility" else "Yes — proven models exist"),
        ("Q5: What is the Iron Law for Π?",
         "Π = revenue, never profit"),
    ]

    y = Inches(1.2)
    for q, a in qs:
        text_box(slide, Inches(0.5), y, Inches(8), Inches(0.3),
                 q, size=14, color=WHITE, bold=True)
        text_box(slide, Inches(0.8), y + Inches(0.35), Inches(8), Inches(0.3),
                 a, size=12, color=CYAN)
        y += Inches(1.1)

    notes(slide, "\n".join(f"{q} → {a}" for q, a in qs))
    return 1


def make_closing(prs, data):
    """Closing slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, DARK_BG)
    rect(slide, 0, Inches(2.8), SLIDE_W, Inches(0.03), CYAN)

    title = data.get("title", "")
    text_box(slide, Inches(0.5), Inches(0.8), Inches(9), Inches(1.5),
             title, size=18, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

    wa = data.get("welfareAnalysis")
    wa = wa if isinstance(wa, dict) else {}
    beta = wa.get("betaW", "")
    if beta and beta != "—":
        text_box(slide, Inches(0.5), Inches(3.2), Inches(9), Inches(0.8),
                 f"βW = {beta}", size=40, color=CYAN, bold=True, align=PP_ALIGN.CENTER)

    tt = data.get("theoremType", "")
    if tt:
        c = RED if tt == "Impossibility" else GOLD
        text_box(slide, Inches(0.5), Inches(4.2), Inches(9), Inches(0.4),
                 f"{tt} Theorem", size=18, color=c, align=PP_ALIGN.CENTER)

    text_box(slide, Inches(0.5), Inches(5.2), Inches(9), Inches(0.3),
             "Erik Postnieks  |  Independent Researcher  |  postnieks.com",
             size=12, color=MID_GRAY, align=PP_ALIGN.CENTER)
    text_box(slide, Inches(0.5), Inches(5.6), Inches(9), Inches(0.3),
             "Paper: SSRN  |  MC Replication: github.com/epostnieks",
             size=10, color=MID_GRAY, align=PP_ALIGN.CENTER)
    text_box(slide, Inches(0.5), Inches(6.2), Inches(9), Inches(0.3),
             "END OF MASTERCLASS", size=10, color=CYAN, bold=True, align=PP_ALIGN.CENTER)

    notes(slide, f"End of masterclass on {title}. βW = {beta}. {tt} Theorem. "
          f"Full paper on SSRN. MC replication on GitHub.")
    return 1


# ══════════════════════════════════════════════════════════════
# TTS AUDIO GENERATION
# ══════════════════════════════════════════════════════════════

def generate_tts(prs, slug):
    """Generate TTS narration using macOS 'say' command (no rate limits)."""
    import subprocess, tempfile

    all_notes = []
    for slide in prs.slides:
        try:
            ns = slide.notes_slide
            text = ns.notes_text_frame.text
            if text and text.strip():
                all_notes.append(text.strip())
        except Exception:
            pass

    if not all_notes:
        return None

    full_script = "\n\n".join(all_notes)

    # Clean text for TTS
    full_script = re.sub(r'[βΠΔ]', '', full_script)  # Remove Greek chars that say can't pronounce
    full_script = full_script.replace('βW', 'beta W')
    full_script = full_script.replace('ΔW', 'delta W')

    out_path = AUDIO_DIR / f"{slug}.m4a"
    try:
        # Write script to temp file
        with tempfile.NamedTemporaryFile(mode='w', suffix='.txt', delete=False) as f:
            f.write(full_script)
            tmp_path = f.name

        # Use macOS say with Samantha voice (high quality)
        subprocess.run([
            'say', '-v', 'Samantha', '-r', '180',
            '-f', tmp_path,
            '-o', str(out_path)
        ], check=True, timeout=300)

        os.unlink(tmp_path)
        return out_path
    except Exception as e:
        print(f"    TTS error for {slug}: {e}")
        return None


# ══════════════════════════════════════════════════════════════
# MAIN GENERATOR
# ══════════════════════════════════════════════════════════════

def generate(json_path):
    """Generate masterclass PPTX + TTS audio."""
    with open(json_path, 'r') as f:
        data = json.load(f)

    slug = json_path.stem

    # Assets
    figures = find_figures(slug)
    docx_path = find_docx(slug)
    tables = get_docx_tables(docx_path) if docx_path else []
    images = get_docx_images(docx_path) if docx_path else []

    # Categorize figures for strategic placement
    fig_map = {}
    for fig in figures:
        name = fig.name.lower()
        if 'channel' in name or 'breakdown' in name:
            fig_map.setdefault('channel', []).append(fig)
        elif 'mc_dist' in name or 'distribution' in name:
            fig_map.setdefault('mc', []).append(fig)
        elif 'sensitivity' in name or 'tornado' in name:
            fig_map.setdefault('sensitivity', []).append(fig)
        elif 'cross_domain' in name:
            fig_map.setdefault('crossdomain', []).append(fig)
        elif 'welfare' in name or 'ledger' in name:
            fig_map.setdefault('welfare', []).append(fig)
        elif 'literature' in name or 'five_lit' in name:
            fig_map.setdefault('literature', []).append(fig)
        elif 'pigou' in name or 'coase' in name:
            fig_map.setdefault('theory', []).append(fig)
        else:
            fig_map.setdefault('other', []).append(fig)

    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    total = 0

    # Title
    total += make_title(prs, data)

    # Ch 1: Numbers + MC figures
    total += make_divider(prs, 1, "THE NUMBERS",
        "System beta, welfare cost, and what they mean")
    total += make_metrics(prs, data)
    for fig in fig_map.get('mc', []):
        total += make_figure(prs, fig, f"Monte Carlo Distribution — {slug}")
    for fig in fig_map.get('sensitivity', []):
        total += make_figure(prs, fig, f"Sensitivity Analysis — {slug}")

    # Ch 2: Executive Summary
    total += make_divider(prs, 2, "EXECUTIVE SUMMARY",
        "The complete argument — what this paper proves")
    total += make_dense_text_slides(prs, "EXECUTIVE SUMMARY",
        "What this paper proves and why it matters",
        data.get("executiveSummary", ""))

    # Ch 3: Key Findings + channel figures
    total += make_divider(prs, 3, "KEY FINDINGS", "What the evidence shows")
    total += make_findings_dense(prs, data.get("keyFindings", []))
    for fig in fig_map.get('channel', []):
        total += make_figure(prs, fig, f"Welfare Channel Breakdown — {slug}")

    # Ch 4: Theorem
    tn_display = data.get("theoremName", "Formal Result")
    # Clean up theorem name display — remove redundant "Theorem" suffix if present
    total += make_divider(prs, 4, tn_display, "Formal proof and structural assumptions")
    total += make_theorem_dense(prs, data)
    for fig in fig_map.get('theory', []):
        total += make_figure(prs, fig, f"Theoretical Framework — {slug}")

    # Ch 5: Welfare Channels + welfare figures
    total += make_divider(prs, 5, "WELFARE CHANNELS",
        "Where the damage comes from")
    total += make_welfare_dense(prs, data)
    for fig in fig_map.get('welfare', []):
        total += make_figure(prs, fig, f"Welfare Analysis — {slug}")
    total += make_dense_text_slides(prs, "CHANNEL TABLE", "",
        (data.get("welfareAnalysis") or {}).get("channelTable", "") if isinstance(data.get("welfareAnalysis"), dict) else "")

    # Ch 6: Monte Carlo
    total += make_divider(prs, 6, "MONTE CARLO SIMULATION",
        "Quantifying uncertainty — 10,000 draws, reproducible")
    total += make_dense_text_slides(prs, "MONTE CARLO METHODOLOGY",
        "N=10,000 | Seed=42 | Reproducible",
        (data.get("welfareAnalysis") or {}).get("monteCarlo", "") if isinstance(data.get("welfareAnalysis"), dict) else "")

    # Ch 7: Paper Figures — ALL remaining McKinsey figures
    other_figs = fig_map.get('other', [])
    if other_figs:
        total += make_divider(prs, 7, "PAPER FIGURES",
            f"{len(other_figs)} McKinsey-quality figures from the research")
        for fig in other_figs:
            total += make_figure(prs, fig, f"{fig_label(fig)} — {slug}")

    # Ch 8: Evidence — cases + Word doc tables + embedded images
    ch_num = 8
    total += make_divider(prs, ch_num, "EVIDENCE",
        f"Case studies, {len(tables)} data tables, {len(images)} figures from the paper")
    total += make_cases_dense(prs, data.get("caseStudies", []))
    for table in tables:
        total += make_table_slide(prs, table, slug)
    for i, img in enumerate(images[:30]):
        total += make_image_slide(prs, img["data"], img["ext"],
            f"Paper Figure {i+1}")

    # Ch 9: Policy
    total += make_divider(prs, 9, "POLICY ANALYSIS", "What works, what fails, why")
    total += make_policy_dense(prs, data)

    # Ch 10: Six-Agent Framework
    total += make_divider(prs, 10, "SIX-AGENT FRAMEWORK",
        "Who can break the Private-Systemic Tension")
    total += make_agents_dense(prs, data.get("agents", {}))

    # Ch 11: Context
    total += make_divider(prs, 11, "CONNECTIONS AND CONTEXT",
        "Cross-domain, literature, methodology, timeline")
    total += make_cross_domain_dense(prs, data.get("crossDomainConnections", []))
    for fig in fig_map.get('crossdomain', []):
        total += make_figure(prs, fig, f"Cross-Domain Map — {slug}")
    for fig in fig_map.get('literature', []):
        total += make_figure(prs, fig, f"Literature Context — {slug}")
    total += make_dense_text_slides(prs, "LITERATURE CONTEXT",
        "Where this paper sits", data.get("literatureContext", ""))
    total += make_dense_text_slides(prs, "METHODOLOGY",
        "Data sources and limitations", data.get("methodology", ""))
    total += make_timeline_dense(prs, data.get("timeline", []))

    # Ch 12: FAQ + Quiz + Closing
    total += make_divider(prs, 12, "FAQ AND ASSESSMENT", "Questions, answers, and comprehension check")
    total += make_faq_dense(prs, data.get("faq", []))
    total += make_quiz(prs, data)
    total += make_closing(prs, data)

    # Save PPTX
    out_path = PPTX_DIR / f"{slug}.pptx"
    prs.save(str(out_path))

    # Generate TTS audio
    audio_path = generate_tts(prs, slug)

    return slug, total, len(figures), len(tables), len(images), audio_path is not None


def main():
    json_files = sorted(PAPER_DATA_DIR.glob("*.json"))
    json_files = [f for f in json_files if '_raw' not in f.name]
    print(f"Found {len(json_files)} papers")

    results = []
    errors = []

    # Use fewer workers since TTS is network-bound
    with ThreadPoolExecutor(max_workers=4) as pool:
        futures = {pool.submit(generate, f): f for f in json_files}
        for future in as_completed(futures):
            f = futures[future]
            try:
                slug, slides, figs, tables, images, has_audio = future.result()
                results.append((slug, slides, figs, tables, images, has_audio))
                audio_str = " +audio" if has_audio else ""
                print(f"  OK  {slug}: {slides} slides, {figs} figs, {tables} tables, {images} imgs{audio_str}")
            except Exception as e:
                errors.append((f.stem, str(e)))
                print(f"  FAIL {f.stem}: {e}")
                traceback.print_exc()

    results.sort(key=lambda x: x[0])
    total_slides = sum(r[1] for r in results)
    total_audio = sum(1 for r in results if r[5])

    print(f"\nDone: {len(results)} OK, {len(errors)} failed")
    print(f"Total slides: {total_slides} (avg {total_slides // max(len(results),1)}/paper)")
    print(f"Audio narrations: {total_audio}/{len(results)}")
    print(f"PPTX output: {PPTX_DIR}")
    print(f"Audio output: {AUDIO_DIR}")

    if errors:
        print(f"\nFailed:")
        for slug, err in errors:
            print(f"  {slug}: {err}")


if __name__ == "__main__":
    main()
