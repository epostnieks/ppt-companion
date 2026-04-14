#!/usr/bin/env python3
"""
PPTX Generator V4 — MASTERCLASS EDITION
150-200 slides per paper. McKinsey figure PNGs as full-bleed anchor slides.
Word doc tables rendered as slides. Rich speaker notes for TTS narration.
Full 20K-word Gemini summary content — NO TRUNCATION.

Architecture:
  1. Load Gemini JSON summary (20K words)
  2. Load McKinsey figure PNGs from deep-research/output/{slug}/figures/
  3. Extract tables from Word docs via python-docx
  4. Build 12-chapter masterclass with figures as anchors
  5. Every slide gets speaker notes for TTS narration
"""

import json, os, re, sys, glob, traceback
from concurrent.futures import ThreadPoolExecutor, as_completed
from pathlib import Path
from html.parser import HTMLParser
from io import BytesIO

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

try:
    from docx import Document as DocxDocument
    HAS_DOCX = True
except ImportError:
    HAS_DOCX = False

# ── Paths ──
PAPER_DATA_DIR = Path(__file__).parent.parent / "src" / "paperData"
OUTPUT_DIR = Path(__file__).parent.parent / "public" / "pptx"
OUTPUT_DIR.mkdir(parents=True, exist_ok=True)
BASE = Path("/Users/erikpostnieks/Projects")
DR_OUTPUT = BASE / "deep-research" / "output"

# ── McKinsey Colors ──
NAVY = RGBColor(0x1F, 0x4E, 0x79)
CREAM = RGBColor(0xFB, 0xF9, 0xF4)
TAN = RGBColor(0xF5, 0xEF, 0xE0)
GOLD = RGBColor(0xC4, 0x97, 0x2A)
SOURCE_BLUE = RGBColor(0x3D, 0x6B, 0x8E)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x00, 0x00, 0x00)
DARK_TEXT = RGBColor(0x33, 0x33, 0x33)
LIGHT_GRAY = RGBColor(0xE8, 0xE0, 0xD0)
MID_GRAY = RGBColor(0x88, 0x88, 0x88)
RED = RGBColor(0xCC, 0x33, 0x33)
GREEN = RGBColor(0x22, 0x88, 0x55)
DARK_BG = RGBColor(0x0F, 0x17, 0x2A)
CYAN = RGBColor(0x00, 0xD4, 0xFF)
SURFACE = RGBColor(0xF7, 0xF5, 0xF0)
SLIDE_BG = RGBColor(0xFF, 0xFF, 0xFF)

FONT = "Liberation Sans"
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


# ══════════════════════════════════════════════════════════════
# TEXT UTILITIES
# ══════════════════════════════════════════════════════════════

class HTMLStripper(HTMLParser):
    def __init__(self):
        super().__init__()
        self.result = []
    def handle_data(self, d):
        self.result.append(d)
    def handle_starttag(self, tag, attrs):
        if tag in ('td', 'th'):
            self.result.append(' | ')
        if tag == 'tr':
            self.result.append('\n')
        if tag == 'li':
            self.result.append('\n• ')
        if tag == 'br':
            self.result.append('\n')
        if tag == 'p':
            self.result.append('\n\n')
    def get_text(self):
        return ''.join(self.result).strip()

def strip_html(text):
    if not text or not isinstance(text, str):
        return str(text) if text else ""
    s = HTMLStripper()
    s.feed(text)
    return s.get_text()

def split_sentences(text):
    """Split text into sentences."""
    text = strip_html(text)
    # Split on sentence boundaries but keep the delimiter
    parts = re.split(r'(?<=[.!?])\s+', text)
    return [p.strip() for p in parts if p.strip()]

def split_into_chunks(text, max_chars=600):
    """Split text into presentation-sized chunks at sentence boundaries."""
    text = strip_html(text)
    if not text:
        return []
    sentences = split_sentences(text)
    chunks = []
    current = ""
    for s in sentences:
        if len(current) + len(s) > max_chars and current:
            chunks.append(current.strip())
            current = s
        else:
            current = (current + " " + s).strip()
    if current:
        chunks.append(current.strip())
    return chunks

def split_into_slide_chunks(text, max_chars=500):
    """Split for slide body text — shorter chunks for readability."""
    return split_into_chunks(text, max_chars)

def split_into_notes_chunks(text, max_chars=2000):
    """Split for speaker notes — longer chunks for narration."""
    return split_into_chunks(text, max_chars)


# ══════════════════════════════════════════════════════════════
# SHAPE HELPERS
# ══════════════════════════════════════════════════════════════

def set_slide_bg(slide, color=SLIDE_BG):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_rect(slide, left, top, width, height, color):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    shape.shadow.inherit = False
    return shape

def add_text(slide, left, top, width, height, text, size=14,
             color=DARK_TEXT, bold=False, italic=False, align=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = str(text)
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.italic = italic
    p.font.name = FONT
    p.alignment = align
    return txBox

def add_multiline(slide, left, top, width, height, text, size=12,
                  color=DARK_TEXT, line_spacing=1.5):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    paragraphs = str(text).split('\n')
    for i, para_text in enumerate(paragraphs):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = para_text
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.font.name = FONT
        p.space_after = Pt(size * 0.5)
    return txBox

def add_notes(slide, text):
    """Add speaker notes for TTS narration."""
    if not text:
        return
    notes_slide = slide.notes_slide
    tf = notes_slide.notes_text_frame
    tf.text = str(text)

def mckinsey_header(slide, title, subtitle=None):
    """Five-layer McKinsey header."""
    add_rect(slide, 0, 0, SLIDE_W, Inches(0.9), NAVY)
    add_text(slide, Inches(0.6), Inches(0.15), Inches(12), Inches(0.6),
             title, size=26, color=WHITE, bold=True)
    if subtitle:
        add_rect(slide, 0, Inches(0.9), SLIDE_W, Inches(0.5), CREAM)
        add_text(slide, Inches(0.6), Inches(0.95), Inches(12), Inches(0.4),
                 subtitle, size=14, color=DARK_TEXT, italic=True)

def key_insight_box(slide, text, y_pos=Inches(5.8)):
    """KEY INSIGHT box with gold left border."""
    add_rect(slide, Inches(0.5), y_pos, Inches(12.3), Inches(1.2), TAN)
    add_rect(slide, Inches(0.5), y_pos, Inches(0.06), Inches(1.2), GOLD)
    add_text(slide, Inches(0.8), y_pos + Inches(0.05), Inches(2), Inches(0.3),
             "KEY INSIGHT", size=11, color=GOLD, bold=True)
    add_text(slide, Inches(0.8), y_pos + Inches(0.3), Inches(11.5), Inches(0.8),
             text[:300], size=13, color=DARK_TEXT)

def source_line(slide, text, y_pos=Inches(7.1)):
    add_text(slide, Inches(0.6), y_pos, Inches(12), Inches(0.3),
             f"Source: {text}", size=11, color=SOURCE_BLUE)

def slide_number(slide, num, total):
    """Add slide number."""
    add_text(slide, Inches(12.0), Inches(7.15), Inches(1.2), Inches(0.3),
             f"{num}/{total}", size=9, color=MID_GRAY, align=PP_ALIGN.RIGHT)


# ══════════════════════════════════════════════════════════════
# FIGURE & TABLE EXTRACTION
# ══════════════════════════════════════════════════════════════

def find_figure_pngs(slug):
    """Find McKinsey figure PNGs from deep-research output."""
    # deep-research uses underscores, convert from hyphenated slug
    dr_slug = slug.replace("-", "_")
    fig_dir = DR_OUTPUT / dr_slug / "figures"
    if not fig_dir.exists():
        # Try hyphenated too
        fig_dir = DR_OUTPUT / slug / "figures"
    if not fig_dir.exists():
        return []

    pngs = sorted(fig_dir.glob("*.png"))
    # Filter out versioned duplicates (_v0, _v1)
    filtered = []
    seen_bases = set()
    for p in pngs:
        base = re.sub(r'_v\d+\.png$', '.png', p.name)
        if base not in seen_bases:
            filtered.append(p)
            seen_bases.add(base)
    return filtered

def find_docx(slug):
    """Find the Word document for a paper."""
    patterns = [
        f"{BASE}/sapm-{slug}/paper/{slug}-current.docx",
        f"{BASE}/sapm-{slug}/paper/*-current.docx",
        f"{BASE}/sapm-{slug}/paper/*-final.docx",
        f"{BASE}/sapm-{slug}/paper/*.docx",
    ]
    for pat in patterns:
        matches = glob.glob(pat)
        matches = [m for m in matches if '~' not in m and 'template' not in m.lower()]
        if matches:
            return matches[0]
    return None

def extract_docx_tables(docx_path, max_tables=50):
    """Extract ALL tables from Word document."""
    if not HAS_DOCX or not docx_path:
        return []
    try:
        doc = DocxDocument(docx_path)
    except Exception:
        return []

    tables = []
    for i, table in enumerate(doc.tables[:max_tables]):
        rows = []
        for row in table.rows:
            cells = [cell.text.strip() for cell in row.cells]
            rows.append(cells)
        if rows and len(rows) > 1:
            # Try to find table title from preceding paragraph
            title = f"Table {i+1}"
            tables.append({"title": title, "header": rows[0], "rows": rows[1:]})
    return tables

def extract_docx_images(docx_path):
    """Extract embedded images from Word document."""
    if not HAS_DOCX or not docx_path:
        return []
    try:
        doc = DocxDocument(docx_path)
    except Exception:
        return []

    images = []
    for rel in doc.part.rels.values():
        if "image" in rel.reltype:
            try:
                img_data = rel.target_part.blob
                ext = rel.target_part.content_type.split('/')[-1]
                if ext in ('png', 'jpeg', 'jpg', 'gif', 'tiff'):
                    images.append({"data": img_data, "ext": ext})
            except Exception:
                pass
    return images

def figure_label(fig_path):
    """Generate a human-readable label from figure filename."""
    name = fig_path.stem
    # Remove figNN_ prefix
    name = re.sub(r'^fig\d+_', '', name)
    # Convert underscores to spaces, title case
    name = name.replace('_', ' ').title()
    return name


# ══════════════════════════════════════════════════════════════
# SLIDE BUILDERS — MASTERCLASS EDITION
# ══════════════════════════════════════════════════════════════

def build_title_slide(prs, data):
    """Opening title slide — dark, cinematic."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)

    add_rect(slide, 0, 0, SLIDE_W, Inches(0.06), CYAN)

    add_text(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(0.4),
             "SYSTEM ASSET PRICING MODEL — MASTERCLASS", size=14, color=CYAN, bold=True)

    title = data.get("title", data.get("slug", ""))
    add_text(slide, Inches(0.8), Inches(1.5), Inches(11.5), Inches(2.5),
             title, size=40, color=WHITE, bold=True)

    tt = data.get("theoremType", "")
    tn = data.get("theoremName", "")
    if tt:
        badge_color = RED if tt == "Impossibility" else GOLD
        add_text(slide, Inches(0.8), Inches(4.2), Inches(6), Inches(0.5),
                 f"{tt} Theorem: {tn}", size=18, color=badge_color, bold=True)

    epigraph = data.get("epigraph", "")
    if epigraph:
        add_text(slide, Inches(0.8), Inches(5.2), Inches(10), Inches(1.2),
                 f'"{strip_html(epigraph)}"', size=16, color=MID_GRAY, italic=True)

    add_text(slide, Inches(0.8), Inches(6.6), Inches(10), Inches(0.4),
             "Erik Postnieks  |  Independent Researcher  |  2026", size=12, color=MID_GRAY)

    notes = f"""Welcome to the SAPM Masterclass on {title}.

{strip_html(epigraph) if epigraph else ''}

This is a {tt} Theorem paper — {tn}. Over the next slides, we will walk through the complete argument: the system beta, the welfare channels, the formal proof, the case studies, the six-agent conflictoring framework, and the policy implications.

This presentation contains approximately 150 to 200 slides with full narration, drawn from the published paper and its 20,000-word analytical summary."""
    add_notes(slide, notes)
    return 1


def build_chapter_divider(prs, num, title, subtitle="", notes_text=""):
    """Chapter divider slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, NAVY)

    add_text(slide, Inches(0.8), Inches(2.0), Inches(11), Inches(0.6),
             f"CHAPTER {num}", size=16, color=CYAN, bold=True)
    add_text(slide, Inches(0.8), Inches(2.8), Inches(11), Inches(1.5),
             title, size=36, color=WHITE, bold=True)
    if subtitle:
        add_text(slide, Inches(0.8), Inches(4.5), Inches(10), Inches(1.0),
                 subtitle, size=16, color=CREAM, italic=True)

    add_notes(slide, notes_text or f"Chapter {num}: {title}. {subtitle}")
    return 1


def build_figure_slide(prs, fig_path, caption="", notes_text=""):
    """Full-bleed McKinsey figure PNG as a slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)

    try:
        # Add figure as large centered image
        # Target: fill most of the slide while maintaining aspect ratio
        pic = slide.shapes.add_picture(
            str(fig_path),
            left=Inches(0.3),
            top=Inches(0.3),
            width=Inches(12.7),
            height=Inches(6.5),
        )
    except Exception as e:
        # Fallback: show label
        add_text(slide, Inches(1), Inches(3), Inches(11), Inches(1),
                 f"[Figure: {fig_path.name}]", size=20, color=MID_GRAY, align=PP_ALIGN.CENTER)

    if caption:
        add_text(slide, Inches(0.5), Inches(7.0), Inches(12), Inches(0.3),
                 caption, size=10, color=SOURCE_BLUE)

    add_notes(slide, notes_text or f"This figure shows {figure_label(fig_path)}. {caption}")
    return 1


def build_embedded_image_slide(prs, img_data, img_ext, title="", notes_text=""):
    """Embedded Word doc image as a slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)

    if title:
        mckinsey_header(slide, title)

    try:
        img_stream = BytesIO(img_data)
        top = Inches(1.2) if title else Inches(0.3)
        height = Inches(5.8) if title else Inches(6.5)
        pic = slide.shapes.add_picture(
            img_stream,
            left=Inches(0.5),
            top=top,
            width=Inches(12.3),
            height=height,
        )
    except Exception:
        add_text(slide, Inches(1), Inches(3), Inches(11), Inches(1),
                 f"[Embedded image — {img_ext}]", size=20, color=MID_GRAY, align=PP_ALIGN.CENTER)

    add_notes(slide, notes_text or f"This is a figure from the published paper. {title}")
    return 1


def build_metrics_slide(prs, data):
    """Key metrics dashboard."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    mckinsey_header(slide, "KEY METRICS DASHBOARD", "System Beta and Welfare Analysis")

    wa = data.get("welfareAnalysis")
    wa = wa if isinstance(wa, dict) else {}

    beta = wa.get("betaW", data.get("betaW", "—"))
    ci = wa.get("ci90", data.get("ci90", "—"))
    pi_raw = wa.get("pi", "—")
    dw_raw = wa.get("deltaW", "—")

    def short_money(val):
        if not val or val == "—":
            return "—"
        val = str(val)
        m = re.search(r'[-−]?\$[\d,.]+\s*(billion|trillion|million)?', val, re.IGNORECASE)
        return m.group(0) if m else val[:60]

    metrics = [
        ("βW (System Beta)", str(beta), "Welfare destruction per dollar of private gain", NAVY),
        ("90% Confidence Interval", str(ci), "Monte Carlo simulation bounds (N=10,000)", DARK_TEXT),
        ("Π (Annual Revenue)", short_money(pi_raw), "Private payoff — always revenue, never profit (Iron Law)", GREEN),
        ("ΔW (Welfare Cost)", short_money(dw_raw), "Annual system welfare destruction", RED),
    ]

    y_start = Inches(1.8)
    card_w = Inches(5.8)
    card_h = Inches(1.2)

    for i, (label, value, desc, accent) in enumerate(metrics):
        col = i % 2
        row = i // 2
        x = Inches(0.5) + col * (card_w + Inches(0.4))
        y = y_start + row * (card_h + Inches(0.3))

        add_rect(slide, x, y, card_w, card_h, SURFACE)
        add_text(slide, x + Inches(0.2), y + Inches(0.08), card_w - Inches(0.4), Inches(0.3),
                 label, size=12, color=MID_GRAY, bold=True)
        val_size = 28 if len(value) < 15 else 20 if len(value) < 30 else 14
        add_text(slide, x + Inches(0.2), y + Inches(0.38), card_w - Inches(0.4), Inches(0.5),
                 value, size=val_size, color=accent, bold=True)
        add_text(slide, x + Inches(0.2), y + Inches(0.88), card_w - Inches(0.4), Inches(0.25),
                 desc, size=10, color=MID_GRAY)

    tt = data.get("theoremType", "")
    if tt:
        key_insight_box(slide,
            f"Classification: {tt} Theorem — {'Physical or biological constraints prevent welfare internalization. No policy can solve this.' if tt == 'Impossibility' else 'Institutional constraints are removable with proven policy reform. At least one country has already solved this.'}",
            y_pos=Inches(5.0))

    source_line(slide, "Monte Carlo simulation (N=10,000, seed=42). Π = annual industry revenue (Iron Law).")

    notes = f"""Here are the key metrics for this domain.

The system beta, beta W, is {beta}. This means that for every dollar of private revenue the industry earns, it destroys {beta} dollars of social welfare.

The 90 percent confidence interval from our Monte Carlo simulation is {ci}. The private payoff Pi is {short_money(pi_raw)} in annual revenue. The total welfare destruction Delta W is {short_money(dw_raw)} per year.

This is classified as a {tt} Theorem. {'This means the welfare cost is driven by physical or biological constraints that no policy can overcome.' if tt == 'Impossibility' else 'This means the welfare cost is driven by institutional constraints. At least one country has demonstrated a proven reform that substantially reduces the welfare gap.'}"""
    add_notes(slide, notes)
    return 1


def build_text_slides(prs, title, subtitle, text, notes_prefix="",
                      max_slide_chars=500, max_notes_chars=1500):
    """Build multiple slides from long text content. Returns slide count.
    Slide shows readable excerpt; speaker notes contain full narration."""
    if not text:
        return 0

    full_text = strip_html(text)
    if not full_text:
        return 0

    slide_chunks = split_into_slide_chunks(full_text, max_slide_chars)
    notes_chunks = split_into_notes_chunks(full_text, max_notes_chars)

    count = 0
    for i, chunk in enumerate(slide_chunks):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        set_slide_bg(slide)
        page_label = f" ({i+1}/{len(slide_chunks)})" if len(slide_chunks) > 1 else ""
        mckinsey_header(slide, f"{title}{page_label}", subtitle)

        add_multiline(slide, Inches(0.6), Inches(1.6), Inches(12), Inches(5.2),
                      chunk, size=14, color=DARK_TEXT, line_spacing=1.6)

        # Speaker notes: corresponding notes chunk (may be longer than slide text)
        note_idx = min(i, len(notes_chunks) - 1) if notes_chunks else 0
        notes_text = notes_chunks[note_idx] if notes_chunks else chunk
        if notes_prefix:
            notes_text = f"{notes_prefix}\n\n{notes_text}"
        add_notes(slide, notes_text)
        count += 1

    return count


def build_finding_slides(prs, findings):
    """One slide per finding with full text + notes."""
    if not findings:
        return 0

    count = 0
    for i, finding in enumerate(findings):
        text = strip_html(finding) if isinstance(finding, str) else strip_html(str(finding))
        if not text:
            continue

        sentences = split_sentences(text)
        headline = sentences[0] if sentences else ""
        body = ' '.join(sentences[1:]) if len(sentences) > 1 else ""

        # If body is long, split across multiple slides
        if len(body) > 600:
            body_chunks = split_into_slide_chunks(body, 500)
        else:
            body_chunks = [body] if body else [""]

        for j, chunk in enumerate(body_chunks):
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            set_slide_bg(slide)
            page_label = f" ({j+1}/{len(body_chunks)})" if len(body_chunks) > 1 else ""
            mckinsey_header(slide, f"KEY FINDING {i+1} OF {len(findings)}{page_label}")

            add_text(slide, Inches(0.6), Inches(1.3), Inches(12), Inches(0.8),
                     headline, size=18, color=NAVY, bold=True)

            if chunk:
                add_multiline(slide, Inches(0.6), Inches(2.3), Inches(12), Inches(3.5),
                              chunk, size=13, color=DARK_TEXT)

            if j == 0:
                key_insight_box(slide, headline, y_pos=Inches(6.0))

            notes = f"Key Finding {i+1}: {headline}\n\n{chunk}"
            add_notes(slide, notes)
            count += 1

    return count


def build_theorem_slides(prs, data):
    """Theorem section — plain language, axioms, formal, proof. Full content."""
    theorem = data.get("theorem", {})
    if not theorem or not isinstance(theorem, dict):
        return 0

    count = 0

    # Plain language — NO TRUNCATION
    plain = theorem.get("plain", "")
    if plain:
        count += build_text_slides(prs, "THE THEOREM — PLAIN LANGUAGE",
                                   data.get("theoremName", ""),
                                   plain, "Here is the theorem in plain language.")

    # Axioms — one slide per axiom for full coverage
    axioms = theorem.get("axioms", [])
    if axioms:
        # Overview slide
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        set_slide_bg(slide)
        mckinsey_header(slide, "AXIOMS", "The structural assumptions that drive the result")

        for j, ax in enumerate(axioms[:5]):
            y = Inches(1.6) + j * Inches(1.1)
            add_rect(slide, Inches(0.5), y, Inches(12.3), Inches(0.95), SURFACE)
            ax_id = ax.get("id", f"A{j+1}")
            ax_name = ax.get("name", "")
            add_text(slide, Inches(0.7), y + Inches(0.05), Inches(3), Inches(0.3),
                     f"{ax_id}: {ax_name}", size=14, color=NAVY, bold=True)
            statement = strip_html(ax.get("statement", ""))
            add_text(slide, Inches(0.7), y + Inches(0.35), Inches(11.5), Inches(0.55),
                     statement[:250], size=11, color=DARK_TEXT)

        axiom_notes = "Here are the axioms that underpin this theorem.\n\n"
        for ax in axioms:
            axiom_notes += f"{ax.get('id', '')}: {ax.get('name', '')} — {strip_html(ax.get('statement', ''))}\n\n"
        add_notes(slide, axiom_notes)
        count += 1

        # Individual axiom detail slides
        for j, ax in enumerate(axioms):
            statement = strip_html(ax.get("statement", ""))
            explanation = strip_html(ax.get("explanation", ax.get("plainEnglish", "")))
            full = f"{statement}\n\n{explanation}" if explanation else statement
            if len(full) > 100:
                ax_id = ax.get("id", f"A{j+1}")
                ax_name = ax.get("name", "")
                count += build_text_slides(prs, f"AXIOM {ax_id}: {ax_name}",
                                           "Deep dive into this structural assumption",
                                           full, f"Let us examine axiom {ax_id}, {ax_name}, in detail.")

    # Formal statement — NO TRUNCATION
    formal = theorem.get("formal", "")
    if formal:
        count += build_text_slides(prs, "FORMAL THEOREM STATEMENT", "",
                                   formal, "Here is the formal statement of the theorem.")

    # Proof sketch — NO TRUNCATION
    proof = theorem.get("proofSketch", "")
    if proof:
        count += build_text_slides(prs, "PROOF SKETCH",
                                   "Step-by-step construction of the result",
                                   proof, "Now we walk through the proof sketch.")

    return count


def build_welfare_slides(prs, data):
    """Welfare channel decomposition — overview + full detail per channel."""
    wa = data.get("welfareAnalysis")
    wa = wa if isinstance(wa, dict) else {}
    channels = wa.get("channels", [])
    if not channels:
        return 0

    count = 0

    # Overview slide
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    mckinsey_header(slide, "WELFARE CHANNEL DECOMPOSITION",
                   "How the total welfare cost breaks down by channel")

    for j, ch in enumerate(channels[:8]):
        col = j % 2
        row = j // 2
        x = Inches(0.5) + col * Inches(6.4)
        y = Inches(1.6) + row * Inches(1.3)

        add_rect(slide, x, y, Inches(6.1), Inches(1.15), SURFACE)
        name = ch.get("name", f"Channel {j+1}")
        value = strip_html(str(ch.get("value", "—")))
        add_text(slide, x + Inches(0.15), y + Inches(0.05), Inches(4), Inches(0.3),
                 name, size=13, color=NAVY, bold=True)
        add_text(slide, x + Inches(4.2), y + Inches(0.05), Inches(1.7), Inches(0.3),
                 value, size=13, color=RED, bold=True, align=PP_ALIGN.RIGHT)
        desc = strip_html(ch.get("description", ""))
        first_sent = desc.split('. ')[0] + '.' if '. ' in desc else desc
        add_text(slide, x + Inches(0.15), y + Inches(0.4), Inches(5.8), Inches(0.7),
                 first_sent[:200], size=10, color=MID_GRAY)

    channel_notes = "This slide shows the welfare channel decomposition.\n\n"
    for ch in channels:
        channel_notes += f"• {ch.get('name', '')}: {strip_html(str(ch.get('value', '')))}\n"
    add_notes(slide, channel_notes)
    count += 1

    # Full detail slides for EVERY channel — NO TRUNCATION
    for j, ch in enumerate(channels):
        desc = strip_html(ch.get("description", ""))
        if len(desc) < 50:
            continue
        name = ch.get("name", f"Channel {j+1}")
        value = strip_html(str(ch.get("value", "")))
        count += build_text_slides(prs, f"CHANNEL: {name.upper()}", value,
                                   desc, f"Channel {j+1}: {name}. {value}.")

    return count


def build_case_study_slides(prs, cases):
    """Case studies — FULL content per case, no truncation."""
    if not cases:
        return 0

    count = 0
    for i, case in enumerate(cases):
        title = case.get("title", f"Case Study {i+1}")
        content = strip_html(case.get("content", case.get("description", "")))
        if not content:
            continue

        count += build_text_slides(prs, f"CASE STUDY: {title.upper()}", "",
                                   content, f"Case study {i+1}: {title}.")
    return count


def build_policy_slides(prs, data):
    """Policy analysis — full content across multiple slides."""
    policy = data.get("policyAnalysis")
    if not policy or not isinstance(policy, dict):
        return 0

    count = 0
    sections = [
        ("CURRENT REGULATORY FRAMEWORK", policy.get("currentFramework", "")),
        ("WHY THE CURRENT FRAMEWORK FAILS", policy.get("failures", "")),
        ("REFORM PATHWAY", policy.get("reform", "")),
    ]

    for title, content in sections:
        if not content:
            continue
        count += build_text_slides(prs, title, "", content)

    return count


def build_agent_slides(prs, agents):
    """Six-agent framework — FULL detail per agent."""
    if not agents or not isinstance(agents, dict):
        return 0

    agent_meta = {
        "whistleblower": ("WHISTLEBLOWER", "Break Information Asymmetry — The insider who sees the Hollow Win from within"),
        "plaintiff": ("PLAINTIFF", "Monetize Liability — Class action, qui tam, securities fraud"),
        "regulator": ("REGULATOR", "Redesign the Game — Structural reform, not just enforcement"),
        "legislator": ("LEGISLATOR", "Change the Law — Close the statutory gaps that enable the Hollow Win"),
        "investor": ("INVESTOR", "Reprice Capital — ESG integration that actually measures welfare beta"),
        "supranational": ("SUPRANATIONAL", "Coordinate Jurisdictions — Treaty-based solutions for cross-border externalities"),
    }

    count = 0
    for key in ["whistleblower", "plaintiff", "regulator", "legislator", "investor", "supranational"]:
        content = agents.get(key)
        if not content:
            # Try "insider" for backward compat
            if key == "whistleblower":
                content = agents.get("insider")
            if not content:
                continue

        label, subtitle = agent_meta.get(key, (key.upper(), ""))
        count += build_text_slides(prs, f"AGENT: {label}", subtitle,
                                   content, f"The {label.lower()} agent.")
    return count


def build_cross_domain_slides(prs, connections):
    """Cross-domain connections — full detail."""
    if not connections:
        return 0

    count = 0
    for conn in connections:
        domain = conn.get("domain", "")
        desc = strip_html(conn.get("connection", conn.get("description", "")))
        if not desc:
            continue

        slide = prs.slides.add_slide(prs.slide_layouts[6])
        set_slide_bg(slide)
        mckinsey_header(slide, "CROSS-DOMAIN CONNECTION", domain)

        add_multiline(slide, Inches(0.6), Inches(1.6), Inches(12), Inches(5.2),
                      desc[:800], size=13, color=DARK_TEXT)

        add_notes(slide, f"Cross-domain connection to {domain}.\n\n{desc}")
        count += 1

    return count


def build_mc_slides(prs, data):
    """Monte Carlo section — full methodology."""
    wa = data.get("welfareAnalysis")
    wa = wa if isinstance(wa, dict) else {}
    mc = wa.get("monteCarlo", "")
    if not mc:
        return 0
    return build_text_slides(prs, "MONTE CARLO SIMULATION",
                             "N=10,000 draws | Seed=42 | Reproducible",
                             mc, "Here is the Monte Carlo methodology.")


def build_channel_table_slides(prs, data):
    """Channel table from Gemini summary."""
    wa = data.get("welfareAnalysis")
    wa = wa if isinstance(wa, dict) else {}
    ct = wa.get("channelTable", "")
    if not ct:
        return 0
    return build_text_slides(prs, "WELFARE CHANNEL TABLE", "", ct)


def build_faq_slides(prs, faq):
    """FAQ — one slide per Q&A pair with full answers."""
    if not faq:
        return 0

    count = 0
    for i, item in enumerate(faq):
        q = strip_html(item.get("q", ""))
        a = strip_html(item.get("a", ""))
        if not q:
            continue

        # If answer is very long, split across slides
        if len(a) > 600:
            a_chunks = split_into_slide_chunks(a, 500)
        else:
            a_chunks = [a]

        for j, chunk in enumerate(a_chunks):
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            set_slide_bg(slide)
            page_label = f" ({j+1}/{len(a_chunks)})" if len(a_chunks) > 1 else ""
            mckinsey_header(slide, f"FAQ {i+1} OF {len(faq)}{page_label}")

            add_text(slide, Inches(0.6), Inches(1.3), Inches(12), Inches(0.8),
                     q, size=18, color=NAVY, bold=True, italic=True)

            add_multiline(slide, Inches(0.6), Inches(2.3), Inches(12), Inches(4.0),
                          chunk, size=13, color=DARK_TEXT)

            add_notes(slide, f"Question: {q}\n\nAnswer: {chunk}")
            count += 1

    return count


def build_methodology_slides(prs, data):
    """Methodology — full text."""
    meth = data.get("methodology", "")
    if not meth:
        return 0
    return build_text_slides(prs, "METHODOLOGY",
                             "Data sources, calibration, and limitations", meth)


def build_literature_slides(prs, data):
    """Literature context — full text."""
    lit = data.get("literatureContext", "")
    if not lit:
        return 0
    return build_text_slides(prs, "LITERATURE CONTEXT",
                             "Where this paper sits in the academic landscape", lit)


def build_timeline_slides(prs, timeline):
    """Timeline events — multiple slides."""
    if not timeline:
        return 0

    count = 0
    # 6 events per slide
    for page in range(0, len(timeline), 6):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        set_slide_bg(slide)
        page_label = f" ({page//6 + 1})" if len(timeline) > 6 else ""
        mckinsey_header(slide, f"TIMELINE{page_label}", "Key events in the domain's history")

        batch = timeline[page:page+6]
        notes = "Timeline:\n\n"
        for j, event in enumerate(batch):
            y = Inches(1.5) + j * Inches(0.9)
            year = str(event.get("year", ""))
            evt = strip_html(event.get("event", ""))
            sig = strip_html(event.get("significance", ""))

            add_text(slide, Inches(0.6), y, Inches(1.2), Inches(0.4),
                     year, size=14, color=NAVY, bold=True)
            add_text(slide, Inches(2.0), y, Inches(4), Inches(0.4),
                     evt[:80], size=12, color=DARK_TEXT, bold=True)
            add_text(slide, Inches(6.2), y, Inches(6.5), Inches(0.8),
                     sig[:200], size=11, color=MID_GRAY)
            notes += f"{year}: {evt}. {sig}\n"

        add_notes(slide, notes)
        count += 1

    return count


def build_docx_table_slide(prs, table, slug):
    """Render a Word doc table as a McKinsey-formatted slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    mckinsey_header(slide, f"DATA TABLE: {table['title'].upper()}",
                   f"{len(table['rows'])} data rows (from published paper)")

    header = table["header"]
    rows = table["rows"][:18]  # Max rows that fit
    n_cols = min(len(header), 8)
    if n_cols == 0:
        return 0

    col_w = min(Inches(12) / n_cols, Inches(3))
    start_x = Inches(0.5)

    # Header row
    y = Inches(1.5)
    add_rect(slide, start_x, y, col_w * n_cols + Inches(0.4), Inches(0.35), NAVY)
    for j, h in enumerate(header[:n_cols]):
        add_text(slide, start_x + Inches(0.1) + j * col_w, y + Inches(0.02),
                 col_w - Inches(0.1), Inches(0.3),
                 h[:40], size=10, color=WHITE, bold=True)

    # Data rows
    for r, row in enumerate(rows):
        y = Inches(1.85) + r * Inches(0.30)
        if r % 2 == 0:
            add_rect(slide, start_x, y, col_w * n_cols + Inches(0.4), Inches(0.30), SURFACE)
        for j, cell in enumerate(row[:n_cols]):
            add_text(slide, start_x + Inches(0.1) + j * col_w, y + Inches(0.02),
                     col_w - Inches(0.1), Inches(0.26),
                     str(cell)[:40], size=9, color=DARK_TEXT)

    source_line(slide, f"Extracted from {slug} Word document")

    # Build notes from table data
    notes = f"This table has {len(table['rows'])} rows.\n\nColumns: {', '.join(header[:n_cols])}\n\n"
    for r, row in enumerate(table['rows'][:10]):
        notes += " | ".join(str(c)[:30] for c in row[:n_cols]) + "\n"
    add_notes(slide, notes)
    return 1


def build_quiz_slides(prs, data):
    """Quiz section — comprehension check."""
    count = 0

    # Quiz intro
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, NAVY)
    add_text(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(0.5),
             "COMPREHENSION CHECK", size=14, color=CYAN, bold=True)

    wa = data.get("welfareAnalysis")
    wa = wa if isinstance(wa, dict) else {}
    beta = wa.get("betaW", "?")
    tt = data.get("theoremType", "?")
    tn = data.get("theoremName", "?")

    questions = [
        (f"Q1: What is the system beta (βW) for this domain?", f"A: βW = {beta}"),
        (f"Q2: Is this an Impossibility or Intractability theorem?", f"A: {tt}"),
        (f"Q3: What is the theorem called?", f"A: {tn}"),
        ("Q4: Can this problem be solved with policy reform alone?",
         f"A: {'No — physical/biological constraints prevent it' if tt == 'Impossibility' else 'Yes — proven models exist'}"),
    ]

    for i, (q, a) in enumerate(questions):
        y = Inches(1.5) + i * Inches(1.3)
        add_text(slide, Inches(0.8), y, Inches(11), Inches(0.5),
                 q, size=16, color=WHITE, bold=True)
        add_text(slide, Inches(1.2), y + Inches(0.5), Inches(10.5), Inches(0.5),
                 a, size=14, color=CYAN)

    add_notes(slide, f"Quiz time. Test your understanding of this domain.\n\n" +
              "\n".join(f"{q}\n{a}" for q, a in questions))
    count += 1

    return count


def build_closing_slide(prs, data):
    """Closing slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)

    add_rect(slide, 0, Inches(3.2), SLIDE_W, Inches(0.04), CYAN)

    title = data.get("title", "")
    add_text(slide, Inches(1), Inches(1.0), Inches(11), Inches(1.5),
             title, size=20, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

    wa = data.get("welfareAnalysis")
    wa = wa if isinstance(wa, dict) else {}
    beta = wa.get("betaW", "")
    if beta and beta != "—":
        add_text(slide, Inches(1), Inches(3.5), Inches(11), Inches(0.8),
                 f"βW = {beta}", size=48, color=CYAN, bold=True, align=PP_ALIGN.CENTER)

    tt = data.get("theoremType", "")
    if tt:
        color = RED if tt == "Impossibility" else GOLD
        add_text(slide, Inches(1), Inches(4.5), Inches(11), Inches(0.5),
                 f"{tt} Theorem", size=20, color=color, align=PP_ALIGN.CENTER)

    add_text(slide, Inches(1), Inches(5.5), Inches(11), Inches(0.4),
             "Erik Postnieks  |  Independent Researcher  |  postnieks.com",
             size=14, color=MID_GRAY, align=PP_ALIGN.CENTER)
    add_text(slide, Inches(1), Inches(6.0), Inches(11), Inches(0.4),
             "Full paper: SSRN  |  MC replication: github.com/epostnieks",
             size=12, color=MID_GRAY, align=PP_ALIGN.CENTER)
    add_text(slide, Inches(1), Inches(6.5), Inches(11), Inches(0.3),
             "END OF MASTERCLASS", size=11, color=CYAN, bold=True, align=PP_ALIGN.CENTER)

    add_notes(slide, f"""Thank you for completing this masterclass on {title}.

The system beta for this domain is {beta}. This is a {tt} theorem.

The full paper is available on SSRN. The Monte Carlo simulation code is available for replication on GitHub at github.com/epostnieks.

This concludes the masterclass. If you found this valuable, explore the other 60 domain papers in the SAPM framework at postnieks.com.""")
    return 1


# ══════════════════════════════════════════════════════════════
# MAIN GENERATOR — MASTERCLASS ASSEMBLY
# ══════════════════════════════════════════════════════════════

def generate_masterclass(json_path):
    """Generate a 150-200 slide masterclass PPTX.

    Structure:
    1. Title slide
    2. Chapter 1: The Numbers (metrics + MC figures)
    3. Chapter 2: Executive Summary (full text)
    4. Chapter 3: Key Findings (one per slide + figure anchors)
    5. Chapter 4: The Theorem (plain, axioms, formal, proof)
    6. Chapter 5: Welfare Channels (overview + detail + figures)
    7. Chapter 6: Monte Carlo (methodology + MC distribution figure)
    8. Chapter 7: Evidence (case studies + Word doc tables + figures)
    9. Chapter 8: Policy Analysis (framework, failures, reform)
    10. Chapter 9: Six-Agent Framework (full detail per agent)
    11. Chapter 10: Cross-Domain + Literature + Timeline
    12. Chapter 11: FAQ (full Q&A)
    13. Chapter 12: Quiz + Closing
    """
    with open(json_path, 'r') as f:
        data = json.load(f)

    slug = json_path.stem

    # Find assets
    figure_pngs = find_figure_pngs(slug)
    docx_path = find_docx(slug)
    docx_tables = extract_docx_tables(docx_path) if docx_path else []
    docx_images = extract_docx_images(docx_path) if docx_path else []

    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    total = 0

    # ── Title ──
    total += build_title_slide(prs, data)

    # ── Chapter 1: The Numbers ──
    total += build_chapter_divider(prs, 1, "THE NUMBERS",
        "System beta, welfare cost, confidence intervals, and what they mean",
        "Chapter 1: The Numbers. We begin with the quantitative foundation — the system beta, the welfare cost, and the Monte Carlo confidence intervals.")
    total += build_metrics_slide(prs, data)

    # Insert MC distribution figure if available
    for fig in figure_pngs:
        if 'mc_distribution' in fig.name or 'mc_dist' in fig.name:
            total += build_figure_slide(prs, fig,
                f"Monte Carlo Distribution — {slug}",
                f"This figure shows the Monte Carlo distribution of beta W for {slug}. The histogram represents 10,000 draws from the calibrated distributions.")
            break

    # Insert sensitivity figure if available
    for fig in figure_pngs:
        if 'sensitivity' in fig.name:
            total += build_figure_slide(prs, fig,
                f"Sensitivity Analysis — {slug}",
                f"This tornado chart shows how beta W responds to changes in each input parameter.")
            break

    # ── Chapter 2: Executive Summary ──
    total += build_chapter_divider(prs, 2, "EXECUTIVE SUMMARY",
        "The complete argument — what this paper proves and why it matters",
        "Chapter 2: Executive Summary. This is the complete argument in the author's own words.")
    es = data.get("executiveSummary", "")
    total += build_text_slides(prs, "EXECUTIVE SUMMARY",
                               "What this paper proves and why it matters", es,
                               "Let me walk you through the executive summary.")

    # ── Chapter 3: Key Findings ──
    total += build_chapter_divider(prs, 3, "KEY FINDINGS",
        "What the evidence shows — one finding at a time",
        "Chapter 3: Key Findings. Each finding is presented individually with full supporting evidence.")
    findings = data.get("keyFindings", [])
    total += build_finding_slides(prs, findings)

    # Insert channel breakdown figure after findings
    for fig in figure_pngs:
        if 'channel_breakdown' in fig.name or 'channel' in fig.name:
            total += build_figure_slide(prs, fig,
                f"Welfare Channel Breakdown — {slug}",
                f"This McKinsey-style figure shows how the total welfare cost decomposes across channels.")
            break

    # ── Chapter 4: The Theorem ──
    total += build_chapter_divider(prs, 4, "THE THEOREM",
        data.get("theoremName", "Formal Result"),
        f"Chapter 4: The Theorem. We now examine the formal impossibility result — the {data.get('theoremName', 'theorem')}.")
    total += build_theorem_slides(prs, data)

    # ── Chapter 5: Welfare Channels ──
    total += build_chapter_divider(prs, 5, "WELFARE CHANNEL DECOMPOSITION",
        "Where the damage comes from — channel by channel",
        "Chapter 5: Welfare Channel Decomposition. We break down the total welfare cost into its constituent channels.")
    total += build_welfare_slides(prs, data)
    total += build_channel_table_slides(prs, data)

    # Insert welfare ledger figure if available
    for fig in figure_pngs:
        if 'welfare_ledger' in fig.name or 'ledger' in fig.name:
            total += build_figure_slide(prs, fig,
                f"Welfare Ledger — {slug}",
                f"The welfare ledger showing GDP recorded versus welfare destroyed.")
            break

    # ── Chapter 6: Monte Carlo ──
    total += build_chapter_divider(prs, 6, "MONTE CARLO SIMULATION",
        "Quantifying uncertainty — 10,000 draws, reproducible results",
        "Chapter 6: Monte Carlo Simulation. Every number in this paper has been tested with 10,000 Monte Carlo draws.")
    total += build_mc_slides(prs, data)

    # Insert ALL remaining figures from deep-research that haven't been used yet
    used_figs = set()
    for fig in figure_pngs:
        if any(kw in fig.name for kw in ['mc_distribution', 'mc_dist', 'sensitivity',
               'channel_breakdown', 'channel', 'welfare_ledger', 'ledger',
               'cross_domain', 'five_literatures', 'pigou_coase']):
            used_figs.add(fig.name)

    unused_figs = [f for f in figure_pngs if f.name not in used_figs]
    if unused_figs:
        total += build_chapter_divider(prs, "6b", "PAPER FIGURES",
            f"{len(unused_figs)} McKinsey-quality figures from the research",
            f"Here are {len(unused_figs)} additional figures from the published research.")
        for fig in unused_figs:
            label = figure_label(fig)
            total += build_figure_slide(prs, fig, f"{label} — {slug}",
                f"This figure illustrates {label.lower()} for the {slug} domain.")

    # ── Chapter 7: Evidence ──
    total += build_chapter_divider(prs, 7, "EVIDENCE",
        "Case studies and data tables from the published paper",
        "Chapter 7: Evidence. Real-world case studies and data tables extracted directly from the published paper.")

    cases = data.get("caseStudies", [])
    total += build_case_study_slides(prs, cases)

    # Word doc tables — ALL of them
    for table in docx_tables:
        total += build_docx_table_slide(prs, table, slug)

    # Embedded Word doc images
    for i, img in enumerate(docx_images[:30]):
        total += build_embedded_image_slide(prs, img["data"], img["ext"],
            f"Paper Figure {i+1}", f"This is figure {i+1} from the published paper.")

    # ── Chapter 8: Policy ──
    total += build_chapter_divider(prs, 8, "POLICY ANALYSIS",
        "What works, what fails, and why",
        "Chapter 8: Policy Analysis. We examine the current regulatory framework, identify why it fails, and propose reform pathways.")
    total += build_policy_slides(prs, data)

    # ── Chapter 9: Six-Agent Framework ──
    total += build_chapter_divider(prs, 9, "SIX-AGENT CONFLICTORING FRAMEWORK",
        "Who can break the Private-Systemic Tension",
        "Chapter 9: The Six-Agent Conflictoring Framework. For each of six agent types, we describe the specific actions available to break the Hollow Win.")
    agents = data.get("agents", {})
    total += build_agent_slides(prs, agents)

    # ── Chapter 10: Context ──
    total += build_chapter_divider(prs, 10, "CONNECTIONS AND CONTEXT",
        "Cross-domain links, literature, and historical timeline",
        "Chapter 10: Connections and Context. How this domain connects to others in the SAPM framework.")

    connections = data.get("crossDomainConnections", [])
    total += build_cross_domain_slides(prs, connections)

    # Cross-domain figure
    for fig in figure_pngs:
        if 'cross_domain' in fig.name:
            total += build_figure_slide(prs, fig,
                f"Cross-Domain Connections — {slug}",
                f"Visual map of how {slug} connects to other domains in the SAPM framework.")
            break

    # Five literatures figure
    for fig in figure_pngs:
        if 'five_literatures' in fig.name:
            total += build_figure_slide(prs, fig,
                f"Five Literatures — {slug}",
                f"Where this paper sits across five bodies of academic literature.")
            break

    total += build_literature_slides(prs, data)
    total += build_methodology_slides(prs, data)

    timeline = data.get("timeline", [])
    total += build_timeline_slides(prs, timeline)

    # ── Chapter 11: FAQ ──
    total += build_chapter_divider(prs, 11, "FREQUENTLY ASKED QUESTIONS",
        "Anticipating and answering the hard questions",
        "Chapter 11: Frequently Asked Questions. These are the questions that skeptics, reviewers, and practitioners most commonly ask.")
    faq = data.get("faq", [])
    total += build_faq_slides(prs, faq)

    # ── Chapter 12: Quiz + Closing ──
    total += build_chapter_divider(prs, 12, "ASSESSMENT",
        "Test your understanding",
        "Chapter 12: Assessment. Let us test your understanding of the key concepts.")
    total += build_quiz_slides(prs, data)
    total += build_closing_slide(prs, data)

    # Save
    out_path = OUTPUT_DIR / f"{slug}.pptx"
    prs.save(str(out_path))

    return slug, total, len(figure_pngs), len(docx_tables), len(docx_images)


# ══════════════════════════════════════════════════════════════
# MAIN — Parallel execution
# ══════════════════════════════════════════════════════════════

def main():
    json_files = sorted(PAPER_DATA_DIR.glob("*.json"))
    json_files = [f for f in json_files if '_raw' not in f.name]
    print(f"Found {len(json_files)} paper JSON files")

    results = []
    errors = []

    with ThreadPoolExecutor(max_workers=8) as pool:
        futures = {pool.submit(generate_masterclass, f): f for f in json_files}
        for future in as_completed(futures):
            f = futures[future]
            try:
                slug, slides, figs, tables, images = future.result()
                results.append((slug, slides, figs, tables, images))
                print(f"  OK  {slug}: {slides} slides, {figs} figures, {tables} tables, {images} images")
            except Exception as e:
                errors.append((f.stem, str(e)))
                print(f"  FAIL {f.stem}: {e}")
                traceback.print_exc()

    results.sort(key=lambda x: x[0])
    total_slides = sum(r[1] for r in results)
    total_figs = sum(r[2] for r in results)
    total_tables = sum(r[3] for r in results)
    total_images = sum(r[4] for r in results)

    print(f"\nDone: {len(results)} OK, {len(errors)} failed")
    print(f"Total: {total_slides} slides, {total_figs} figures, {total_tables} tables, {total_images} images")
    print(f"Average: {total_slides / max(len(results), 1):.0f} slides per paper")
    print(f"Output: {OUTPUT_DIR}")

    if errors:
        print(f"\nFailed:")
        for slug, err in errors:
            print(f"  {slug}: {err}")


if __name__ == "__main__":
    main()
