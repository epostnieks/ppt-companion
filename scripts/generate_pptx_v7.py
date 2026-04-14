#!/usr/bin/env python3
"""
PPTX Generator V7 — STONECASTLE DESIGN SYSTEM
Adapted from Steadfastly_StoneCastle_Deck_v30.pptx (proven Claude output, 10× better than V5).
Design: gold accent line, navy number squares, alternating-row tables, dense grid, action titles.
"""

import json, os, re, sys, glob, traceback, copy
from concurrent.futures import ThreadPoolExecutor, as_completed
from pathlib import Path
from html.parser import HTMLParser
from io import BytesIO
from lxml import etree

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, MSO_AUTO_SIZE
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn, nsmap

try:
    from docx import Document as DocxDocument
    HAS_DOCX = True
except ImportError:
    HAS_DOCX = False

# ── Paths ──
SCRIPT_DIR = Path(__file__).parent
PROJECT_DIR = SCRIPT_DIR.parent
PAPER_DATA_DIR = PROJECT_DIR / "src" / "paperData"
PPTX_DIR = PROJECT_DIR / "public" / "pptx"
AUDIO_DIR = PROJECT_DIR / "public" / "audio"
PPTX_DIR.mkdir(parents=True, exist_ok=True)
AUDIO_DIR.mkdir(parents=True, exist_ok=True)
BASE = Path("/Users/erikpostnieks/Projects")
DR_OUTPUT = BASE / "deep-research" / "output"

# ── Dimensions (Stonecastle: 10" × 5.625") ──
SLIDE_W = Inches(10)
SLIDE_H = Inches(5.625)
MARGIN = Inches(0.5)
CONTENT_W = Inches(9.0)

# ── McKinsey Figure Spec v2.2 Color Palette ──
NAVY = RGBColor(0x1F, 0x4E, 0x79)        # Title bars, headers (#1F4E79)
DARK_NAVY = RGBColor(0x0F, 0x1A, 0x30)   # Dark title slide bg
BLUE_GRAY = RGBColor(0x3D, 0x6B, 0x8E)   # Secondary headers, source text (#3D6B8E)
GOLD = RGBColor(0xC4, 0x97, 0x2A)        # Accent, punchline border (#C4972A)
CRIMSON = RGBColor(0x8B, 0x1A, 0x1A)     # Warning, impossibility (#8B1A1A)
GREEN = RGBColor(0x2E, 0x7D, 0x32)       # Positive indicators (#2E7D32)
CREAM = RGBColor(0xFB, 0xF9, 0xF4)       # Subtitle bar, even rows (#FBF9F4)
TAN = RGBColor(0xF5, 0xEF, 0xE0)         # Punchline box, highlights (#F5EFE0)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x2D, 0x2D, 0x2D)       # Primary text (#2D2D2D)
MID_GRAY = RGBColor(0x76, 0x76, 0x76)
BORDER = RGBColor(0xE0, 0xD8, 0xC8)      # Subtle borders (#E0D8C8)
ROW_BORDER = RGBColor(0xE8, 0xE0, 0xD0)  # Table row separators (#E8E0D0)

FONT = "Arial"
FONT_TITLE = "Georgia"

# ── Grid constants (Stonecastle uses 0.50" margins, 0.27" row height) ──
ROW_H = Inches(0.27)
GOLD_LINE_H = Inches(0.06)
FOOTER_Y = Inches(5.05)
SQ_SIZE = Inches(0.30)  # Navy number squares


# ══════════════════════════════════════════════════════════════
# TEXT UTILITIES
# ══════════════════════════════════════════════════════════════

class HTMLStripper(HTMLParser):
    def __init__(self):
        super().__init__()
        self.result = []
    def handle_data(self, d):
        self.result.append(d)
    def handle_starttag(self, tag, attrs):
        if tag in ('td', 'th'):
            self.result.append(' | ')
        if tag == 'tr':
            self.result.append('\n')
        if tag == 'li':
            self.result.append('\n• ')
        if tag in ('br', 'p'):
            self.result.append('\n')
    def get_text(self):
        return re.sub(r'\n{3,}', '\n\n', ''.join(self.result)).strip()

def strip_html(text):
    if not text or not isinstance(text, str):
        return str(text) if text else ""
    s = HTMLStripper()
    s.feed(text)
    return s.get_text()

def first_sentence(text):
    """Extract the first sentence as an action title."""
    text = strip_html(text)
    m = re.match(r'^(.+?[.!?])\s', text)
    return m.group(1) if m else text[:120]

def split_sentences(text):
    text = strip_html(text)
    parts = re.split(r'(?<=[.!?])\s+', text)
    return [p.strip() for p in parts if p.strip()]

def paginate_text(text, max_chars=1400):
    """Split text into page-sized chunks. 1400 chars = dense but readable at 16pt."""
    text = strip_html(text)
    if not text:
        return []
    sentences = split_sentences(text)
    pages = []
    current = ""
    for s in sentences:
        if len(current) + len(s) + 1 > max_chars and current:
            pages.append(current.strip())
            current = s
        else:
            current = (current + " " + s).strip()
    if current:
        pages.append(current.strip())
    return pages

def action_title(section_label, content=""):
    """Generate an action title — a complete sentence stating the insight.
    If content provides a good first sentence, use it. Otherwise use the label."""
    if content:
        fs = first_sentence(content)
        if len(fs) > 20 and len(fs) < 140 and '.' in fs:
            return fs
    return section_label

def clean_theorem_display(theorem_type, theorem_name):
    """Remove redundant theorem type from name."""
    if not theorem_type or not theorem_name:
        return theorem_name or theorem_type or ""
    # "The X Impossibility Theorem" already contains "Impossibility"
    if theorem_type.lower() in theorem_name.lower():
        return theorem_name
    return f"{theorem_name} ({theorem_type})"


# ══════════════════════════════════════════════════════════════
# SHAPE HELPERS — Minimal, non-overlapping
# ══════════════════════════════════════════════════════════════

def set_bg(slide, color=WHITE):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_rect(slide, left, top, width, height, color):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    shape.shadow.inherit = False
    return shape

def add_line(slide, left, top, width, color=MID_GRAY, thickness=Pt(0.5)):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, thickness)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    shape.shadow.inherit = False
    return shape

def add_text(slide, left, top, width, height, text, size=12,
             color=BLACK, bold=False, italic=False, align=PP_ALIGN.LEFT,
             font_name=None, line_spacing=1.0):
    """Text box — word-wraps, no overflow."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.NONE

    p = tf.paragraphs[0]
    p.text = str(text)[:3000]
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.italic = italic
    p.font.name = font_name or FONT
    p.alignment = align

    if line_spacing != 1.0:
        pPr = p._p.get_or_add_pPr()
        lnSpc = pPr.makeelement(qn('a:lnSpc'), {})
        spcPct = lnSpc.makeelement(qn('a:spcPct'), {'val': str(int(line_spacing * 100000))})
        lnSpc.append(spcPct)
        pPr.append(lnSpc)

    return txBox

def add_body_text(slide, left, top, width, height, text, size=11,
                  color=BLACK, line_spacing=1.15):
    """Multi-paragraph body text."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.NONE

    paragraphs = str(text).split('\n')
    first = True
    for para_text in paragraphs:
        para_text = para_text.strip()
        if not para_text:
            continue
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.text = para_text
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.font.name = FONT
        p.space_after = Pt(3)

        pPr = p._p.get_or_add_pPr()
        lnSpc = pPr.makeelement(qn('a:lnSpc'), {})
        spcPct = lnSpc.makeelement(qn('a:spcPct'), {'val': str(int(line_spacing * 100000))})
        lnSpc.append(spcPct)
        pPr.append(lnSpc)

    return txBox

def add_notes(slide, text):
    if not text:
        return
    ns = slide.notes_slide
    ns.notes_text_frame.text = str(text)[:5000]

def add_num_square(slide, left, top, num, color=NAVY):
    """Navy numbered square (Stonecastle signature element)."""
    add_rect(slide, left, top, SQ_SIZE, SQ_SIZE, color)
    add_text(slide, left, top + Inches(0.02), SQ_SIZE, SQ_SIZE - Inches(0.04),
             str(num), size=11, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

def add_gold_line(slide):
    """Gold accent line across top of slide (Stonecastle signature)."""
    add_rect(slide, 0, 0, SLIDE_W, GOLD_LINE_H, GOLD)

def add_footer(slide, page_num, label="SAPM Masterclass"):
    """Footer: label + page number (Stonecastle style)."""
    add_text(slide, Inches(3.5), FOOTER_Y, Inches(3.0), Inches(0.30),
             label, size=8, color=MID_GRAY, align=PP_ALIGN.CENTER)
    add_text(slide, Inches(9.2), FOOTER_Y, Inches(0.50), Inches(0.30),
             str(page_num), size=8, color=MID_GRAY, align=PP_ALIGN.RIGHT)

def add_source(slide, text):
    """Source citation above footer."""
    add_text(slide, MARGIN, Inches(4.85), CONTENT_W, Inches(0.18),
             text, size=7, color=MID_GRAY, italic=True)


# ══════════════════════════════════════════════════════════════
# SLIDE TEMPLATES — McKinsey layouts
# ══════════════════════════════════════════════════════════════

def slide_title_dark(prs, title, subtitle="", badge="", epigraph="", notes_text="", page_num=1):
    """Dark title slide — paper opener (Stonecastle style)."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, DARK_NAVY)
    add_gold_line(slide)

    # Series label
    add_text(slide, Inches(0.50), Inches(0.30), Inches(8), Inches(0.25),
             "SYSTEM ASSET PRICING MODEL — MASTERCLASS", size=10, color=GOLD, bold=True)

    # Main title
    add_text(slide, Inches(0.50), Inches(0.80), Inches(8.5), Inches(1.5),
             title, size=24, color=WHITE, bold=True, font_name=FONT_TITLE,
             line_spacing=1.1)

    # Theorem badge
    if badge:
        badge_color = CRIMSON if 'Impossibility' in badge else GOLD
        add_text(slide, Inches(0.50), Inches(2.5), Inches(8), Inches(0.30),
                 badge, size=13, color=badge_color, bold=True)

    # Epigraph
    if epigraph:
        add_text(slide, Inches(0.50), Inches(3.2), Inches(7.5), Inches(0.80),
                 f'"{epigraph}"', size=11, color=MID_GRAY, italic=True,
                 font_name=FONT_TITLE, line_spacing=1.2)

    # Author line + divider
    add_line(slide, Inches(0.50), Inches(4.35), Inches(3.0), color=MID_GRAY)
    add_text(slide, Inches(0.50), Inches(4.45), Inches(8), Inches(0.25),
             "Erik Postnieks  |  Independent Researcher  |  2026", size=9, color=MID_GRAY)

    add_footer(slide, page_num, "Confidential")
    add_notes(slide, notes_text or f"Welcome to the SAPM analysis of {title}. {badge}")
    return slide

def slide_section(prs, chapter_num, title, subtitle="", notes_text="", page_num=0):
    """Chapter divider — navy background (Stonecastle style)."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, NAVY)
    add_rect(slide, 0, 0, SLIDE_W, GOLD_LINE_H, GOLD)

    add_num_square(slide, Inches(0.50), Inches(1.5), chapter_num, GOLD)
    add_text(slide, Inches(0.95), Inches(1.47), Inches(4), Inches(0.35),
             f"CHAPTER {chapter_num}", size=11, color=GOLD, bold=True)

    add_text(slide, Inches(0.50), Inches(2.1), Inches(8.5), Inches(1.2),
             title, size=24, color=WHITE, bold=True, font_name=FONT_TITLE,
             line_spacing=1.1)

    if subtitle:
        add_text(slide, Inches(0.50), Inches(3.3), Inches(8), Inches(0.50),
                 subtitle, size=11, color=RGBColor(0xBB, 0xCC, 0xDD), italic=True)

    add_footer(slide, page_num, "SAPM Masterclass")
    add_notes(slide, notes_text or f"Chapter {chapter_num}: {title}. {subtitle}")
    return slide

def slide_content(prs, title, body, subtitle=None, notes_text=None, source=None, page_num=0):
    """Standard content slide — Stonecastle design: gold line, action title, dense body."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_gold_line(slide)

    # Action title (full sentence)
    add_text(slide, MARGIN, Inches(0.15), Inches(9.0), Inches(0.55),
             title[:150], size=16, color=NAVY, bold=True, line_spacing=1.05)

    # Subtitle
    body_top = Inches(0.80)
    if subtitle:
        add_text(slide, MARGIN, Inches(0.65), Inches(9.0), Inches(0.20),
                 subtitle, size=9, color=MID_GRAY, italic=True)
        body_top = Inches(0.90)

    # Body text — constrained to safe zone above footer
    body_h = Inches(3.8) if not source else Inches(3.6)
    add_body_text(slide, MARGIN, body_top, CONTENT_W, body_h,
                  body[:2400], size=11, color=BLACK, line_spacing=1.15)

    if source:
        add_source(slide, source)
    add_footer(slide, page_num, "SAPM Masterclass")
    add_notes(slide, notes_text or body[:3000])
    return slide

def slide_figure(prs, fig_path, caption="", notes_text="", page_num=0):
    """Full-page figure — figure IS the slide. McKinsey figs are 1400×1050 (4:3).
    Slide is 10"×5.625" (16:9). Center the figure, filling width."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    try:
        # 1400×1050 = 4:3 aspect. At 10" wide → 7.5" tall (taller than slide).
        # Scale to fit height: 5.625"/7.5" = 0.75 → width = 10"×0.75 = 7.5"
        fig_w = Inches(7.5)
        fig_h = SLIDE_H
        fig_left = Inches(1.25)  # Center horizontally
        slide.shapes.add_picture(str(fig_path),
            left=fig_left, top=0,
            width=fig_w, height=fig_h)
    except Exception:
        add_text(slide, Inches(2), Inches(2.5), Inches(6), Inches(0.5),
                 f"[Figure: {fig_path.name}]", size=14, color=MID_GRAY,
                 align=PP_ALIGN.CENTER)
    add_notes(slide, notes_text or f"Figure: {fig_label(fig_path)}. {caption}")
    return slide

def slide_metrics(prs, data, page_num=0):
    """Key metrics dashboard — 2×2 grid, McKinsey layout, 5.625" height."""
    wa = data.get("welfareAnalysis")
    wa = wa if isinstance(wa, dict) else {}
    beta = wa.get("betaW", data.get("betaW", "—"))
    ci = wa.get("ci90", data.get("ci90", "—"))
    pi_raw = wa.get("pi", "—")
    dw_raw = wa.get("deltaW", "—")

    def money(val):
        if not val or val == "—": return "—"
        m = re.search(r'[-−]?\$[\d,.]+\s*(billion|trillion|million)?', str(val), re.I)
        return m.group(0) if m else str(val)[:60]

    tt = data.get("theoremType", "")
    title = f"Every dollar of revenue destroys ${beta} in social welfare"
    if beta == "—":
        title = "Welfare metrics for this domain"

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_gold_line(slide)
    add_text(slide, MARGIN, Inches(0.15), Inches(9.0), Inches(0.50),
             title, size=16, color=NAVY, bold=True, font_name=FONT_TITLE)

    metrics = [
        ("System Beta (βW)", str(beta), "Welfare loss per $ of private gain", NAVY),
        ("90% Confidence Interval", str(ci), "Monte Carlo bounds (N=10,000)", BLUE_GRAY),
        ("Industry Revenue (Π)", money(pi_raw), "Annual revenue — never profit (Iron Law)", GREEN),
        ("Welfare Cost (ΔW)", money(dw_raw), "Annual welfare destruction", CRIMSON),
    ]

    for i, (label, value, desc, accent) in enumerate(metrics):
        col, row = i % 2, i // 2
        x = Inches(0.5) + col * Inches(4.7)
        y = Inches(0.80) + row * Inches(1.5)
        w, h = Inches(4.4), Inches(1.3)

        add_rect(slide, x, y, w, h, CREAM)
        add_rect(slide, x, y, Inches(0.06), h, accent)

        add_text(slide, x + Inches(0.20), y + Inches(0.08), w - Inches(0.3), Inches(0.18),
                 label, size=9, color=MID_GRAY, bold=True)
        vs = 22 if len(value) < 12 else 16 if len(value) < 20 else 12
        add_text(slide, x + Inches(0.20), y + Inches(0.32), w - Inches(0.3), Inches(0.45),
                 value, size=vs, color=accent, bold=True)
        add_text(slide, x + Inches(0.20), y + Inches(0.90), w - Inches(0.3), Inches(0.30),
                 desc, size=8, color=MID_GRAY)

    # Classification box
    tn = data.get("theoremName", "")
    if tt:
        y = Inches(3.90)
        box_color = TAN if tt == "Impossibility" else CREAM
        border_color = CRIMSON if tt == "Impossibility" else GOLD
        add_rect(slide, MARGIN, y, CONTENT_W, Inches(0.65), box_color)
        add_rect(slide, MARGIN, y, Inches(0.06), Inches(0.65), border_color)
        class_text = (f"{tn} — Physical constraints prevent welfare internalization."
                      if tt == "Impossibility"
                      else f"{tn} — Institutional constraints. Proven reforms exist.")
        add_text(slide, Inches(0.7), y + Inches(0.10), Inches(8.5), Inches(0.45),
                 class_text, size=10, color=BLACK)

    add_source(slide, "Source: Monte Carlo simulation (N=10,000, seed=42). Π = annual revenue (Iron Law).")
    add_footer(slide, page_num, "SAPM Masterclass")
    add_notes(slide, f"System Beta: {beta}. CI: {ci}. Π: {money(pi_raw)}. ΔW: {money(dw_raw)}. {tt}.")
    return slide

def slide_two_column(prs, title, left_title, left_body, right_title, right_body, notes_text=None, page_num=0):
    """Two-column comparison (Stonecastle dual-panel style)."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_gold_line(slide)
    add_text(slide, MARGIN, Inches(0.15), Inches(9.0), Inches(0.50),
             title, size=16, color=NAVY, bold=True, font_name=FONT_TITLE)

    # Left panel
    add_rect(slide, Inches(0.50), Inches(0.75), Inches(4.20), Inches(0.28), NAVY)
    add_text(slide, Inches(0.60), Inches(0.76), Inches(4.00), Inches(0.26),
             left_title, size=10, color=WHITE, bold=True)
    add_body_text(slide, Inches(0.50), Inches(1.10), Inches(4.20), Inches(3.50),
                  left_body[:1200], size=10, color=BLACK, line_spacing=1.15)

    # Right panel
    add_rect(slide, Inches(5.30), Inches(0.75), Inches(4.20), Inches(0.28), NAVY)
    add_text(slide, Inches(5.40), Inches(0.76), Inches(4.00), Inches(0.26),
             right_title, size=10, color=WHITE, bold=True)
    add_body_text(slide, Inches(5.30), Inches(1.10), Inches(4.20), Inches(3.50),
                  right_body[:1200], size=10, color=BLACK, line_spacing=1.15)

    add_footer(slide, page_num, "SAPM Masterclass")
    add_notes(slide, notes_text or f"{left_title}: {left_body}\n\n{right_title}: {right_body}")
    return slide

def slide_key_insight(prs, insight_text, supporting_text="", notes_text=None, page_num=0):
    """Punchline slide — tan background, gold border (McKinsey Spec punchline box)."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, TAN)
    add_gold_line(slide)

    add_text(slide, MARGIN, Inches(0.25), Inches(2.0), Inches(0.20),
             "KEY INSIGHT", size=9, color=GOLD, bold=True)

    add_rect(slide, Inches(0.50), Inches(0.55), Inches(0.06), Inches(1.8), GOLD)
    add_text(slide, Inches(0.75), Inches(0.55), Inches(8.5), Inches(1.8),
             insight_text, size=18, color=DARK_NAVY, bold=True, font_name=FONT_TITLE,
             line_spacing=1.25)

    if supporting_text:
        add_line(slide, Inches(0.75), Inches(2.60), Inches(2.5), MID_GRAY)
        add_body_text(slide, Inches(0.75), Inches(2.80), Inches(8.2), Inches(1.8),
                      supporting_text[:1000], size=11, color=BLACK, line_spacing=1.2)

    add_footer(slide, page_num, "SAPM Masterclass")
    add_notes(slide, notes_text or f"Key insight: {insight_text}. {supporting_text}")
    return slide

def slide_table(prs, title, header, rows, source_text="", notes_text=None, page_num=0):
    """Data table — native python-pptx table, McKinsey colors."""
    n_cols = min(len(header), 7)
    n_rows = min(len(rows) + 1, 16)  # header + data (fits 5.625" slide)
    if n_cols == 0 or n_rows <= 1:
        return None

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_gold_line(slide)
    add_text(slide, MARGIN, Inches(0.15), Inches(9.0), Inches(0.50),
             title, size=16, color=NAVY, bold=True, font_name=FONT_TITLE)

    total_w = Inches(9.0)
    col_w = total_w // n_cols

    tbl_shape = slide.shapes.add_table(
        n_rows, n_cols,
        left=MARGIN, top=Inches(0.70),
        width=total_w, height=Inches(4.0)
    )
    tbl = tbl_shape.table

    # Header row — navy bg, white text
    for j in range(n_cols):
        cell = tbl.cell(0, j)
        cell.text = str(header[j])[:40] if j < len(header) else ""
        p = cell.text_frame.paragraphs[0]
        p.font.size = Pt(9)
        p.font.color.rgb = WHITE
        p.font.bold = True
        p.font.name = FONT
        tcPr = cell._tc.get_or_add_tcPr()
        solidFill = tcPr.makeelement(qn('a:solidFill'), {})
        srgbClr = solidFill.makeelement(qn('a:srgbClr'), {'val': '1F4E79'})
        solidFill.append(srgbClr)
        tcPr.append(solidFill)

    # Data rows — alternating white/cream
    for i, row in enumerate(rows[:n_rows - 1]):
        for j in range(n_cols):
            cell = tbl.cell(i + 1, j)
            cell.text = str(row[j])[:40] if j < len(row) else ""
            p = cell.text_frame.paragraphs[0]
            p.font.size = Pt(8)
            p.font.color.rgb = BLACK
            p.font.name = FONT
            if i % 2 == 0:
                tcPr = cell._tc.get_or_add_tcPr()
                solidFill = tcPr.makeelement(qn('a:solidFill'), {})
                srgbClr = solidFill.makeelement(qn('a:srgbClr'), {'val': 'FBF9F4'})
                solidFill.append(srgbClr)
                tcPr.append(solidFill)

    if source_text:
        add_source(slide, f"Source: {source_text}")
    add_footer(slide, page_num, "SAPM Masterclass")
    add_notes(slide, notes_text or f"Table: {title}. {len(rows)} rows, {n_cols} columns.")
    return slide


# ══════════════════════════════════════════════════════════════
# ASSET EXTRACTION (from V5, proven correct)
# ══════════════════════════════════════════════════════════════

def find_figures(slug):
    sapm_fig_dir = BASE / f"sapm-{slug}" / "figures" / "png"
    if sapm_fig_dir.exists():
        pngs = sorted(sapm_fig_dir.glob("*.png"))
        if pngs:
            seen = set()
            out = []
            for p in pngs:
                base = re.sub(r'_v\d+\.png$', '.png', p.name)
                if base not in seen:
                    out.append(p)
                    seen.add(base)
            return out
    sapm_fig_dir2 = BASE / f"sapm-{slug}" / "figures"
    if sapm_fig_dir2.exists():
        pngs = sorted(sapm_fig_dir2.glob("*.png"))
        if pngs:
            return pngs
    dr_slug = slug.replace("-", "_")
    fig_dir = DR_OUTPUT / dr_slug / "figures"
    if not fig_dir.exists():
        fig_dir = DR_OUTPUT / slug / "figures"
    if not fig_dir.exists():
        return []
    pngs = sorted(fig_dir.glob("*.png"))
    seen = set()
    out = []
    for p in pngs:
        base = re.sub(r'_v\d+\.png$', '.png', p.name)
        if base not in seen:
            out.append(p)
            seen.add(base)
    return out

def find_docx(slug):
    patterns = [
        f"{BASE}/sapm-{slug}/paper/{slug}-current.docx",
        f"{BASE}/sapm-{slug}/paper/*-current.docx",
        f"{BASE}/sapm-{slug}/paper/*-final.docx",
        f"{BASE}/sapm-{slug}/paper/*.docx",
    ]
    for pat in patterns:
        matches = [m for m in glob.glob(pat) if '~' not in m and 'template' not in m.lower()]
        if matches:
            return matches[0]
    return None

def get_docx_tables(path, max_tables=50):
    if not HAS_DOCX or not path:
        return []
    try:
        doc = DocxDocument(path)
    except Exception:
        return []
    tables = []
    for i, t in enumerate(doc.tables[:max_tables]):
        rows = [[c.text.strip() for c in r.cells] for r in t.rows]
        if len(rows) > 1:
            tables.append({"title": f"Table {i+1}", "header": rows[0], "rows": rows[1:]})
    return tables

def get_docx_images(path):
    if not HAS_DOCX or not path:
        return []
    try:
        doc = DocxDocument(path)
    except Exception:
        return []
    images = []
    for rel in doc.part.rels.values():
        if "image" in rel.reltype:
            try:
                blob = rel.target_part.blob
                ext = rel.target_part.content_type.split('/')[-1]
                if ext in ('png', 'jpeg', 'jpg', 'gif'):
                    images.append({"data": blob, "ext": ext})
            except Exception:
                pass
    return images

def fig_label(path):
    name = re.sub(r'^fig\d+_', '', path.stem)
    return name.replace('_', ' ').title()

def categorize_figures(figures):
    fig_map = {}
    for fig in figures:
        name = fig.name.lower()
        if 'channel' in name or 'breakdown' in name:
            fig_map.setdefault('channel', []).append(fig)
        elif 'mc_dist' in name or 'distribution' in name or 'histogram' in name:
            fig_map.setdefault('mc', []).append(fig)
        elif 'sensitivity' in name or 'tornado' in name:
            fig_map.setdefault('sensitivity', []).append(fig)
        elif 'cross_domain' in name or 'radar' in name:
            fig_map.setdefault('crossdomain', []).append(fig)
        elif 'welfare' in name or 'ledger' in name:
            fig_map.setdefault('welfare', []).append(fig)
        elif 'literature' in name or 'five_lit' in name:
            fig_map.setdefault('literature', []).append(fig)
        elif 'pigou' in name or 'coase' in name:
            fig_map.setdefault('theory', []).append(fig)
        elif 'timeline' in name:
            fig_map.setdefault('timeline', []).append(fig)
        elif 'outcome' in name or 'payoff' in name or 'game' in name:
            fig_map.setdefault('game', []).append(fig)
        elif 'policy' in name or 'reform' in name:
            fig_map.setdefault('policy', []).append(fig)
        elif 'agent' in name:
            fig_map.setdefault('agent', []).append(fig)
        else:
            fig_map.setdefault('other', []).append(fig)
    return fig_map


# ══════════════════════════════════════════════════════════════
# CONTENT SLIDE BUILDERS
# ══════════════════════════════════════════════════════════════

def build_text_slides(prs, section_title, text, notes_prefix=""):
    """Build dense text slides with action titles from content."""
    if not text:
        return 0
    full = strip_html(text)
    if not full:
        return 0

    pages = paginate_text(full, 1400)
    count = 0
    for i, page in enumerate(pages):
        # Action title: first sentence of each page
        title = action_title(section_title, page) if i == 0 else section_title
        pg = f" ({i+1}/{len(pages)})" if len(pages) > 1 else ""

        slide_content(prs, f"{title}{pg}", page,
                      notes_text=f"{notes_prefix}\n\n{page}" if notes_prefix and i == 0 else page)
        count += 1
    return count

def build_findings_slides(prs, findings):
    """Key findings — grouped 2-3 per slide with numbered entries."""
    if not findings:
        return 0

    # Group findings
    groups = []
    grp = []
    grp_len = 0
    for f in findings:
        txt = strip_html(f) if isinstance(f, str) else strip_html(str(f))
        if grp_len + len(txt) > 1200 and grp:
            groups.append(grp)
            grp = [txt]
            grp_len = len(txt)
        else:
            grp.append(txt)
            grp_len += len(txt)
    if grp:
        groups.append(grp)

    count = 0
    finding_num = 0
    for gi, group in enumerate(groups):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        set_bg(slide)

        # Action title from first finding in group
        title = action_title("Key findings from the analysis", group[0])
        pg = f" ({gi+1}/{len(groups)})" if len(groups) > 1 else ""
        add_gold_line(slide)
        add_text(slide, MARGIN, Inches(0.15), Inches(9.0), Inches(0.50),
                 f"{title}{pg}", size=14, color=NAVY, bold=True, font_name=FONT_TITLE)

        y = Inches(1.1)
        slide_notes = []
        for txt in group:
            finding_num += 1
            sentences = split_sentences(txt)
            headline = sentences[0] if sentences else txt[:120]
            body = ' '.join(sentences[1:]) if len(sentences) > 1 else ""

            # Number circle
            add_rect(slide, MARGIN, y, Inches(0.4), Inches(0.4), NAVY)
            add_text(slide, MARGIN, y + Inches(0.05), Inches(0.4), Inches(0.3),
                     str(finding_num), size=14, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

            # Headline + body
            add_text(slide, Inches(1.1), y, Inches(8.2), Inches(0.4),
                     headline[:180], size=14, color=NAVY, bold=True)

            if body:
                body_text = body[:500]
                est_lines = len(body_text) // 80 + 1
                body_h = min(Inches(0.2 * est_lines + 0.2), Inches(1.5))
                add_body_text(slide, Inches(1.1), y + Inches(0.45), Inches(8.2), body_h,
                              body_text, size=12, color=BLACK, line_spacing=1.15)
                y += Inches(0.5) + body_h + Inches(0.15)
            else:
                y += Inches(0.7)

            slide_notes.append(f"Finding {finding_num}: {txt}")
            if y > Inches(4.5):
                break

        add_notes(slide, "\n\n".join(slide_notes))
        count += 1

    return count

def build_theorem_slides(prs, data):
    """Theorem — plain language, axioms, formal statement, proof sketch."""
    theorem = data.get("theorem", {})
    if not theorem or not isinstance(theorem, dict):
        return 0
    count = 0

    # Plain language
    plain = theorem.get("plain", "")
    if plain:
        count += build_text_slides(prs, "The theorem in plain language", plain)

    # Axioms
    axioms = theorem.get("axioms", [])
    if axioms:
        # Group axioms: 3-4 per slide
        for i in range(0, len(axioms), 4):
            batch = axioms[i:i+4]
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            set_bg(slide)

            pg = f" ({i//4+1}/{(len(axioms)+3)//4})" if len(axioms) > 4 else ""
            add_gold_line(slide)
            add_text(slide, MARGIN, Inches(0.15), Inches(9.0), Inches(0.50),
                     f"Structural axioms driving the result{pg}",
                     size=14, color=NAVY, bold=True, font_name=FONT_TITLE)

            y = Inches(0.75)
            ax_notes = []
            for ax in batch:
                ax_id = ax.get("id", "A?")
                ax_name = ax.get("name", "")
                stmt = strip_html(ax.get("statement", ""))

                est_h = max(Inches(1.0), Inches(0.3 * (len(stmt) // 90 + 2)))
                est_h = min(est_h, Inches(1.5))

                add_rect(slide, MARGIN, y, CONTENT_W, est_h, CREAM)
                add_rect(slide, MARGIN, y, Inches(0.06), est_h, GOLD)

                add_text(slide, Inches(0.7), y + Inches(0.08), Inches(8.3), Inches(0.25),
                         f"{ax_id}: {ax_name}", size=13, color=NAVY, bold=True)
                add_body_text(slide, Inches(0.7), y + Inches(0.35), Inches(8.3), est_h - Inches(0.4),
                              stmt[:400], size=11, color=BLACK, line_spacing=1.15)

                y += est_h + Inches(0.12)
                ax_notes.append(f"{ax_id}: {ax_name}. {stmt}")
                if y > Inches(4.5):
                    break

            add_notes(slide, "Axioms:\n\n" + "\n\n".join(ax_notes))
            count += 1

    # Formal + proof
    formal = theorem.get("formal", "")
    if formal:
        count += build_text_slides(prs, "Formal theorem statement", formal)

    proof = theorem.get("proofSketch", "")
    if proof:
        count += build_text_slides(prs, "Proof sketch — step-by-step construction", proof)

    return count

def build_welfare_slides(prs, data):
    """Welfare channels — overview table + detail slides."""
    wa = data.get("welfareAnalysis")
    wa = wa if isinstance(wa, dict) else {}
    channels = wa.get("channels", [])
    if not channels:
        return 0

    count = 0

    # Overview slide: all channels as a table-like layout
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)

    n_ch = len(channels)
    title = f"{n_ch} channels of welfare destruction sum to the total social cost"
    add_gold_line(slide)
    add_text(slide, MARGIN, Inches(0.15), Inches(9.0), Inches(0.50),
             title, size=14, color=NAVY, bold=True, font_name=FONT_TITLE)

    # Column headers
    y = Inches(1.0)
    add_rect(slide, MARGIN, y, CONTENT_W, Inches(0.3), NAVY)
    add_text(slide, Inches(0.6), y + Inches(0.03), Inches(3.0), Inches(0.24),
             "Channel", size=10, color=WHITE, bold=True)
    add_text(slide, Inches(3.8), y + Inches(0.03), Inches(1.5), Inches(0.24),
             "Cost", size=10, color=WHITE, bold=True, align=PP_ALIGN.RIGHT)
    add_text(slide, Inches(5.5), y + Inches(0.03), Inches(4.0), Inches(0.24),
             "Description", size=10, color=WHITE, bold=True)
    y += Inches(0.3)

    ch_notes = []
    row_h = Inches(0.45)
    for j, ch in enumerate(channels):
        name = ch.get("name", f"Channel {j+1}")
        value = strip_html(str(ch.get("value", "—")))
        desc = strip_html(ch.get("description", ""))
        first_sent = (desc.split('. ')[0] + '.') if '. ' in desc else desc[:100]

        if j % 2 == 0:
            add_rect(slide, MARGIN, y, CONTENT_W, row_h, CREAM)

        add_text(slide, Inches(0.6), y + Inches(0.05), Inches(3.0), Inches(0.35),
                 name, size=11, color=NAVY, bold=True)
        add_text(slide, Inches(3.8), y + Inches(0.05), Inches(1.5), Inches(0.35),
                 value, size=11, color=CRIMSON, bold=True, align=PP_ALIGN.RIGHT)
        add_text(slide, Inches(5.5), y + Inches(0.05), Inches(4.0), Inches(0.35),
                 first_sent[:100], size=9, color=MID_GRAY)

        ch_notes.append(f"{name}: {value}. {first_sent}")
        y += row_h
        if y > Inches(4.5):
            break

    add_notes(slide, "Welfare channels:\n" + "\n".join(ch_notes))
    count += 1

    # Detail slides for substantial channels
    for ch in channels:
        desc = strip_html(ch.get("description", ""))
        if len(desc) > 200:
            name = ch.get("name", "Channel")
            value = strip_html(str(ch.get("value", "")))
            title = action_title(f"{name}: {value}", desc)
            count += build_text_slides(prs, title, desc)

    return count

def build_cases_slides(prs, cases):
    if not cases:
        return 0
    count = 0
    for case in cases:
        title = case.get("title", "Case Study")
        content = strip_html(case.get("content", case.get("description", "")))
        if content:
            at = action_title(f"Case: {title}", content)
            count += build_text_slides(prs, at, content)
    return count

def build_policy_slides(prs, data):
    policy = data.get("policyAnalysis")
    if not policy or not isinstance(policy, dict):
        return 0
    count = 0

    sections = [
        ("currentFramework", "The current regulatory framework and its structural gaps"),
        ("failures", "Why the existing framework fails to internalize welfare costs"),
        ("reform", "Reform pathway: what would actually work"),
    ]

    for key, label in sections:
        content = policy.get(key, "")
        if content:
            at = action_title(label, content)
            count += build_text_slides(prs, at, content)
    return count

def build_agents_slides(prs, agents):
    if not agents or not isinstance(agents, dict):
        return 0

    meta = {
        "whistleblower": ("Whistleblower", "Breaking information asymmetry"),
        "plaintiff": ("Plaintiff", "Monetizing liability to realign incentives"),
        "regulator": ("Regulator", "Redesigning the rules of the game"),
        "legislator": ("Legislator", "Changing the law to close structural gaps"),
        "investor": ("Investor", "Repricing capital to reflect true costs"),
        "supranational": ("Supranational", "Coordinating across jurisdictions"),
    }

    count = 0
    for key in ["whistleblower", "plaintiff", "regulator", "legislator", "investor", "supranational"]:
        content = agents.get(key) or (agents.get("insider") if key == "whistleblower" else None)
        if not content:
            continue
        label, sub = meta.get(key, (key.title(), ""))
        at = action_title(f"{label}: {sub}", content)
        count += build_text_slides(prs, at, content)
    return count

def build_cross_domain_slides(prs, connections):
    if not connections:
        return 0
    count = 0
    for i in range(0, len(connections), 2):
        batch = connections[i:i+2]
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        set_bg(slide)

        add_gold_line(slide)
        add_text(slide, MARGIN, Inches(0.15), Inches(9.0), Inches(0.50),
                 "Cross-domain connections reveal shared structural patterns",
                 size=14, color=NAVY, bold=True, font_name=FONT_TITLE)

        y = Inches(1.1)
        conn_notes = []
        for conn in batch:
            domain = conn.get("domain", "")
            desc = strip_html(conn.get("connection", conn.get("description", "")))[:600]

            est_h = max(Inches(1.8), Inches(0.3 * (len(desc) // 80 + 2)))
            est_h = min(est_h, Inches(2.8))

            add_rect(slide, MARGIN, y, CONTENT_W, est_h, CREAM)
            add_rect(slide, MARGIN, y, Inches(0.06), est_h, BLUE_GRAY)
            add_text(slide, Inches(0.7), y + Inches(0.1), Inches(8.3), Inches(0.3),
                     domain, size=14, color=NAVY, bold=True)
            add_body_text(slide, Inches(0.7), y + Inches(0.45), Inches(8.3), est_h - Inches(0.55),
                          desc, size=12, color=BLACK, line_spacing=1.15)
            y += est_h + Inches(0.15)
            conn_notes.append(f"{domain}: {desc}")

        add_notes(slide, "\n\n".join(conn_notes))
        count += 1

    return count

def build_faq_slides(prs, faq):
    if not faq:
        return 0
    count = 0
    groups = []
    grp = []
    grp_len = 0
    for item in faq:
        q = strip_html(item.get("q", ""))
        a = strip_html(item.get("a", ""))
        entry_len = len(q) + len(a)
        if grp_len + entry_len > 1400 and grp:
            groups.append(grp)
            grp = [(q, a)]
            grp_len = entry_len
        else:
            grp.append((q, a))
            grp_len += entry_len
    if grp:
        groups.append(grp)

    for gi, group in enumerate(groups):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        set_bg(slide)

        pg = f" ({gi+1}/{len(groups)})" if len(groups) > 1 else ""
        add_gold_line(slide)
        add_text(slide, MARGIN, Inches(0.15), Inches(9.0), Inches(0.50),
                 f"Frequently asked questions{pg}",
                 size=14, color=NAVY, bold=True, font_name=FONT_TITLE)

        y = Inches(1.1)
        faq_notes = []
        for q, a in group:
            add_text(slide, MARGIN, y, CONTENT_W, Inches(0.35),
                     f"Q: {q[:180]}", size=13, color=NAVY, bold=True)

            a_text = a[:600]
            est_h = max(Inches(0.8), Inches(0.25 * (len(a_text) // 80 + 1)))
            est_h = min(est_h, Inches(1.8))

            add_body_text(slide, Inches(0.7), y + Inches(0.4), Inches(8.5), est_h,
                          a_text, size=12, color=BLACK, line_spacing=1.15)
            faq_notes.append(f"Q: {q}\nA: {a}")
            y += Inches(0.45) + est_h + Inches(0.15)
            if y > Inches(4.3):
                break

        add_notes(slide, "\n\n".join(faq_notes))
        count += 1

    return count

def build_timeline_slides(prs, timeline):
    if not timeline:
        return 0
    count = 0
    for i in range(0, len(timeline), 10):
        batch = timeline[i:i+10]
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        set_bg(slide)

        add_gold_line(slide)
        add_text(slide, MARGIN, Inches(0.15), Inches(9.0), Inches(0.50),
                 "Key events in the evolution of this domain",
                 size=14, color=NAVY, bold=True, font_name=FONT_TITLE)

        y = Inches(1.0)
        tl_notes = []
        for ev in batch:
            year = str(ev.get("year", ""))
            evt = strip_html(ev.get("event", ""))
            sig = strip_html(ev.get("significance", ""))

            # Timeline dot
            add_rect(slide, Inches(0.7), y + Inches(0.08), Inches(0.12), Inches(0.12), NAVY)

            add_text(slide, Inches(0.5), y, Inches(0.9), Inches(0.3),
                     year, size=11, color=NAVY, bold=True)
            add_text(slide, Inches(1.5), y, Inches(3.5), Inches(0.3),
                     evt[:80], size=11, color=BLACK, bold=True)
            add_text(slide, Inches(5.2), y, Inches(4.3), Inches(0.5),
                     sig[:160], size=10, color=MID_GRAY)
            tl_notes.append(f"{year}: {evt}. {sig}")
            y += Inches(0.55)
            if y > Inches(4.5):
                break

        add_notes(slide, "\n".join(tl_notes))
        count += 1

    return count

def build_quiz_slide(prs, data, page_num=0):
    wa = data.get("welfareAnalysis")
    wa = wa if isinstance(wa, dict) else {}
    beta = wa.get("betaW", "?")
    tt = data.get("theoremType", "?")
    tn = data.get("theoremName", "?")

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, DARK_NAVY)
    add_rect(slide, 0, 0, SLIDE_W, GOLD_LINE_H, GOLD)

    add_text(slide, Inches(0.50), Inches(0.25), Inches(8), Inches(0.25),
             "COMPREHENSION CHECK", size=10, color=GOLD, bold=True)

    qs = [
        (f"What is βW for this domain?", f"βW = {beta}"),
        (f"Impossibility or Intractability?", f"{tt}"),
        (f"Theorem name?", f"{tn}"),
        ("Solvable with policy?",
         "No — physical constraints" if tt == "Impossibility"
         else "Yes — proven models exist"),
        ("Iron Law for Π?", "Revenue, never profit."),
    ]

    y = Inches(0.65)
    for q, a in qs:
        add_text(slide, Inches(0.50), y, Inches(8.5), Inches(0.28),
                 q, size=12, color=WHITE, bold=True)
        add_text(slide, Inches(0.80), y + Inches(0.30), Inches(8.2), Inches(0.25),
                 a, size=11, color=GOLD)
        y += Inches(0.80)

    add_footer(slide, page_num, "SAPM Masterclass")
    add_notes(slide, "\n".join(f"{q} → {a}" for q, a in qs))
    return 1

def build_closing_slide(prs, data, page_num=0):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, DARK_NAVY)
    add_rect(slide, 0, 0, SLIDE_W, GOLD_LINE_H, GOLD)

    title = data.get("title", "")
    wa = data.get("welfareAnalysis")
    wa = wa if isinstance(wa, dict) else {}
    beta = wa.get("betaW", "")
    tt = data.get("theoremType", "")

    add_text(slide, Inches(0.5), Inches(0.5), Inches(9), Inches(1.0),
             title, size=16, color=WHITE, bold=True, font_name=FONT_TITLE,
             align=PP_ALIGN.CENTER, line_spacing=1.15)

    add_line(slide, Inches(3.5), Inches(1.7), Inches(3.0), GOLD, Pt(2))

    if beta and beta != "—":
        add_text(slide, Inches(0.5), Inches(2.0), Inches(9), Inches(0.6),
                 f"βW = {beta}", size=32, color=GOLD, bold=True,
                 align=PP_ALIGN.CENTER)

    if tt:
        c = CRIMSON if tt == "Impossibility" else GOLD
        add_text(slide, Inches(0.5), Inches(2.8), Inches(9), Inches(0.35),
                 f"{tt} Theorem", size=14, color=c, align=PP_ALIGN.CENTER)

    add_text(slide, Inches(0.5), Inches(3.6), Inches(9), Inches(0.25),
             "Erik Postnieks  |  Independent Researcher  |  postnieks.com",
             size=9, color=MID_GRAY, align=PP_ALIGN.CENTER)
    add_text(slide, Inches(0.5), Inches(3.9), Inches(9), Inches(0.25),
             "Paper: SSRN  |  MC Replication: github.com/epostnieks",
             size=8, color=MID_GRAY, align=PP_ALIGN.CENTER)

    add_footer(slide, page_num, "SAPM Masterclass")
    add_notes(slide, f"End: {title}. βW = {beta}. {tt} Theorem.")
    return 1


# ══════════════════════════════════════════════════════════════
# TTS AUDIO
# ══════════════════════════════════════════════════════════════

def generate_tts(prs, slug):
    import subprocess, tempfile
    all_notes = []
    for slide in prs.slides:
        try:
            ns = slide.notes_slide
            text = ns.notes_text_frame.text
            if text and text.strip():
                all_notes.append(text.strip())
        except Exception:
            pass

    if not all_notes:
        return None

    full_script = "\n\n".join(all_notes)
    full_script = full_script.replace('βW', 'beta W')
    full_script = full_script.replace('ΔW', 'delta W')
    full_script = full_script.replace('Π', 'Pi')
    full_script = re.sub(r'[βΠΔ]', '', full_script)

    out_path = AUDIO_DIR / f"{slug}.m4a"
    try:
        with tempfile.NamedTemporaryFile(mode='w', suffix='.txt', delete=False) as f:
            f.write(full_script)
            tmp_path = f.name

        subprocess.run([
            'say', '-v', 'Samantha', '-r', '190',
            '-f', tmp_path,
            '-o', str(out_path)
        ], check=True, timeout=600)

        os.unlink(tmp_path)
        return out_path
    except Exception as e:
        print(f"    TTS error for {slug}: {e}")
        return None


# ══════════════════════════════════════════════════════════════
# MAIN GENERATOR
# ══════════════════════════════════════════════════════════════

def generate(json_path):
    with open(json_path, 'r') as f:
        data = json.load(f)

    slug = json_path.stem
    figures = find_figures(slug)
    docx_path = find_docx(slug)
    tables = get_docx_tables(docx_path) if docx_path else []
    images = get_docx_images(docx_path) if docx_path else []
    fig_map = categorize_figures(figures)

    # Create presentation
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    total = 0
    tt = data.get("theoremType", "")
    tn = data.get("theoremName", "")

    # ── Title ──
    badge = clean_theorem_display(tt, tn)
    ep = strip_html(data.get("epigraph", ""))
    slide_title_dark(prs, data.get("title", slug), badge=badge, epigraph=ep)
    total += 1

    # ── Ch 1: The Numbers ──
    slide_section(prs, 1, "THE NUMBERS",
        "System beta, welfare cost, and Monte Carlo quantification")
    total += 1
    slide_metrics(prs, data)
    total += 1
    for fig in fig_map.get('mc', []):
        slide_figure(prs, fig, f"Monte Carlo Distribution — {slug}")
        total += 1
    for fig in fig_map.get('sensitivity', []):
        slide_figure(prs, fig, f"Sensitivity Analysis — {slug}")
        total += 1

    # ── Ch 2: Executive Summary ──
    exec_sum = data.get("executiveSummary", "")
    if exec_sum:
        slide_section(prs, 2, "EXECUTIVE SUMMARY",
            "The complete argument in one chapter")
        total += 1
        total += build_text_slides(prs, "What this paper proves and why it matters", exec_sum)

        # Key insight slide from first 2 sentences
        sentences = split_sentences(exec_sum)
        if len(sentences) >= 2:
            slide_key_insight(prs, sentences[0],
                              ' '.join(sentences[1:4]))
            total += 1

    # ── Ch 3: Key Findings ──
    findings = data.get("keyFindings", [])
    if findings:
        slide_section(prs, 3, "KEY FINDINGS",
            f"{len(findings)} findings from the research")
        total += 1
        total += build_findings_slides(prs, findings)
        for fig in fig_map.get('channel', []):
            slide_figure(prs, fig, f"Welfare Channel Breakdown")
            total += 1

    # ── Ch 4: Theorem ──
    slide_section(prs, 4, clean_theorem_display(tt, tn) or "THE THEOREM",
        "Formal proof and structural assumptions")
    total += 1
    total += build_theorem_slides(prs, data)
    for fig in fig_map.get('theory', []):
        slide_figure(prs, fig, "Theoretical Framework")
        total += 1
    for fig in fig_map.get('game', []):
        slide_figure(prs, fig, "Game Theory and Payoff Analysis")
        total += 1

    # ── Ch 5: Welfare Channels ──
    slide_section(prs, 5, "WELFARE CHANNELS",
        "Where the damage comes from — channel by channel")
    total += 1
    total += build_welfare_slides(prs, data)
    for fig in fig_map.get('welfare', []):
        slide_figure(prs, fig, "Welfare Analysis")
        total += 1

    # Additional welfare content
    wa = data.get("welfareAnalysis")
    if isinstance(wa, dict):
        ct = wa.get("channelTable", "")
        if ct:
            total += build_text_slides(prs, "Welfare channel data table", ct)

    # ── Ch 6: Monte Carlo ──
    mc_text = ""
    if isinstance(wa, dict):
        mc_text = wa.get("monteCarlo", "")
    if mc_text:
        slide_section(prs, 6, "MONTE CARLO SIMULATION",
            "Quantifying uncertainty: 10,000 draws, reproducible")
        total += 1
        total += build_text_slides(prs, "Monte Carlo methodology and results", mc_text)

    # ── Ch 7: Paper Figures ──
    remaining_figs = (fig_map.get('other', []) + fig_map.get('policy', []) +
                      fig_map.get('agent', []) + fig_map.get('timeline', []))
    if remaining_figs:
        slide_section(prs, 7, "PAPER FIGURES",
            f"{len(remaining_figs)} figures from the published research")
        total += 1
        for fig in remaining_figs:
            slide_figure(prs, fig, fig_label(fig))
            total += 1

    # ── Ch 8: Evidence ──
    cases = data.get("caseStudies", [])
    if cases or tables or images:
        slide_section(prs, 8, "EVIDENCE",
            f"{len(cases)} case studies, {len(tables)} data tables, {len(images)} figures")
        total += 1
        total += build_cases_slides(prs, cases)

        for table in tables:
            s = slide_table(prs, table['title'], table['header'], table['rows'],
                            source_text=f"{slug} published paper")
            if s:
                total += 1

        for i, img in enumerate(images[:20]):
            slide_img = prs.slides.add_slide(prs.slide_layouts[6])
            set_bg(slide_img)
            add_gold_line(slide_img)
            add_text(slide_img, MARGIN, Inches(0.15), Inches(9.0), Inches(0.50),
                     f"Paper figure {i+1}", size=14, color=NAVY, bold=True, font_name=FONT_TITLE)
            try:
                slide_img.shapes.add_picture(BytesIO(img["data"]),
                    left=Inches(0.3), top=Inches(0.70),
                    width=Inches(9.4), height=Inches(4.5))
            except Exception:
                pass
            add_notes(slide_img, f"Embedded figure {i+1} from the published paper.")
            total += 1

    # ── Ch 9: Policy ──
    if data.get("policyAnalysis"):
        slide_section(prs, 9, "POLICY ANALYSIS",
            "What works, what fails, and why")
        total += 1
        total += build_policy_slides(prs, data)
        for fig in fig_map.get('policy', []):
            slide_figure(prs, fig, "Policy Analysis")
            total += 1

    # ── Ch 10: Six-Agent Framework ──
    agents = data.get("agents", {})
    if agents:
        slide_section(prs, 10, "SIX-AGENT FRAMEWORK",
            "Who can break the Private-Systemic Tension")
        total += 1
        total += build_agents_slides(prs, agents)

    # ── Ch 11: Context ──
    connections = data.get("crossDomainConnections", [])
    lit = data.get("literatureContext", "")
    method = data.get("methodology", "")
    timeline = data.get("timeline", [])
    if connections or lit or method or timeline:
        slide_section(prs, 11, "CONNECTIONS AND CONTEXT",
            "Cross-domain patterns, literature, methodology")
        total += 1
        total += build_cross_domain_slides(prs, connections)
        for fig in fig_map.get('crossdomain', []):
            slide_figure(prs, fig, "Cross-Domain Map")
            total += 1
        for fig in fig_map.get('literature', []):
            slide_figure(prs, fig, "Literature Context")
            total += 1
        if lit:
            total += build_text_slides(prs, "Where this paper sits in the literature", lit)
        if method:
            total += build_text_slides(prs, "Methodology, data sources, and limitations", method)
        total += build_timeline_slides(prs, timeline)

    # ── Ch 12: FAQ + Quiz + Close ──
    faq = data.get("faq", [])
    if faq:
        slide_section(prs, 12, "QUESTIONS AND ASSESSMENT",
            "FAQ and comprehension check")
        total += 1
        total += build_faq_slides(prs, faq)

    total += build_quiz_slide(prs, data)
    build_closing_slide(prs, data)
    total += 1

    # Save
    out_path = PPTX_DIR / f"{slug}.pptx"
    prs.save(str(out_path))

    # TTS
    audio_path = generate_tts(prs, slug)

    return slug, total, len(figures), len(tables), len(images), audio_path is not None


def main():
    json_files = sorted(PAPER_DATA_DIR.glob("*.json"))
    json_files = [f for f in json_files if '_raw' not in f.name]
    print(f"Found {len(json_files)} papers")

    results = []
    errors = []

    with ThreadPoolExecutor(max_workers=6) as pool:
        futures = {pool.submit(generate, f): f for f in json_files}
        for future in as_completed(futures):
            f = futures[future]
            try:
                slug, slides, figs, tbls, imgs, has_audio = future.result()
                results.append((slug, slides, figs, tbls, imgs, has_audio))
                audio_str = " +audio" if has_audio else ""
                print(f"  OK  {slug}: {slides} slides, {figs} figs, {tbls} tables{audio_str}")
            except Exception as e:
                errors.append((f.stem, str(e)))
                print(f"  FAIL {f.stem}: {e}")
                traceback.print_exc()

    results.sort(key=lambda x: x[0])
    total_slides = sum(r[1] for r in results)
    total_audio = sum(1 for r in results if r[5])

    print(f"\nDone: {len(results)} OK, {len(errors)} failed")
    print(f"Total slides: {total_slides} (avg {total_slides // max(len(results),1)}/paper)")
    print(f"Audio: {total_audio}/{len(results)}")

    if errors:
        print(f"\nFailed:")
        for slug, err in errors:
            print(f"  {slug}: {err}")


if __name__ == "__main__":
    main()
